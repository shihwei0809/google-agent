import os
import io
import uuid
import shutil
from datetime import datetime
from PIL import Image

import uvicorn
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel

import easyocr
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import cv2   # 用於影像修補 (Inpainting) 去除原圖文字
import numpy as np

app = FastAPI(title="DeckEdit Offline Converter")

# 啟用 CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 建立輸出暫存區
OUTPUT_DIR = "outputs"
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

# 初始化 EasyOCR 模型
print("[*] 正在載入 EasyOCR 繁體中文與英文模型...")
reader = easyocr.Reader(['ch_tra', 'en'], gpu=False)
print("[+] OCR 模型載入完成，網頁伺服器啟動中...")

# 智慧文字背景色偵測
def get_background_color(image, box):
    try:
        left = int(min(p[0] for p in box))
        top = int(min(p[1] for p in box))
        right = int(max(p[0] for p in box))
        bottom = int(max(p[1] for p in box))
        
        w, h = image.size
        left = max(0, left - 2)
        top = max(0, top - 2)
        right = min(w - 1, right + 2)
        bottom = min(h - 1, bottom + 2)
        
        if right <= left or bottom <= top:
            return (255, 255, 255)
            
        cropped = image.crop((left, top, right, bottom))
        cw, ch = cropped.size
        
        edge_pixels = []
        for x in range(cw):
            edge_pixels.append(cropped.getpixel((x, 0)))
            edge_pixels.append(cropped.getpixel((x, ch - 1)))
        for y in range(ch):
            edge_pixels.append(cropped.getpixel((0, y)))
            edge_pixels.append(cropped.getpixel((cw - 1, y)))
            
        r_sum = g_sum = b_sum = 0
        count = 0
        for p in edge_pixels:
            if isinstance(p, tuple):
                r_sum += p[0]
                g_sum += p[1]
                b_sum += p[2]
            else:
                r_sum += p
                g_sum += p
                b_sum += p
            count += 1
            
        if count == 0:
            return (255, 255, 255)
            
        return (int(r_sum / count), int(g_sum / count), int(b_sum / count))
    except Exception:
        return (255, 255, 255)

# 影像修補，抹除底圖文字的函數
def remove_text_from_image_bytes(img_bytes, results):
    try:
        nparr = np.frombuffer(img_bytes, np.uint8)
        image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        if image is None:
            return img_bytes
            
        mask = np.zeros(image.shape[:2], dtype=np.uint8)
        for box, text, conf in results:
            pts = np.array(box, dtype=np.int32)
            rect = cv2.boundingRect(pts)
            x, y, w, h = rect
            # 稍微外擴以確保完全覆蓋邊緣
            x = max(0, x - 2)
            y = max(0, y - 2)
            w = w + 4
            h = h + 4
            cv2.rectangle(mask, (x, y), (x + w, y + h), 255, -1)
            
        inpainted = cv2.inpaint(image, mask, 3, cv2.INPAINT_TELEA)
        
        is_success, encoded_img = cv2.imencode(".png", inpainted)
        if is_success:
            return encoded_img.tobytes()
    except Exception as e:
        print(f"[-] Inpainting 失敗: {e}")
    return img_bytes

# 首頁接口
@app.get("/")
async def get_index():
    if os.path.exists("web_index.html"):
        return FileResponse("web_index.html")
    else:
        return JSONResponse(status_code=404, content={"message": "web_index.html not found"})

# 下載接口
@app.get("/download/{file_id}")
async def download_file(file_id: str, name: str = "result.docx"):
    file_path = os.path.join(OUTPUT_DIR, file_id)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="檔案不存在或已過期")
    return FileResponse(file_path, filename=name, media_type="application/octet-stream")

# 轉換核心接口
@app.post("/convert")
async def convert_file(
    file: UploadFile = File(...),
    format: str = Form(...)
):
    try:
        filename = file.filename
        file_bytes = await file.read()
        is_pdf = filename.lower().endswith('.pdf')
        
        unique_id = str(uuid.uuid4())
        ext = ".pptx" if format == "ppt" else ".docx"
        output_filename = unique_id + ext
        output_filepath = os.path.join(OUTPUT_DIR, output_filename)
        
        if format == "word":
            doc_word = Document()
            results_all = []
            
            if is_pdf:
                pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
                total_pages = len(pdf_doc)
                doc_word.add_heading(f"PDF OCR 辨識結果 - {filename}", level=1)
                
                for page_num in range(total_pages):
                    page = pdf_doc[page_num]
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("png")
                    
                    results = reader.readtext(img_bytes, detail=0)
                    results_all.extend(results)
                    
                    doc_word.add_heading(f"--- 第 {page_num + 1} 頁 ---", level=2)
                    for line in results:
                        doc_word.add_paragraph(line)
                    
                    if page_num < total_pages - 1:
                        doc_word.add_page_break()
            else:
                results = reader.readtext(file_bytes, detail=0)
                results_all.extend(results)
                
                doc_word.add_heading(f"圖片 OCR 辨識結果 - {filename}", level=1)
                doc_word.add_paragraph(f"辨識時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                doc_word.add_paragraph("=" * 50)
                for line in results:
                    doc_word.add_paragraph(line)
                    
            doc_word.save(output_filepath)
            
        elif format == "ppt":
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]
            
            slide_w_in = 13.333
            slide_h_in = 7.5
            
            results_all_text = []
            
            if is_pdf:
                pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
                total_pages = len(pdf_doc)
                
                for page_num in range(total_pages):
                    page = pdf_doc[page_num]
                    pix = page.get_pixmap(dpi=150)
                    img_width = pix.width
                    img_height = pix.height
                    img_bytes = pix.tobytes("png")
                    
                    # 辨識文字座標
                    results = reader.readtext(img_bytes, detail=1)
                    results_all_text.extend([item[1] for item in results])
                    
                    # 影像修補，抹除文字底圖
                    cleaned_img_bytes = remove_text_from_image_bytes(img_bytes, results)
                    
                    pil_img = Image.open(io.BytesIO(img_bytes))
                    
                    slide = prs.slides.add_slide(blank_layout)
                    
                    img_stream = io.BytesIO(cleaned_img_bytes)
                    slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                    
                    # 疊加可編輯文字框 (維持透明背景)
                    for box, text, conf in results:
                        left_px = min(p[0] for p in box)
                        top_px = min(p[1] for p in box)
                        right_px = max(p[0] for p in box)
                        bottom_px = max(p[1] for p in box)
                        
                        width_px = right_px - left_px
                        height_px = bottom_px - top_px
                        
                        left_in = (left_px / img_width) * slide_w_in
                        top_in = (top_px / img_height) * slide_h_in
                        width_in = max(0.05, (width_px / img_width) * slide_w_in)
                        height_in = max(0.05, (height_px / img_height) * slide_h_in)
                        
                        txBox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        tf.margin_left = Inches(0)
                        tf.margin_right = Inches(0)
                        tf.margin_top = Inches(0)
                        tf.margin_bottom = Inches(0)
                        
                        p = tf.paragraphs[0]
                        p.text = text
                        p.font.name = 'Microsoft JhengHei'
                        p.font.size = Pt(11)
                        
                        # 依照明暗自動變更文字顏色
                        bg_color = get_background_color(pil_img, box)
                        brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
                        if brightness < 128:
                            p.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            p.font.color.rgb = RGBColor(30, 41, 59)
            else:
                # 圖片轉 PPT
                img_stream = io.BytesIO(file_bytes)
                pil_img = Image.open(img_stream)
                img_width, img_height = pil_img.size
                
                results = reader.readtext(file_bytes, detail=1)
                results_all_text.extend([item[1] for item in results])
                
                # 影像修補，抹除底圖文字
                cleaned_img_bytes = remove_text_from_image_bytes(file_bytes, results)
                
                slide = prs.slides.add_slide(blank_layout)
                
                # 新增修補後的無字底圖
                img_stream = io.BytesIO(cleaned_img_bytes)
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                
                # 疊加編輯方塊 (透明背景)
                for box, text, conf in results:
                    left_px = min(p[0] for p in box)
                    top_px = min(p[1] for p in box)
                    right_px = max(p[0] for p in box)
                    bottom_px = max(p[1] for p in box)
                    
                    width_px = right_px - left_px
                    height_px = bottom_px - top_px
                    
                    left_in = (left_px / img_width) * slide_w_in
                    top_in = (top_px / img_height) * slide_h_in
                    width_in = max(0.05, (width_px / img_width) * slide_w_in)
                    height_in = max(0.05, (height_px / img_height) * slide_h_in)
                    
                    txBox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0)
                    tf.margin_right = Inches(0)
                    tf.margin_top = Inches(0)
                    tf.margin_bottom = Inches(0)
                    
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(12)
                    
                    bg_color = get_background_color(pil_img, box)
                    brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
                    if brightness < 128:
                        p.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        p.font.color.rgb = RGBColor(30, 41, 59)
                    
            prs.save(output_filepath)
            
        # 整理出所有辨識到的純文字
        flat_text = ""
        if format == "word":
            flat_text = "\n".join(results_all)
        elif format == "ppt":
            flat_text = "\n".join(results_all_text)

        # 產出下載檔名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        download_name = os.path.splitext(filename)[0] + f"_ocr_{timestamp}" + (".pptx" if format == "ppt" else ".docx")
        
        return {
            "success": True,
            "download_url": f"/download/{output_filename}",
            "filename": download_name,
            "text": flat_text
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JSONResponse(status_code=500, content={"success": False, "message": str(e)})

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
