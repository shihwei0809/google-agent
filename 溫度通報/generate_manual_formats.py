import os
import re
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors

# ---------------------------------------------------------
# DOCX Generation
# ---------------------------------------------------------
def md_to_docx(md_path, docx_path):
    doc = docx.Document()
    
    # Set margins to 1 inch
    for section in doc.sections:
        section.top_margin = docx.shared.Inches(1)
        section.bottom_margin = docx.shared.Inches(1)
        section.left_margin = docx.shared.Inches(1)
        section.right_margin = docx.shared.Inches(1)

    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    in_code_block = False
    code_text = []
    
    # Table parser state
    in_table = False
    table_rows = []

    for line in lines:
        stripped = line.strip()
        
        # Table detection
        if stripped.startswith('|') and not in_code_block:
            if re.match(r'^\|[\s:-|]+\|$', stripped):
                continue
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            table_rows.append(cells)
            in_table = True
            continue
        elif in_table:
            if table_rows:
                # Add table to doc
                cols_count = len(table_rows[0])
                table = doc.add_table(rows=0, cols=cols_count)
                table.style = 'Light Shading Accent 1'
                for row_cells in table_rows:
                    row = table.add_row()
                    for c_idx, cell_value in enumerate(row_cells):
                        cell = row.cells[c_idx]
                        cell.text = cell_value
                # Add spacing
                doc.add_paragraph()
            table_rows = []
            in_table = False

        # Code blocks
        if stripped.startswith('```'):
            if in_code_block:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = docx.shared.Inches(0.3)
                p_run = p.add_run(''.join(code_text))
                p_run.font.name = 'Consolas'
                p_run.font.size = docx.shared.Pt(9.5)
                p.paragraph_format.space_after = docx.shared.Pt(6)
                code_text = []
                in_code_block = False
            else:
                in_code_block = True
            continue
            
        if in_code_block:
            code_text.append(line)
            continue
            
        # Headers
        if stripped.startswith('# '):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith('## '):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith('### '):
            doc.add_heading(stripped[4:], level=3)
        # Bullet points
        elif stripped.startswith('* ') or stripped.startswith('- '):
            p = doc.add_paragraph(style='List Bullet')
            p.add_run(stripped[2:])
        # Numbered list
        elif re.match(r'^\d+\.\s', stripped):
            match = re.match(r'^(\d+)\.\s(.*)', stripped)
            p = doc.add_paragraph(style='List Number')
            p.add_run(match.group(2))
        # Empty line
        elif not stripped:
            continue
        # Blockquote
        elif stripped.startswith('> '):
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = docx.shared.Inches(0.2)
            p_run = p.add_run(stripped[2:])
            p_run.font.italic = True
        # Normal text
        else:
            doc.add_paragraph(stripped)
            
    doc.save(docx_path)
    print(f"Docx generated successfully at {docx_path}")


# ---------------------------------------------------------
# PDF Generation
# ---------------------------------------------------------
def md_to_pdf(md_path, pdf_path):
    font_path = "C:\\Windows\\Fonts\\msjh.ttc"
    if os.path.exists(font_path):
        pdfmetrics.registerFont(TTFont('MSJH', font_path))
        font_name = 'MSJH'
    else:
        font_name = 'Helvetica'
        
    doc = SimpleDocTemplate(pdf_path, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'ManualTitle',
        parent=styles['Title'],
        fontName=font_name,
        fontSize=18,
        leading=22,
        spaceAfter=15,
        alignment=1 # Center
    )
    h1_style = ParagraphStyle(
        'ManualH1',
        parent=styles['Heading1'],
        fontName=font_name,
        fontSize=13,
        leading=16,
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    h2_style = ParagraphStyle(
        'ManualH2',
        parent=styles['Heading2'],
        fontName=font_name,
        fontSize=11,
        leading=14,
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    body_style = ParagraphStyle(
        'ManualBody',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=9.5,
        leading=14,
        spaceAfter=8
    )
    code_style = ParagraphStyle(
        'ManualCode',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=8,
        leading=11,
        leftIndent=15,
        spaceAfter=4,
        backColor="#f4f4f4"
    )
    table_cell_style = ParagraphStyle(
        'TableCell',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=8.5,
        leading=11
    )
    
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    story = []
    in_code_block = False
    in_table = False
    table_rows = []
    
    for line in lines:
        stripped = line.strip()
        
        # Table detection
        if stripped.startswith('|') and not in_code_block:
            if re.match(r'^\|[\s:-|]+\|$', stripped):
                continue
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            table_rows.append(cells)
            in_table = True
            continue
        elif in_table:
            if table_rows:
                data = []
                for row_data in table_rows:
                    data.append([Paragraph(cell, table_cell_style) for cell in row_data])
                
                num_cols = len(table_rows[0])
                available_width = A4[0] - 80 # margin left & right
                col_width = available_width / num_cols
                
                t = Table(data, colWidths=[col_width]*num_cols)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#e2e8f0')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.black),
                    ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
                    ('TOPPADDING', (0,0), (-1,-1), 5),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                    ('LEFTPADDING', (0,0), (-1,-1), 5),
                    ('RIGHTPADDING', (0,0), (-1,-1), 5),
                ]))
                story.append(t)
                story.append(Spacer(1, 10))
            table_rows = []
            in_table = False

        # Code blocks
        if stripped.startswith('```'):
            in_code_block = not in_code_block
            continue
            
        if in_code_block:
            escaped = line.replace(' ', '&nbsp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
            story.append(Paragraph(escaped, code_style))
            continue
            
        # Headers
        if stripped.startswith('# '):
            story.append(Paragraph(stripped[2:], title_style))
            story.append(Spacer(1, 10))
        elif stripped.startswith('## '):
            story.append(Paragraph(stripped[3:], h1_style))
        elif stripped.startswith('### '):
            story.append(Paragraph(stripped[4:], h2_style))
        # Bullet points
        elif stripped.startswith('* ') or stripped.startswith('- '):
            bullet_text = f"• {stripped[2:]}"
            story.append(Paragraph(bullet_text, body_style))
        # Numbered list
        elif re.match(r'^\d+\.\s', stripped):
            story.append(Paragraph(stripped, body_style))
        # Blockquote
        elif stripped.startswith('> '):
            quote_style = ParagraphStyle(
                'BlockQuote', parent=body_style, fontName=font_name, fontSize=9, leftIndent=10, textColor=colors.HexColor('#475569')
            )
            story.append(Paragraph(stripped[2:], quote_style))
        # Empty line
        elif not stripped:
            story.append(Spacer(1, 6))
        # Normal text
        else:
            story.append(Paragraph(stripped, body_style))
            
    doc.build(story)
    print(f"PDF generated successfully at {pdf_path}")


# ---------------------------------------------------------
# PPTX Generation
# ---------------------------------------------------------
def create_pptx(pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette: Slate & Cyan (Premium Dark Mode)
    bg_color = RGBColor(15, 23, 42)      # Slate 900
    card_color = RGBColor(30, 41, 59)    # Slate 800
    text_white = RGBColor(255, 255, 255)
    text_muted = RGBColor(148, 163, 184) # Slate 400
    accent_cyan = RGBColor(34, 211, 238) # Cyan 400
    accent_green = RGBColor(34, 197, 94) # Green 500
    accent_yellow = RGBColor(234, 179, 8)# Yellow 500
    accent_red = RGBColor(239, 68, 68)   # Red 500

    def apply_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_header(slide, title_text):
        apply_bg(slide)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = text_white
        
        # Draw cyan underline
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(2.0), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_cyan
        line.line.fill.background()

    def add_card(slide, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.color.rgb = RGBColor(51, 65, 85) # Slate 700
        card.line.width = Pt(1)
        return card

    # Slide 1: Cover
    slide = prs.slides.add_slide(blank_layout)
    apply_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "彰化縣線西鄉環境溫度監控與雙軌備援通報系統"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "操作與維護說明書 (v4.7 整合版)"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(26)
    p2.font.color.rgb = accent_cyan
    p2.space_before = Pt(10)
    
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "整合 GAS 後端安全設定管理與 Firebase HMI 即時雲端監控儀表板，支援任意行動與桌面裝置。"
    p_sub.font.name = 'Microsoft JhengHei'
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = text_muted
    p_sub.space_after = Pt(20)
    
    p_date = tf_sub.add_paragraph()
    p_date.text = "更新日期：2026-06-12 | 系統版本：v4.7 整合版"
    p_date.font.name = 'Microsoft JhengHei'
    p_date.font.size = Pt(14)
    p_date.font.color.rgb = text_muted

    # Slide 2: Architecture
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🏗️ 系統架構（v4.7 整合雙軌備援）")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    lbl_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_left = lbl_left.text_frame
    tf_left.word_wrap = True
    p_l1 = tf_left.paragraphs[0]
    p_l1.text = "雙軌主要與備援監測端"
    p_l1.font.name = 'Microsoft JhengHei'
    p_l1.font.size = Pt(20)
    p_l1.font.bold = True
    p_l1.font.color.rgb = accent_cyan
    p_l1.space_after = Pt(15)
    
    bullets_l = [
        "主要端 (本機 Python)：每天 07:58 起每 10 分鐘執行，抓取 CWA 最新溫度，將即時資料與心跳雙寫至 Firebase 與 GAS。",
        "備援端 (GAS 雲端腳本)：每小時自動執行，偵測本機異常時接手。在 Firebase 寫入 cloud 來源心跳並發送 LINE/Email 警報。",
        "日常維護與手動同步：提供 Google 試算表選單與本機 Python 同步腳本，一鍵同步兩端資料庫。"
    ]
    for b in bullets_l:
        p = tf_left.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    lbl_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_right = lbl_right.text_frame
    tf_right.word_wrap = True
    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "Firebase Firestore 與人機網頁"
    p_r1.font.name = 'Microsoft JhengHei'
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = accent_cyan
    p_r1.space_after = Pt(15)
    
    bullets_r = [
        "Google 試算表：作為設定檔與歷史紀錄歸檔主資料庫。",
        "Firebase Firestore：即時雲端資料庫，即時同步溫度、心跳、警報狀態與聯絡人資訊。",
        "Firebase HMI 儀表板：靜態網頁託管 (Hosting)，可使用任意瀏覽器開啟，資料每秒即時同步更新。"
    ]
    for b in bullets_r:
        p = tf_right.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    # Slide 3: Two Portals HMI
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🖥️ 雙軌人機網頁 (Firebase & GAS HMI)")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(1.2))
    url_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(0.8))
    tf_url = url_box.text_frame
    tf_url.word_wrap = True
    p_u = tf_url.paragraphs[0]
    p_u.text = "即時監控儀表板網址（任意裝置可開）："
    p_u.font.name = 'Microsoft JhengHei'
    p_u.font.size = Pt(16)
    p_u.font.color.rgb = text_muted
    p_link = tf_url.add_paragraph()
    p_link.text = "https://hongsheng-temp-523.web.app"
    p_link.font.name = 'Consolas'
    p_link.font.size = Pt(22)
    p_link.font.bold = True
    p_link.font.color.rgb = accent_cyan
    p_link.space_before = Pt(5)

    add_card(slide, Inches(0.6), Inches(3.0), Inches(5.8), Inches(3.8))
    c1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(5.4), Inches(3.4))
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    p_c1_t = tf_c1.paragraphs[0]
    p_c1_t.text = "💡 免安裝與即時更新"
    p_c1_t.font.name = 'Microsoft JhengHei'
    p_c1_t.font.size = Pt(20)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = accent_cyan
    p_c1_t.space_after = Pt(10)
    bullets_c1 = [
        "無須安裝 App，使用手機、電腦、平板之瀏覽器即可監控。",
        "基於 Firestore onSnapshot 技術，資料變更時每秒自動推送至網頁，無須手動重新整理。",
        "響應式網頁設計 (RWD)，完美適配各種螢幕尺寸。"
    ]
    for b in bullets_c1:
        p = tf_c1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    add_card(slide, Inches(6.8), Inches(3.0), Inches(5.8), Inches(3.8))
    c2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.4), Inches(3.4))
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    p_c2_t = tf_c2.paragraphs[0]
    p_c2_t.text = "🔒 權限與安全登入後端"
    p_c2_t.font.name = 'Microsoft JhengHei'
    p_c2_t.font.size = Pt(20)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = accent_cyan
    p_c2_t.space_after = Pt(10)
    bullets_c2 = [
        "公開顯示：首頁即時溫度、三態心跳監控與歷史通報紀錄對所有人開放。",
        "密碼防護：系統設定、聯絡名冊新增/刪除及測試工具均需輸入密碼解鎖（預設 admin888）。",
        "手機適配：窄螢幕下選單自動轉底部懸浮列，表格支援手勢橫向滑動防止截斷。"
    ]
    for b in bullets_c2:
        p = tf_c2.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    # Slide 4: Dashboard & Heartbeats
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "📊 儀表板與即時監控")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    hb_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_hbl = hb_left.text_frame
    tf_hbl.word_wrap = True
    p_hbl = tf_hbl.paragraphs[0]
    p_hbl.text = "📊 首頁顯示資訊"
    p_hbl.font.name = 'Microsoft JhengHei'
    p_hbl.font.size = Pt(20)
    p_hbl.font.bold = True
    p_hbl.font.color.rgb = accent_cyan
    p_hbl.space_after = Pt(15)
    bullets_hbl = [
        "即時環境溫度：環形溫度計隨狀態（正常/高溫）自動變色，大字級即時顯示。",
        "警報閾值：當前設定的警報啟動溫度（°C）。",
        "監控時段與頻率：標示監測時間區段以及 GAS 的雲端備援監控頻率。",
        "本機監控端心跳：顯示最後一次收到心跳的時間，為判斷系統在線與否的關鍵指標。"
    ]
    for b in bullets_hbl:
        p = tf_hbl.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    hb_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_hbr = hb_right.text_frame
    tf_hbr.word_wrap = True
    p_hbr = tf_hbr.paragraphs[0]
    p_hbr.text = "💓 本機心跳三段監控狀態"
    p_hbr.font.name = 'Microsoft JhengHei'
    p_hbr.font.size = Pt(20)
    p_hbr.font.bold = True
    p_hbr.font.color.rgb = accent_cyan
    p_hbr.space_after = Pt(15)
    
    p_s1 = tf_hbr.add_paragraph()
    p_s1.text = "🟢 在線 (正常)"
    p_s1.font.name = 'Microsoft JhengHei'
    p_s1.font.size = Pt(16)
    p_s1.font.bold = True
    p_s1.font.color.rgb = accent_green
    p_s1_d = tf_hbr.add_paragraph()
    p_s1_d.text = "    說明：本機 Python 程式正常運行並定時上報心跳。"
    p_s1_d.font.name = 'Microsoft JhengHei'
    p_s1_d.font.size = Pt(13)
    p_s1_d.font.color.rgb = text_white
    p_s1_d.space_after = Pt(10)
    
    p_s2 = tf_hbr.add_paragraph()
    p_s2.text = "🟡 雲端監控（本機異常）"
    p_s2.font.name = 'Microsoft JhengHei'
    p_s2.font.size = Pt(16)
    p_s2.font.bold = True
    p_s2.font.color.rgb = accent_yellow
    p_s2_d = tf_hbr.add_paragraph()
    p_s2_d.text = "    說明：本機電腦無心跳，雲端 GAS 備援已接手（琥珀色閃爍）。"
    p_s2_d.font.name = 'Microsoft JhengHei'
    p_s2_d.font.size = Pt(13)
    p_s2_d.font.color.rgb = text_white
    p_s2_d.space_after = Pt(10)

    p_s3 = tf_hbr.add_paragraph()
    p_s3.text = "🔴 離線警報"
    p_s3.font.name = 'Microsoft JhengHei'
    p_s3.font.size = Pt(16)
    p_s3.font.bold = True
    p_s3.font.color.rgb = accent_red
    p_s3_d = tf_hbr.add_paragraph()
    p_s3_d.text = "    說明：本機與雲端備援皆無上報心跳，請立即前往檢查。"
    p_s3_d.font.name = 'Microsoft JhengHei'
    p_s3_d.font.size = Pt(13)
    p_s3_d.font.color.rgb = text_white

    # Slide 5: Settings & Contacts
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "⚙️ 系統設定與 👥 聯絡名冊管理")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    sc_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_scl = sc_left.text_frame
    tf_scl.word_wrap = True
    p_scl = tf_scl.paragraphs[0]
    p_scl.text = "⚙️ HMI 系統設定調整"
    p_scl.font.name = 'Microsoft JhengHei'
    p_scl.font.size = Pt(20)
    p_scl.font.bold = True
    p_scl.font.color.rgb = accent_cyan
    p_scl.space_after = Pt(15)
    bullets_scl = [
        "管理驗證：需輸入密碼以解鎖修改權限（預設 admin888）。",
        "可調參數：警報溫度閥值（°C）、監測開始與結束時間、雲端監測頻率、變更管理密碼。",
        "即時套用：點「儲存設定」後，資訊立即更新至 Firestore 雲端，前端網頁即時變更。",
        "雙軌同步：在 Google 試算表選單按「同步資料至 Firebase」或執行 python sync_sheet_to_firebase.py 將試算表與 Firebase 設定同步。"
    ]
    for b in bullets_scl:
        p = tf_scl.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    sc_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_scr = sc_right.text_frame
    tf_scr.word_wrap = True
    p_scr = tf_scr.paragraphs[0]
    p_scr.text = "👥 聯絡人名冊即時更新"
    p_scr.font.name = 'Microsoft JhengHei'
    p_scr.font.size = Pt(20)
    p_scr.font.bold = True
    p_scr.font.color.rgb = accent_cyan
    p_scr.space_after = Pt(15)
    bullets_scr = [
        "即時管理：可於網頁直接新增聯絡人、停用或刪除現有對象。",
        "通報類型支援：LINE（輸入使用者或群組 ID）、Email（收件者信箱）。",
        "一鍵啟用／停用：可勾選開關切換特定人員是否接收通知，不需刪除資料。",
        "即時同步：儲存後立即寫入 Firestore，本機與雲端下次執行時會直接讀取最新名單發送。"
    ]
    for b in bullets_scr:
        p = tf_scr.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    # Slide 6: History Logs & Alert Colors
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "📋 歷史通報紀錄與狀態顏色標示")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    rec_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_recl = rec_left.text_frame
    tf_recl.word_wrap = True
    p_recl = tf_recl.paragraphs[0]
    p_recl.text = "📋 通報紀錄功能"
    p_recl.font.name = 'Microsoft JhengHei'
    p_recl.font.size = Pt(20)
    p_recl.font.bold = True
    p_recl.font.color.rgb = accent_cyan
    p_recl.space_after = Pt(15)
    bullets_recl = [
        "歷史明細：顯示最近 100 筆通報紀錄，含時間、觀測溫度、閾值、狀態、管道等。",
        "即時篩選：右上角搜尋框可依任何關鍵字（如「本機」、「LINE」、「高溫」）即時過濾。",
        "匯出功能：點「匯出 CSV」可將歷史紀錄下載備份為 Excel 可開的 CSV 檔。"
    ]
    for b in bullets_recl:
        p = tf_recl.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    rec_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_recr = rec_right.text_frame
    tf_recr.word_wrap = True
    p_recr = tf_recr.paragraphs[0]
    p_recr.text = "🎨 警報字體與狀態顏色"
    p_recr.font.name = 'Microsoft JhengHei'
    p_recr.font.size = Pt(20)
    p_recr.font.bold = True
    p_recr.font.color.rgb = accent_cyan
    p_recr.space_after = Pt(15)
    
    p_n = tf_recr.add_paragraph()
    p_n.text = "⬛ 正常（未超標） — 灰藍色"
    p_n.font.name = 'Microsoft JhengHei'
    p_n.font.size = Pt(15)
    p_n.font.bold = True
    p_n.font.color.rgb = text_muted
    p_n_d = tf_recr.add_paragraph()
    p_n_d.text = "    狀態正常，僅寫入心跳或例行記錄，不再誤標為紅色。"
    p_n_d.font.name = 'Microsoft JhengHei'
    p_n_d.font.size = Pt(13)
    p_n_d.font.color.rgb = text_white
    p_n_d.space_after = Pt(10)
    
    p_g = tf_recr.add_paragraph()
    p_g.text = "🟢 溫度回落正常（警報解除） — 綠色"
    p_g.font.name = 'Microsoft JhengHei'
    p_g.font.size = Pt(15)
    p_g.font.bold = True
    p_g.font.color.rgb = accent_green
    p_g_d = tf_recr.add_paragraph()
    p_g_d.text = "    原高溫警報解除，溫度已降回閾值以下，發送解除通知。"
    p_g_d.font.name = 'Microsoft JhengHei'
    p_g_d.font.size = Pt(13)
    p_g_d.font.color.rgb = text_white
    p_g_d.space_after = Pt(10)

    p_r = tf_recr.add_paragraph()
    p_r.text = "🔴 高溫超標警報 — 紅色"
    p_r.font.name = 'Microsoft JhengHei'
    p_r.font.size = Pt(15)
    p_r.font.bold = True
    p_r.font.color.rgb = accent_red
    p_r_d = tf_recr.add_paragraph()
    p_r_d.text = "    溫度高於警報閾值，觸發 LINE/Email 警報通知。"
    p_r_d.font.name = 'Microsoft JhengHei'
    p_r_d.font.size = Pt(13)
    p_r_d.font.color.rgb = text_white

    # Slide 7: Python Monitor Config
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🛠️ 本機 Python 程式配置與必要檔案")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    cfg_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_cfgl = cfg_left.text_frame
    tf_cfgl.word_wrap = True
    p_cfgl = tf_cfgl.paragraphs[0]
    p_cfgl.text = "📂 必要檔案清單 (路徑 C:\\GOOGLE ANGET\\溫度通報)"
    p_cfgl.font.name = 'Microsoft JhengHei'
    p_cfgl.font.size = Pt(20)
    p_cfgl.font.bold = True
    p_cfgl.font.color.rgb = accent_cyan
    p_cfgl.space_after = Pt(15)
    bullets_cfgl = [
        "config.json：本機設定檔（含 LINE / SMTP 信箱設定）。",
        "firebase_key.json：Firebase 連線私鑰，用以直接極速讀寫 Firestore 即時狀態與紀錄。",
        "weather_monitor.py：主程式，讀取氣象 API 並進行通報與上報心跳。",
        "sync_sheet_to_firebase.py：手動補資料同步工具。"
    ]
    for b in bullets_cfgl:
        p = tf_cfgl.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(15)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    cfg_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_cfgr = cfg_right.text_frame
    tf_cfgr.word_wrap = True
    p_cfgr = tf_cfgr.paragraphs[0]
    p_cfgr.text = "⚙️ config.json 設定檔重點欄位"
    p_cfgr.font.name = 'Microsoft JhengHei'
    p_cfgr.font.size = Pt(20)
    p_cfgr.font.bold = True
    p_cfgr.font.color.rgb = accent_cyan
    p_cfgr.space_after = Pt(15)
    bullets_cfgr = [
        "cwa_api_key & cwa_station_id：氣象署開放資料 API Key 與站點編號 (C2G870)。",
        "temperature_threshold：溫度警報臨界點（例如 28.0）。",
        "web_app_url：GAS Web App URL。本機程式定時向其發送 POST 請求以更新試算表數據。",
        "line & email：LINE Bot 憑證/收件人、SMTP 伺服器與發信/收信 Gmail 帳密。"
    ]
    for b in bullets_cfgr:
        p = tf_cfgr.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = text_white
        p.space_after = Pt(8)

    # Slide 8: Windows Task Scheduler
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🛠️ Windows 工作排程器自動化設定")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2))
    sched_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    tf_sched = sched_box.text_frame
    tf_sched.word_wrap = True
    p_sched_t = tf_sched.paragraphs[0]
    p_sched_t.text = "⏰ 建立 Windows 背景定時排程步驟"
    p_sched_t.font.name = 'Microsoft JhengHei'
    p_sched_t.font.size = Pt(20)
    p_sched_t.font.bold = True
    p_sched_t.font.color.rgb = accent_cyan
    p_sched_t.space_after = Pt(15)
    bullets_sched = [
        "開啟排程器：按 Win + R 輸入 taskschd.msc 開啟工作排程器。",
        "建立基本工作：於右側點擊「建立基本工作」，設定名稱為「環境溫度監控通報系統」。",
        "觸發程序與動作：設定為「每天」執行。動作選擇「啟動程式」。",
        "設定程式引數與路徑：程式或指令碼輸入 pythonw (無視窗背景版)，引數輸入 weather_monitor.py，開始位置輸入 C:\\GOOGLE ANGET\\溫度通報。",
        "設定執行頻率：點擊該排程右鍵 -> 內容 -> 觸發程序 -> 編輯 -> 勾選「重複工作間隔」設為 10 分鐘，持續時間設為 1 天。",
        "時間校正：排程開始時間設為 07:58 AM，重複間隔為 10 分鐘，將於每天的 07:58 到 23:58 全程高頻精準觀測。"
    ]
    for b in bullets_sched:
        p = tf_sched.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    # Slide 9: GAS Cloud Backup
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "☁️ Google Apps Script (GAS) 雲端備援設定")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    cloud_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_cl = cloud_left.text_frame
    tf_cl.word_wrap = True
    p_cl = tf_cl.paragraphs[0]
    p_cl.text = "☁️ 雲端備援端角色與功能"
    p_cl.font.name = 'Microsoft JhengHei'
    p_cl.font.size = Pt(20)
    p_cl.font.bold = True
    p_cl.font.color.rgb = accent_cyan
    p_cl.space_after = Pt(15)
    bullets_cl = [
        "備援機制：當本機電腦斷電、當機、網路中斷時，雲端 GAS 會自動接手進行監控。",
        "即時狀態上報：當備援端執行且偵測到本機異常時，會在 Firebase 中寫入 heartbeat_source: \"cloud\"。",
        "HMI 狀態對應：此時 HMI 儀表板會將本機心跳燈亮起 🟡 琥珀色閃爍 並顯示「雲端監控（本機異常）」，方便管理員巡檢。"
    ]
    for b in bullets_cl:
        p = tf_cl.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(15)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    cloud_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_cr = cloud_right.text_frame
    tf_cr.word_wrap = True
    p_cr = tf_cr.paragraphs[0]
    p_cr.text = "⚙️ GAS 部署與排程設定"
    p_cr.font.name = 'Microsoft JhengHei'
    p_cr.font.size = Pt(20)
    p_cr.font.bold = True
    p_cr.font.color.rgb = accent_cyan
    p_cr.space_after = Pt(15)
    bullets_cr = [
        "代碼更新：開啟 Apps Script 編輯器，複製最新的 Code.gs 並存檔。",
        "管理部署：部署 -> 管理部署 -> 選取現有版本 -> 點編輯選「建立新版本」後部署。複製產出網址填入本機 config.json。",
        "新增觸發條件：點選左側時鐘圖示，新增執行 checkWeatherAndNotify 的時間型觸發程序，間隔設為每小時一次。"
    ]
    for b in bullets_cr:
        p = tf_cr.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(12)

    # Slide 10: Maintenance & Security
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🔒 日常維護與安全管理")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
    maint_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_ml = maint_left.text_frame
    tf_ml.word_wrap = True
    p_ml = tf_ml.paragraphs[0]
    p_ml.text = "🕹️ 日常維護操作方式"
    p_ml.font.name = 'Microsoft JhengHei'
    p_ml.font.size = Pt(20)
    p_ml.font.bold = True
    p_ml.font.color.rgb = accent_cyan
    p_ml.space_after = Pt(15)
    bullets_ml = [
        "試算表自訂選單（🌡️ 溫度通報系統）：測試即時通報、重置防重複鎖定、同步資料至 Firebase、重設欄寬。",
        "手動補歷史資料工具：執行 python sync_sheet_to_firebase.py，將試算表的歷史通報數據，補寫入 Firebase 供 HMI 網頁顯示。"
    ]
    for b in bullets_ml:
        p = tf_ml.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
    maint_right = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.4), Inches(4.8))
    tf_mr = maint_right.text_frame
    tf_mr.word_wrap = True
    p_mr = tf_mr.paragraphs[0]
    p_mr.text = "🔒 安全管理指引"
    p_mr.font.name = 'Microsoft JhengHei'
    p_mr.font.size = Pt(20)
    p_mr.font.bold = True
    p_mr.font.color.rgb = accent_cyan
    p_mr.space_after = Pt(15)
    bullets_mr = [
        "Firebase 金鑰：firebase_key.json 為雲端資料庫私鑰，權限等同最高管理員，嚴禁分享或公開。",
        "設定檔安全：config.json 包含 LINE Access Token、Gmail 帳密等敏感欄位，切勿公開。",
        "HMI 管理密碼：預設密碼為 admin888。強烈建議在系統上線後，第一時間於系統設定頁面修改為強密碼。"
    ]
    for b in bullets_mr:
        p = tf_mr.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(15)

    # Slide 11: v4.5 Updates
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🆕 v4.5 版本更新摘要")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2))
    upd_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    tf_upd = upd_box.text_frame
    tf_upd.word_wrap = True
    p_upd_t = tf_upd.paragraphs[0]
    p_upd_t.text = "🆕 系統 v4.5 新增與修改亮點"
    p_upd_t.font.name = 'Microsoft JhengHei'
    p_upd_t.font.size = Pt(20)
    p_upd_t.font.bold = True
    p_upd_t.font.color.rgb = accent_cyan
    p_upd_t.space_after = Pt(15)
    bullets_upd = [
        "雙軌寫入整合：本機端 weather_monitor.py 同時雙寫 GAS Web App 與 Firebase，完美融合兩者優勢。",
        "三態心跳燈標示：HMI 儀表板以 🟢在線、🟡雲端接管、🔴完全離線 清晰呈現系統健康狀態。",
        "高頻監控對齊：Windows 工作排程器改為 07:58 AM 開始每 10 分鐘執行，排除延遲，全程高頻精準監控。",
        "介面安全防護：GAS 網頁端新增密碼驗證（預設 admin888）與安全登出機制，行動端版面全面適配。",
        "手動同步加強：sync_sheet_to_firebase.py 整合歷史紀錄 schema，確保資料統一。"
    ]
    for b in bullets_upd:
        p = tf_upd.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    prs.save(pptx_path)
    print(f"PPTX generated successfully at {pptx_path}")

    # Slide 12: v4.7 Updates
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🆕 v4.7 版本更新摘要")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2))
    upd_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    tf_upd2 = upd_box2.text_frame
    tf_upd2.word_wrap = True
    p_upd_t2 = tf_upd2.paragraphs[0]
    p_upd_t2.text = "🆕 系統 v4.7 新增與修改亮點"
    p_upd_t2.font.name = 'Microsoft JhengHei'
    p_upd_t2.font.size = Pt(20)
    p_upd_t2.font.bold = True
    p_upd_t2.font.color.rgb = accent_cyan
    p_upd_t2.space_after = Pt(15)
    bullets_upd2 = [
        "雙模溫度趨勢圖：不論是 Firebase HMI 還是 GAS 雲端後台，皆全面整合 Chart.js 溫度趨勢圖，支援即時 24H 與歷史範圍查詢，並內建等距降採樣演算法，流暢不卡頓。",
        "精準整點心跳排程 (XX:58)：放寬監測時段至 07:00 - 24:00，搭配 60 分鐘本機節流，精確於 07:58, 08:58, ..., 23:58 發送心跳，完成每日最後一次安全監控。",
        "試算表刪除與修改即時同步：同步邏輯改為雙向比對，刪除的行會即時 DELETE 移出 Firestore。",
        "安裝式 onChange 監聽：註冊雲端 onChange 觸發器，試算表一有變更（如手動增刪）就即時自動同步至 Firebase，無須手動按鈕。",
        "HMI 功能相容性修正：修正全域變數重複宣告 JS 語法錯誤，並加入 Cache-busting 資源載入（v4.6 版本），徹底修復歷史查詢無反應問題。"
    ]
    for b in bullets_upd2:
        p = tf_upd2.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    # Slide 13: v4.8 Updates
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🆕 v4.8 版本更新摘要")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2))
    upd_box3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    tf_upd3 = upd_box3.text_frame
    tf_upd3.word_wrap = True
    p_upd_t3 = tf_upd3.paragraphs[0]
    p_upd_t3.text = "🆕 系統 v4.8 新增與修改亮點"
    p_upd_t3.font.name = 'Microsoft JhengHei'
    p_upd_t3.font.size = Pt(20)
    p_upd_t3.font.bold = True
    p_upd_t3.font.color.rgb = accent_cyan
    p_upd_t3.space_after = Pt(15)
    bullets_upd3 = [
        "新增 Microsoft Teams Webhook 通報：本機 Python 程式與雲端 GAS 備援程式皆已支援發送高溫與回落警報至 Teams 頻道。",
        "試算表設定檔動態載入：雲端 GAS 支援從「系統設定」工作表動態讀取「Teams Webhook URL」，免修改程式即可更新網址。",
        "MessageCard 卡片美化：傳送至 Teams 頻道的訊息採用結構化 MessageCard 設計，高溫警報以紅色警示，正常回落以綠色標示。"
    ]
    for b in bullets_upd3:
        p = tf_upd3.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    # Slide 14: v4.9 Updates
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "🆕 v4.9 版本更新摘要")
    
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.2))
    upd_box4 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    tf_upd4 = upd_box4.text_frame
    tf_upd4.word_wrap = True
    p_upd_t4 = tf_upd4.paragraphs[0]
    p_upd_t4.text = "🆕 系統 v4.9 新增與修改亮點"
    p_upd_t4.font.name = 'Microsoft JhengHei'
    p_upd_t4.font.size = Pt(20)
    p_upd_t4.font.bold = True
    p_upd_t4.font.color.rgb = accent_cyan
    p_upd_t4.space_after = Pt(15)
    bullets_upd4 = [
        "獨立的月度 24 小時記錄分頁：雲端 Apps Script 重構日誌寫入路由，日常每 10 分鐘心跳與溫度觀測，自動分流寫入「YYYY年X月-24小時記錄」分頁，避免主通報紀錄分頁被洗版且能完美防範跨年度名稱衝突。",
        "專屬綠色系排版風格：針對「24小時記錄」分頁套用專屬的深綠色主題表頭（#2E7D32）與排版格式，利於視覺區分，畫面更乾淨整潔。"
    ]
    for b in bullets_upd4:
        p = tf_upd4.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = text_white
        p.space_after = Pt(10)

    prs.save(pptx_path)
    print(f"PPTX generated successfully at {pptx_path}")


# ---------------------------------------------------------
# Main Execution
# ---------------------------------------------------------
if __name__ == "__main__":
    md = "溫度通報系統操作說明.md"
    docx_out = "溫度通報系統操作說明.docx"
    pdf_out = "溫度通報系統操作說明.pdf"
    pptx_out = "溫度通報系統操作說明.pptx"
    
    # 1. Generate DOCX
    md_to_docx(md, docx_out)
    
    # 2. Generate PDF
    md_to_pdf(md, pdf_out)
    
    # 3. Generate PPTX
    create_pptx(pptx_out)
    
    print("All formats (docx, pdf, pptx) v4.9 generated successfully!")

