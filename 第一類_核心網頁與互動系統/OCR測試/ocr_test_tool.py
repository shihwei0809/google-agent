import os
import io
from datetime import datetime
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from PIL import Image
import easyocr
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF 用於處理 PDF
import cv2   # 用於影像修補 (Inpainting) 去除原圖文字
import numpy as np

def get_background_color(image, box):
    try:
        left = int(min(p[0] for p in box))
        top = int(min(p[1] for p in box))
        right = int(max(p[0] for p in box))
        bottom = int(max(p[1] for p in box))
        
        w, h = image.size
        left = max(0, left - 2)
        top = max(0, top - 2)
        right = min(w - 1, right + 2)
        bottom = min(h - 1, bottom + 2)
        
        if right <= left or bottom <= top:
            return (255, 255, 255)
            
        cropped = image.crop((left, top, right, bottom))
        cw, ch = cropped.size
        
        edge_pixels = []
        for x in range(cw):
            edge_pixels.append(cropped.getpixel((x, 0)))
            edge_pixels.append(cropped.getpixel((x, ch - 1)))
        for y in range(ch):
            edge_pixels.append(cropped.getpixel((0, y)))
            edge_pixels.append(cropped.getpixel((cw - 1, y)))
            
        r_sum = g_sum = b_sum = 0
        count = 0
        for p in edge_pixels:
            if isinstance(p, tuple):
                r_sum += p[0]
                g_sum += p[1]
                b_sum += p[2]
            else:
                r_sum += p
                g_sum += p
                b_sum += p
            count += 1
            
        if count == 0:
            return (255, 255, 255)
            
        return (int(r_sum / count), int(g_sum / count), int(b_sum / count))
    except Exception:
        return (255, 255, 255)

def remove_text_from_image_bytes(img_bytes, results):
    try:
        # 解碼為 OpenCV 影像
        nparr = np.frombuffer(img_bytes, np.uint8)
        image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        if image is None:
            return img_bytes
            
        # 建立黑色遮罩
        mask = np.zeros(image.shape[:2], dtype=np.uint8)
        for box, text, conf in results:
            pts = np.array(box, dtype=np.int32)
            rect = cv2.boundingRect(pts)
            x, y, w, h = rect
            # 稍微外擴以確保完全覆蓋邊緣
            x = max(0, x - 2)
            y = max(0, y - 2)
            w = w + 4
            h = h + 4
            cv2.rectangle(mask, (x, y), (x + w, y + h), 255, -1)
            
        # 執行修補 (Inpainting) 去除文字
        inpainted = cv2.inpaint(image, mask, 3, cv2.INPAINT_TELEA)
        
        # 編碼回 PNG bytes
        is_success, encoded_img = cv2.imencode(".png", inpainted)
        if is_success:
            return encoded_img.tobytes()
    except Exception as e:
        print(f"[-] Inpainting 失敗: {e}")
    return img_bytes

class OCRTestApp:
    def __init__(self, root):
        self.root = root
        self.root.title("EasyOCR 圖片與 PDF 文字測試工具")
        self.root.geometry("800x650")
        
        # 狀態列
        self.label_status = tk.Label(root, text="正在載入 OCR 模型 (繁中、英文)...", fg="orange", font=("Microsoft JhengHei", 12, "bold"))
        self.label_status.pack(pady=15)
        self.root.update()
        
        # 載入 EasyOCR
        try:
            self.reader = easyocr.Reader(['ch_tra', 'en'], gpu=False)
            self.label_status.config(text="✨ OCR 模型載入成功！請選擇圖片或 PDF 測試", fg="green")
        except Exception as e:
            self.label_status.config(text=f"❌ 模型載入失敗: {e}", fg="red")
            messagebox.showerror("錯誤", f"無法載入 EasyOCR: {e}")
        
        # 操作按鈕區
        frame_top = ttk.Frame(root, padding=10)
        frame_top.pack(fill="x")
        
        self.btn_select = ttk.Button(frame_top, text="1. 選擇測試檔案 (圖片/PDF)", command=self.select_file)
        self.btn_select.pack(side="left", padx=10, ipady=4)
        
        self.btn_run = ttk.Button(frame_top, text="2. 開始辨識並存為 Word", command=self.run_ocr_word, state="disabled")
        self.btn_run.pack(side="left", padx=10, ipady=4)
        
        self.btn_run_ppt = ttk.Button(frame_top, text="3. 開始辨識並存為 PPT (無字底圖版)", command=self.run_ocr_ppt, state="disabled")
        self.btn_run_ppt.pack(side="left", padx=10, ipady=4)
        
        self.label_path = ttk.Label(root, text="尚未選取任何圖片或 PDF 檔案...", wraplength=750, font=("Microsoft JhengHei", 9))
        self.label_path.pack(pady=10)
        
        # 辨識結果顯示區
        frame_bottom = ttk.LabelFrame(root, text=" 辨識結果輸出區 (已同步存為同資料夾對應格式之檔案) ", padding=10)
        frame_bottom.pack(fill="both", expand=True, padx=20, pady=15)
        
        self.text_output = tk.Text(frame_bottom, font=("Consolas", 11), wrap="word")
        self.text_output.pack(side="left", fill="both", expand=True)
        
        scrollbar = ttk.Scrollbar(frame_bottom, command=self.text_output.yview)
        scrollbar.pack(side="right", fill="y")
        self.text_output.config(yscrollcommand=scrollbar.set)
        
        self.image_path = None

    def select_file(self):
        path = filedialog.askopenfilename(filetypes=[("支援的格式", "*.jpg;*.jpeg;*.png;*.bmp;*.pdf")])
        if path:
            self.image_path = path
            self.label_path.config(text=f"已選擇檔案：{path}")
            self.btn_run.config(state="normal")
            self.btn_run_ppt.config(state="normal")
            self.text_output.delete("1.0", tk.END)
            self.label_status.config(text="檔案已選取，請選擇要匯出的格式開始辨識", fg="blue")

    def toggle_buttons(self, state):
        self.btn_run.config(state=state)
        self.btn_run_ppt.config(state=state)
        self.btn_select.config(state=state)

    def run_ocr_word(self):
        if not self.image_path:
            return
        
        self.toggle_buttons("disabled")
        self.text_output.delete("1.0", tk.END)
        self.root.update()
        
        is_pdf = self.image_path.lower().endswith('.pdf')
        results_all = []
        
        try:
            doc_word = Document()
            
            if is_pdf:
                self.label_status.config(text="⏳ 正在開啟 PDF 檔案...", fg="blue")
                self.root.update()
                
                pdf_doc = fitz.open(self.image_path)
                total_pages = len(pdf_doc)
                doc_word.add_heading(f"PDF OCR 辨識結果 - {os.path.basename(self.image_path)}", level=1)
                
                for page_num in range(total_pages):
                    self.label_status.config(text=f"⏳ 正在辨識 PDF 第 {page_num + 1}/{total_pages} 頁...", fg="blue")
                    self.root.update()
                    
                    page = pdf_doc[page_num]
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("png")
                    
                    results = self.reader.readtext(img_bytes, detail=0)
                    results_all.extend(results)
                    
                    doc_word.add_heading(f"--- 第 {page_num + 1} 頁 ---", level=2)
                    for line in results:
                        doc_word.add_paragraph(line)
                    
                    if page_num < total_pages - 1:
                        doc_word.add_page_break()
            else:
                self.label_status.config(text="⏳ 正在辨識單張圖片...", fg="blue")
                self.root.update()
                
                with open(self.image_path, "rb") as f:
                    img_bytes = f.read()
                results = self.reader.readtext(img_bytes, detail=0)
                results_all.extend(results)
                
                doc_word.add_heading("OCR 辨識結果測試", level=1)
                doc_word.add_paragraph(f"來源圖片：{self.image_path}")
                doc_word.add_paragraph(f"辨識時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                doc_word.add_paragraph("=" * 50)
                for line in results:
                    doc_word.add_paragraph(line)
            
            # 將辨識純文字輸出至 GUI 畫面
            full_text = "\n".join(results_all)
            self.text_output.insert(tk.END, full_text)
            
            # 儲存 Word 檔
            base_path, _ = os.path.splitext(self.image_path)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            docx_path = f"{base_path}_ocr_result_{timestamp}.docx"
            doc_word.save(docx_path)
            
            self.label_status.config(text="✅ 辨識完成！已存成 Word 並為您開啟中...", fg="green")
            os.startfile(docx_path)
            
        except Exception as e:
            messagebox.showerror("錯誤", f"辨識或寫入 Word 失敗：{e}")
            self.label_status.config(text="❌ 辨識或存檔失敗", fg="red")
        finally:
            self.toggle_buttons("normal")

    def run_ocr_ppt(self):
        if not self.image_path:
            return
            
        self.toggle_buttons("disabled")
        self.text_output.delete("1.0", tk.END)
        self.root.update()
        
        is_pdf = self.image_path.lower().endswith('.pdf')
        results_all_text = []
        
        try:
            prs = Presentation()
            # 設定寬螢幕 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]
            
            slide_w_in = 13.333
            slide_h_in = 7.5
            
            if is_pdf:
                self.label_status.config(text="⏳ 正在開啟 PDF 檔案...", fg="blue")
                self.root.update()
                
                pdf_doc = fitz.open(self.image_path)
                total_pages = len(pdf_doc)
                
                for page_num in range(total_pages):
                    self.label_status.config(text=f"⏳ 正在處理 PDF 第 {page_num + 1}/{total_pages} 頁...", fg="blue")
                    self.root.update()
                    
                    page = pdf_doc[page_num]
                    pix = page.get_pixmap(dpi=150)
                    img_width = pix.width
                    img_height = pix.height
                    img_bytes = pix.tobytes("png")
                    
                    # 辨識文字與座標
                    results = self.reader.readtext(img_bytes, detail=1)
                    results_all_text.extend([item[1] for item in results])
                    
                    # 進行 Inpainting 修補，抹除原圖上的硬編碼文字
                    self.label_status.config(text=f"⏳ 正在修補抹除 PDF 第 {page_num + 1} 頁的背景文字...", fg="blue")
                    self.root.update()
                    cleaned_img_bytes = remove_text_from_image_bytes(img_bytes, results)
                    
                    # 轉換為 PIL Image 用於計算文字對比顏色
                    pil_img = Image.open(io.BytesIO(img_bytes))
                    
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # 鋪設無文字修補後的底圖
                    img_stream = io.BytesIO(cleaned_img_bytes)
                    slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                    
                    for box, text, conf in results:
                        left_px = min(p[0] for p in box)
                        top_px = min(p[1] for p in box)
                        right_px = max(p[0] for p in box)
                        bottom_px = max(p[1] for p in box)
                        
                        width_px = right_px - left_px
                        height_px = bottom_px - top_px
                        
                        left_in = (left_px / img_width) * slide_w_in
                        top_in = (top_px / img_height) * slide_h_in
                        width_in = max(0.05, (width_px / img_width) * slide_w_in)
                        height_in = max(0.05, (height_px / img_height) * slide_h_in)
                        
                        # 新增文字框 (背景維持透明，無遮擋框)
                        txBox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        tf.margin_left = Inches(0)
                        tf.margin_right = Inches(0)
                        tf.margin_top = Inches(0)
                        tf.margin_bottom = Inches(0)
                        
                        p = tf.paragraphs[0]
                        p.text = text
                        p.font.name = 'Microsoft JhengHei'
                        p.font.size = Pt(11)
                        
                        # 偵測背景明暗，給予適合的字型顏色
                        bg_color = get_background_color(pil_img, box)
                        brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
                        if brightness < 128:
                            p.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            p.font.color.rgb = RGBColor(30, 41, 59)
            else:
                self.label_status.config(text="⏳ 正在辨識單張圖片...", fg="blue")
                self.root.update()
                
                pil_img = Image.open(self.image_path)
                img_width, img_height = pil_img.size
                
                with open(self.image_path, "rb") as f:
                    img_bytes = f.read()
                    
                results = self.reader.readtext(img_bytes, detail=1)
                results_all_text.extend([item[1] for item in results])
                
                # 影像修補，抹除文字
                self.label_status.config(text="⏳ 正在修補抹除底圖文字...", fg="blue")
                self.root.update()
                cleaned_img_bytes = remove_text_from_image_bytes(img_bytes, results)
                
                slide = prs.slides.add_slide(blank_layout)
                
                # 鋪設無字底圖
                img_stream = io.BytesIO(cleaned_img_bytes)
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                
                for box, text, conf in results:
                    left_px = min(p[0] for p in box)
                    top_px = min(p[1] for p in box)
                    right_px = max(p[0] for p in box)
                    bottom_px = max(p[1] for p in box)
                    
                    width_px = right_px - left_px
                    height_px = bottom_px - top_px
                    
                    left_in = (left_px / img_width) * slide_w_in
                    top_in = (top_px / img_height) * slide_h_in
                    width_in = max(0.05, (width_px / img_width) * slide_w_in)
                    height_in = max(0.05, (height_px / img_height) * slide_h_in)
                    
                    txBox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0)
                    tf.margin_right = Inches(0)
                    tf.margin_top = Inches(0)
                    tf.margin_bottom = Inches(0)
                    
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(12)
                    
                    bg_color = get_background_color(pil_img, box)
                    brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
                    if brightness < 128:
                        p.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        p.font.color.rgb = RGBColor(30, 41, 59)
            
            # 將辨識文字輸出至介面
            full_text = "\n".join(results_all_text)
            self.text_output.insert(tk.END, full_text)
            
            # 儲存 PPTX
            base_path, _ = os.path.splitext(self.image_path)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            ppt_path = f"{base_path}_ocr_result_{timestamp}.pptx"
            prs.save(ppt_path)
            
            self.label_status.config(text="✅ 辨識與修補完成！已存成 PPT 並開啟中...", fg="green")
            os.startfile(ppt_path)
            
        except Exception as e:
            messagebox.showerror("錯誤", f"辨識或寫入 PPT 失敗：{e}")
            self.label_status.config(text="❌ 辨識或存檔失敗", fg="red")
        finally:
            self.toggle_buttons("normal")

if __name__ == "__main__":
    root = tk.Tk()
    app = OCRTestApp(root)
    root.mainloop()
