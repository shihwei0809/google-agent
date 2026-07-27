"""
Patch script: rewrites the export_pptx function in app.py
Run once with: venv\Scripts\python.exe patch_pptx.py
"""
import pathlib

app_path = pathlib.Path("app.py")
content = app_path.read_text(encoding="utf-8")
lines = content.split("\n")

# Find start (line with @app.post("/api/export-pptx")) and end
start_idx = None
end_idx = None
for i, line in enumerate(lines):
    if '@app.post("/api/export-pptx")' in line:
        start_idx = i
    if start_idx is not None and i > start_idx and line.startswith("@app.post("):
        end_idx = i
        break

print(f"Found export_pptx at lines {start_idx+1}–{end_idx} (0-indexed: {start_idx}–{end_idx-1})")

NEW_FUNC = r'''@app.post("/api/export-pptx")
async def export_pptx(job_id: str = Form(...)):
    import time, re
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt, Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    job_dir = JOBS_DIR / job_id
    with open(job_dir / "script.json", "r", encoding="utf-8") as f:
        job_data = json.load(f)

    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    W = prs.slide_width
    H = prs.slide_height

    C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
    C_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    C_BG     = RGBColor(0xF1, 0xF5, 0xF9)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
    C_MUTED  = RGBColor(0x64, 0x74, 0x8B)

    def hex_rgb(c):
        return f"{c.red:02X}{c.green:02X}{c.blue:02X}"

    def force_fill(shape, color):
        """直接操作 XML 注入 solidFill，繞過 theme 繼承"""
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(shape._element, qn("p:spPr"))
        for tag in [qn("a:noFill"), qn("a:gradFill"), qn("a:solidFill"), qn("a:pattFill")]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        sr = etree.SubElement(sf, qn("a:srgbClr"))
        sr.set("val", hex_rgb(color))

    def add_rect(sl, l, t, w, h, color):
        shp = sl.shapes.add_shape(1, int(l), int(t), int(w), int(h))
        force_fill(shp, color)
        shp.line.fill.background()
        return shp

    def add_text(sl, l, t, w, h, text, size=20, bold=False, color=None,
                 align=PP_ALIGN.LEFT, line_sp=1.4):
        if color is None:
            color = C_TEXT
        tb = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in text.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spc   = etree.SubElement(lnSpc, qn("a:spcPct"))
            spc.set("val", str(int(line_sp * 100000)))
            r = p.add_run()
            r.text = line
            r.font.name  = "\u5fae\u8edf\u6b63\u9ed1\u9ad4"
            r.font.size  = Pt(size)
            r.font.bold  = bold
            r.font.color.rgb = color
        return tb

    blank = prs.slide_layouts[6]  # Blank

    # ─── 封面頁 ───────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = C_DARK
    add_rect(cover, 0, H * 0.60, W, H * 0.40, C_BG)
    add_rect(cover, 0, H * 0.10, Emu(14000), H * 0.45, C_ACCENT)
    title_str = job_data.get("title", job_id)
    add_text(cover, Emu(100000), H * 0.12, W * 0.92, H * 0.42,
             title_str, size=38, bold=True, color=C_WHITE, line_sp=1.4)
    add_text(cover, Emu(100000), H * 0.67, W * 0.60, H * 0.18,
             "\u54e1\u5de5\u6559\u80b2\u8a13\u7df4\u6559\u6750", size=17, color=C_MUTED)

    # ─── 內容頁 ───────────────────────────────────────────
    for slide_data in job_data.get("slides", []):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = C_BG
        add_rect(sl, 0, 0, W, H * 0.19, C_DARK)
        add_rect(sl, 0, H * 0.19, Emu(13000), H * 0.81, C_ACCENT)

        t_text = coerce_str(slide_data.get("title", ""))
        add_text(sl, Emu(70000), H * 0.02, W * 0.94, H * 0.16,
                 t_text, size=30, bold=True, color=C_WHITE, line_sp=1.2)

        raw = coerce_str(slide_data.get("content", ""))
        bullets = [s.strip() for s in re.split(r"[\u3002\uff1b;\n]", raw) if s.strip()]
        body = "\n".join(f"\u2022  {b}" for b in bullets) if bullets else raw
        add_text(sl, Emu(80000), H * 0.23, W * 0.94, H * 0.72,
                 body, size=19, color=C_TEXT, line_sp=1.7)

        sl.notes_slide.notes_text_frame.text = coerce_str(slide_data.get("narration", ""))

    out_path = job_dir / f"{job_id}.pptx"
    prs.save(str(out_path))
    return {"url": f"/jobs/{job_id}/{job_id}.pptx?t={int(time.time())}"}

'''

new_lines = lines[:start_idx] + NEW_FUNC.split("\n") + lines[end_idx:]
app_path.write_text("\n".join(new_lines), encoding="utf-8")
print(f"Done! New line count: {len(new_lines)}")
