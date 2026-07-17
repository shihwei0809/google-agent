import asyncio
import base64
import json
import logging
import os
import shutil
import uuid
import wave
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

try:
    from dotenv import load_dotenv
    load_dotenv(override=False)
except ImportError:
    pass

import edge_tts
import fitz  # PyMuPDF
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from PIL import Image

# For PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
import docx

# ─── Config ───────────────────────────────────────────────────────────────────
JOBS_DIR = Path("jobs")
JOBS_DIR.mkdir(exist_ok=True)
TEMPLATES_DIR = Path("templates")
TEMPLATES_DIR.mkdir(exist_ok=True)

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

executor = ThreadPoolExecutor(max_workers=2)

# ─── Voice Configurations ───────────────────────────────────────────────────────
VOICE_GROUPS = {
    "🇹🇼 繁體中文（台灣）": [
        {"id": "zh-TW-HsiaoChenNeural", "label": "曉臻（女聲・自然）"},
        {"id": "zh-TW-HsiaoYuNeural",   "label": "曉雨（女聲・活潑）"},
        {"id": "zh-TW-YunJheNeural",    "label": "雲哲（男聲）"},
    ],
    "🇨🇳 普通話（大陸）": [
        {"id": "zh-CN-XiaoxiaoNeural",  "label": "曉曉（女聲）"},
        {"id": "zh-CN-YunxiNeural",     "label": "雲希（男聲）"},
    ],
    "🇺🇸 英文（美國）": [
        {"id": "en-US-JennyNeural",  "label": "Jenny（Female）"},
        {"id": "en-US-GuyNeural",    "label": "Guy（Male）"},
    ]
}

GEMINI_TTS_VOICES = [
    {"id": "Zephyr", "label": "Zephyr"},
    {"id": "Puck",   "label": "Puck"},
    {"id": "Charon", "label": "Charon"},
    {"id": "Kore",   "label": "Kore"}
]

# ─── Helpers ──────────────────────────────────────────────────────────────────
def parse_api_keys(raw_keys: str) -> list[str]:
    if not raw_keys:
        return []
    keys = []
    for line in raw_keys.replace(",", "\n").splitlines():
        k = line.strip()
        if k:
            keys.append(k)
    return keys

# Gemini Integration
import requests

def call_gemini_vision(api_key: str, model: str, prompt: str, image_path: str = None, text_content: str = None) -> str:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    
    parts = [{"text": prompt}]
    if image_path:
        with open(image_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode("utf-8")
        parts.append({
            "inlineData": {
                "mimeType": "image/png",
                "data": img_b64
            }
        })
    if text_content:
        parts.append({"text": f"\n\n參考內容：\n{text_content}"})
        
    payload = {"contents": [{"parts": parts}]}
    
    response = requests.post(url, json=payload, headers=headers, timeout=60)
    response.raise_for_status()
    data = response.json()
    return data["candidates"][0]["content"]["parts"][0]["text"].strip()

def pcm_to_wav(pcm_data: bytes, sample_rate: int = 24000, channels: int = 1, sample_width: int = 2) -> bytes:
    import io
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(channels)
        wf.setsampwidth(sample_width)
        wf.setframerate(sample_rate)
        wf.writeframes(pcm_data)
    return buf.getvalue()

def call_gemini_tts(api_key: str, model: str, voice: str, text: str, output_path: str) -> None:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    payload = {
        "contents": [{"parts": [{"text": text}]}],
        "generationConfig": {
            "responseModalities": ["AUDIO"],
            "speechConfig": {
                "voiceConfig": {
                    "prebuiltVoiceConfig": {
                        "voiceName": voice
                    }
                }
            }
        }
    }
    response = requests.post(url, json=payload, headers=headers, timeout=120)
    response.raise_for_status()
    data = response.json()

    audio_b64 = None
    mime_type = ""
    for cand in data.get("candidates", []):
        content = cand.get("content", {})
        if isinstance(content, dict):
            for part in content.get("parts", []):
                inline = part.get("inlineData", {})
                if inline.get("mimeType", "").startswith("audio/"):
                    audio_b64 = inline["data"]
                    mime_type = inline["mimeType"]
                    break
    if not audio_b64:
        raise ValueError(f"Gemini TTS: no audio data in response. Raw: {str(data)[:300]}")

    raw_bytes = base64.b64decode(audio_b64)
    # Gemini 回傳 audio/L16 (raw PCM) 時需要包成 WAV
    if "L16" in mime_type or "pcm" in mime_type.lower():
        raw_bytes = pcm_to_wav(raw_bytes)
    with open(output_path, "wb") as f:
        f.write(raw_bytes)

# ─── FastAPI Application ──────────────────────────────────────────────────────
app = FastAPI(title="SOP Interactive Builder")
templates = Jinja2Templates(directory="templates")
app.mount("/jobs", StaticFiles(directory="jobs"), name="jobs")

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse(request=request, name="index.html")

@app.get("/api/voices")
async def get_voices():
    return {"edge": VOICE_GROUPS, "gemini": GEMINI_TTS_VOICES}

@app.get("/api/jobs")
async def list_jobs():
    """載入歷史專案清單"""
    jobs = []
    for d in JOBS_DIR.iterdir():
        if d.is_dir() and (d / "script.json").exists():
            jobs.append(d.name)
    return {"jobs": sorted(jobs, reverse=True)}

@app.get("/api/jobs/{job_id}")
async def get_job_data(job_id: str):
    job_dir = JOBS_DIR / job_id
    script_path = job_dir / "script.json"
    if not script_path.exists():
        raise HTTPException(status_code=404, detail="找不到該專案")
    with open(script_path, "r", encoding="utf-8") as f:
        return json.load(f)

@app.post("/api/extract")
async def extract_document(
    file: UploadFile = File(...),
    api_keys: str = Form(""),
    model: str = Form("gemini-3.5-flash"),
    extract_mode: str = Form("gemini")
):
    """解析 PDF / DOCX 並切分為投影片腳本"""
    import datetime, re
    safe_name = re.sub(r'[^\w\u4e00-\u9fa5\-]', '_', Path(file.filename).stem)
    job_id = f"{safe_name}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    job_dir = JOBS_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    
    keys = parse_api_keys(api_keys)
    env_keys = parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
    keys = keys or env_keys
    
    if extract_mode == "gemini" and not keys:
        extract_mode = "local"  # 自動退回到本機模式
    
    content_text = ""
    file_bytes = await file.read()
    
    local_slides = []
    
    if file.filename.lower().endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            content_text += f"\n--- 第 {i+1} 頁 ---\n" + text
            
            # 準備本機模式的腳本
            if text:
                local_slides.append({
                    "title": f"投影片 {i+1}",
                    "content": text[:150] + ("..." if len(text)>150 else ""),
                    "narration": text.replace("\n", "，")[:300]
                })
                
            # 同時產出圖片做備用/參考
            pix = page.get_pixmap()
            pix.save(str(job_dir / f"page_{i}.png"))
    elif file.filename.lower().endswith(".docx"):
        import io
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for i, p_text in enumerate(paragraphs):
            content_text += p_text + "\n"
            local_slides.append({
                "title": f"段落 {i+1}",
                "content": p_text[:150] + ("..." if len(p_text)>150 else ""),
                "narration": p_text.replace("\n", "，")[:300]
            })
    else:
        raise HTTPException(status_code=400, detail="僅支援 PDF 或 DOCX")

    if not local_slides:
        local_slides.append({
            "title": "無內容",
            "content": "無法解析出文字",
            "narration": "系統無法從此檔案解析出文字"
        })

    if extract_mode == "local":
        parsed_data = {"slides": local_slides, "quiz": []}
    else:
        # 使用 Gemini 分析並輸出 JSON 格式的大綱與旁白，以及測驗題目
        prompt = (
            "請將以下教育訓練教材內容整理成多頁的『簡報腳本』，並根據內容生成 3 到 5 題『選擇題測驗』。\n"
            "請以嚴格的 JSON 格式回傳，包含兩個主要屬性：\n"
            '1. "slides": 陣列，每個元素代表一頁簡報，包含:\n'
            '   - "title": 投影片標題\n'
            '   - "content": 投影片重點內容（條列式）\n'
            '   - "narration": 口語化的朗讀旁白腳本（適合語音播放）\n'
            '2. "quiz": 陣列，每個元素代表一題選擇題，包含:\n'
            '   - "question": 題目敘述\n'
            '   - "options": 長度為 4 的字串陣列 (選項 A, B, C, D)\n'
            '   - "answer": 正確選項的索引值 (0 到 3)\n'
            "只回傳 JSON，不要加 markdown 標記。"
        )
    
        # Model fallback 清單：選到的 model 排第一，失敗自動往下試
        FALLBACK_MODELS = [
            "gemini-3.5-flash",
            "gemini-2.5-flash",
            "gemini-2.0-flash",
            "gemini-2.0-flash-lite",
        ]
        if model in FALLBACK_MODELS:
            model_queue = [model] + [m for m in FALLBACK_MODELS if m != model]
        else:
            model_queue = [model] + FALLBACK_MODELS

        success, result_text, last_err = False, "", None
        for try_model in model_queue:
            for key in keys:
                try:
                    result_text = call_gemini_vision(key, try_model, prompt, text_content=content_text)
                    logger.info(f"Gemini 成功：model={try_model}")
                    success = True
                    break
                except Exception as e:
                    last_err = e
                    logger.warning(f"Gemini 失敗：model={try_model} err={e}")
            if success:
                break

        if not success:
            raise HTTPException(status_code=500, detail=f"AI 分析失敗（已試所有 model 與 key）: {last_err}")
        
        try:
            if result_text.startswith("```json"):
                result_text = result_text.strip("`").replace("json\n", "")
            parsed_data = json.loads(result_text)
            
            if isinstance(parsed_data, list):
                parsed_data = {"slides": parsed_data, "quiz": []}
            elif "slides" not in parsed_data:
                parsed_data = {"slides": local_slides, "quiz": []}
                
            slides = parsed_data.get("slides", [])
            quiz = parsed_data.get("quiz", [])
            
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"解析 AI JSON 失敗: {str(e)}\n\n{result_text}")

    slides = parsed_data.get("slides", [])
    quiz   = parsed_data.get("quiz", [])

    job_data = {
        "job_id": job_id,
        "slides": slides,
        "quiz": quiz,
        "voice": "zh-TW-HsiaoChenNeural",
        "tts_engine": "edge"
    }
    with open(job_dir / "script.json", "w", encoding="utf-8") as f:
        json.dump(job_data, f, ensure_ascii=False, indent=2)
        
    return {"job_id": job_id, "slides": slides, "quiz": quiz}

@app.post("/api/generate-quiz")
async def generate_quiz_endpoint(
    api_keys: str = Form(""),
    model: str = Form("gemini-3.5-flash"),
    content: str = Form(...),
    count: int = Form(5),
):
    """根據腳本內容 AI 生成考題"""
    keys = parse_api_keys(api_keys) or parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
    if not keys:
        raise HTTPException(status_code=400, detail="需要 Gemini API 金鑰才能 AI 生成考題")

    prompt = (
        f"請根據以下教育訓練腳本內容，生成 {count} 題繁體中文選擇題測驗。\n"
        "每題必須包含 4 個選項（A/B/C/D）且只有一個正確答案。\n"
        "以 JSON 陣列格式回傳，每個元素包含：\n"
        '- "question": 題目敘述\n'
        '- "options": 長度為 4 的字串陣列\n'
        '- "answer": 正確選項索引（0=A, 1=B, 2=C, 3=D）\n'
        "只回傳 JSON 陣列，不要加 markdown 標記。\n\n"
        f"腳本內容：\n{content[:6000]}"
    )

    FALLBACK_MODELS = [
        "gemini-3.5-flash", "gemini-2.5-flash",
        "gemini-2.0-flash", "gemini-2.0-flash-lite",
    ]
    model_queue = [model] + [m for m in FALLBACK_MODELS if m != model]

    result_text, last_err, success = "", None, False
    for try_model in model_queue:
        for key in keys:
            try:
                result_text = call_gemini_vision(key, try_model, prompt)
                logger.info(f"generate-quiz 成功：model={try_model}")
                success = True
                break
            except Exception as e:
                last_err = e
                logger.warning(f"generate-quiz 失敗：model={try_model} err={e}")
        if success:
            break

    if not success:
        raise HTTPException(status_code=500, detail=f"AI 生成考題失敗: {last_err}")

    try:
        if result_text.startswith("```"):
            result_text = result_text.strip("`").replace("json\n", "").strip()
        quiz = json.loads(result_text)
        if not isinstance(quiz, list):
            quiz = quiz.get("quiz", [])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"解析考題 JSON 失敗: {e}\n{result_text[:300]}")

    return {"quiz": quiz}

@app.post("/api/generate-audio")
async def generate_audio(
    job_id: str = Form(...),
    page_index: int = Form(...),
    text: str = Form(...),
    tts_engine: str = Form("edge"),
    voice: str = Form(...),
    api_keys: str = Form("")
):
    """生成單頁語音"""
    job_dir = JOBS_DIR / job_id
    if not job_dir.exists():
        raise HTTPException(status_code=404, detail="找不到專案")
        
    audio_ext = "wav" if tts_engine == "gemini" else "mp3"
    audio_path = job_dir / f"audio_{page_index}.{audio_ext}"
    
    try:
        if tts_engine == "gemini":
            keys = parse_api_keys(api_keys) or parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
            if not keys:
                raise HTTPException(status_code=400, detail="Gemini TTS 需要 API 金鑰")
            success = False
            for key in keys:
                try:
                    call_gemini_tts(key, "gemini-2.5-flash-preview-tts", voice, text, str(audio_path))
                    success = True
                    break
                except:
                    continue
            if not success:
                raise ValueError("所有金鑰皆失效")
        else:
            communicate = edge_tts.Communicate(text, voice)
            await communicate.save(str(audio_path))
            
        import time
        return {"url": f"/jobs/{job_id}/audio_{page_index}.{audio_ext}?t={int(time.time())}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/save-script")
async def save_script(request: Request):
    data = await request.json()
    job_id = data.get("job_id")
    job_dir = JOBS_DIR / job_id
    if job_dir.exists():
        with open(job_dir / "script.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    return {"status": "ok"}

@app.post("/api/export-pptx")
async def export_pptx(job_id: str = Form(...)):
    job_dir = JOBS_DIR / job_id
    script_path = job_dir / "script.json"
    with open(script_path, "r", encoding="utf-8") as f:
        job_data = json.load(f)
        
    prs = Presentation()
    for slide_data in job_data.get("slides", []):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data.get("title", "")
        content = slide.placeholders[1]
        content.text = slide_data.get("content", "")
        
        # Add narration to notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("narration", "")
        
    out_path = job_dir / f"{job_id}.pptx"
    prs.save(str(out_path))
    return {"url": f"/jobs/{job_id}/{job_id}.pptx"}

@app.post("/api/export-web")
async def export_web(job_id: str = Form(...)):
    job_dir = JOBS_DIR / job_id
    with open(job_dir / "script.json", "r", encoding="utf-8") as f:
        job_data = json.load(f)
        
    template_path = TEMPLATES_DIR / "player_template.html"
    if not template_path.exists():
        raise HTTPException(status_code=500, detail="Player template missing")
        
    with open(template_path, "r", encoding="utf-8") as f:
        html = f.read()
        
    # Inject slides and quiz config
    slides_config = []
    for i, s in enumerate(job_data.get("slides", [])):
        audio_ext = "wav" if job_data.get("tts_engine") == "gemini" else "mp3"
        audio_file = f"audio_{i}.{audio_ext}"
        slides_config.append({
            "title": s.get("title", ""),
            "content": s.get("content", ""),
            "audio": audio_file if (job_dir / audio_file).exists() else None
        })
        
    html = html.replace("{{SLIDES_JSON}}", json.dumps(slides_config, ensure_ascii=False))
    html = html.replace("{{QUIZ_JSON}}", json.dumps(job_data.get("quiz", []), ensure_ascii=False))
    
    out_path = job_dir / "index.html"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)
        
    return {"url": f"/jobs/{job_id}/index.html"}
