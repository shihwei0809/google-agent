import sys
import os
import re
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color Palette Definitions
DARK_BG = RGBColor(15, 23, 42)      # Slate Navy
DARK_TITLE = RGBColor(255, 255, 255) # White
DARK_ACCENT = RGBColor(14, 165, 233) # Sky Blue

LIGHT_BG = RGBColor(248, 250, 252)  # Light Slate Gray
LIGHT_TITLE = RGBColor(15, 23, 42)   # Navy
LIGHT_BODY = RGBColor(51, 65, 85)    # Slate Gray
LIGHT_ACCENT = RGBColor(37, 99, 235)  # Royal Blue
LIGHT_ACCENT_LINE = RGBColor(226, 232, 240) # Slate 200 (for footer line)
MUTED_TEXT = RGBColor(148, 163, 184) # Slate 400 (for footer text)

FONT_NAME = "Microsoft JhengHei"

def clean_markdown(text):
    """
    Cleans markdown formatting syntax for plain text display in PPTX.
    """
    # Remove bold markers
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)
    # Remove italic markers
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'_(.*?)_', r'\1', text)
    # Remove inline code markers
    text = re.sub(r'`(.*?)`', r'\1', text)
    # Remove markdown link syntax [label](url)
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    # Clean up escaped dollar signs and double hyphens
    text = text.replace(r'\$', '$')
    text = text.replace('--', '～')
    return text.strip()

def is_separator_line(line):
    """
    Detects if a line is a table markdown separator (e.g. |---| or ---- ----).
    """
    stripped = line.strip()
    if not stripped:
        return False
    if stripped.count('-') < 3 and stripped.count('=') < 3:
        return False
    return all(c in '-=|: \t' for c in stripped)

def get_list_level(indent, indent_stack):
    """
    Tracks and computes nested list levels based on leading indentation spaces.
    """
    if not indent_stack:
        indent_stack.append(indent)
        return 0
    if indent > indent_stack[-1]:
        indent_stack.append(indent)
        return len(indent_stack) - 1
    while indent_stack and indent < indent_stack[-1]:
        indent_stack.pop()
    if not indent_stack or indent > indent_stack[-1]:
        indent_stack.append(indent)
    return len(indent_stack) - 1

def parse_slide_content(lines):
    """
    Parses the raw lines of a slide body into structured blocks (lists, tables, paragraphs, images).
    """
    blocks = []
    i = 0
    n = len(lines)
    list_indent_stack = []
    
    while i < n:
        line = lines[i]
        stripped = line.strip()
        
        # Skip empty lines, and reset the list indentation stack
        if not stripped:
            list_indent_stack = []
            i += 1
            continue
            
        # 1. Image Detection
        image_match = re.match(r'^!\[(.*?)\]\((.*?)\)', stripped)
        if image_match:
            alt, img_path = image_match.groups()
            # Strip file:/// prefixes for local paths on Windows
            if img_path.startswith('file:///'):
                img_path = img_path[8:]
                img_path = img_path.replace('/', '\\')
                from urllib.parse import unquote
                img_path = unquote(img_path)
            blocks.append({
                "type": "image",
                "alt": alt,
                "path": img_path
            })
            i += 1
            list_indent_stack = []
            continue
            
        # 2. Table Detection (GFM or pandoc space-separated tables)
        if i + 1 < n and is_separator_line(lines[i+1]):
            header_line = line
            sep_line = lines[i+1]
            
            table_rows = []
            j = i + 2
            while j < n:
                row_stripped = lines[j].strip()
                if not row_stripped:
                    break
                if is_separator_line(lines[j]):
                    j += 1
                    continue
                # Stop table parsing if encountering a header, list item or image
                if row_stripped.startswith('#') or row_stripped.startswith('- ') or row_stripped.startswith('* ') or re.match(r'^\d+\.\s+', row_stripped) or row_stripped.startswith('!['):
                    break
                table_rows.append(lines[j])
                j += 1
                
            has_pipe = '|' in sep_line
            if has_pipe:
                def split_pipe_row(r):
                    parts = r.split('|')
                    if r.strip().startswith('|'):
                        parts = parts[1:]
                    if r.strip().endswith('|'):
                        parts = parts[:-1]
                    return [clean_markdown(c) for c in parts]
                
                headers = split_pipe_row(header_line)
                rows = [split_pipe_row(r) for r in table_rows]
            else:
                # Split by 2 or more spaces
                headers = [clean_markdown(c) for c in re.split(r'\s{2,}', header_line.strip())]
                rows = [[clean_markdown(c) for c in re.split(r'\s{2,}', r.strip())] for r in table_rows]
            
            blocks.append({
                "type": "table",
                "headers": headers,
                "rows": rows
            })
            i = j
            list_indent_stack = []
            continue
            
        # 3. List Item Detection (bullet or numbered)
        list_match = re.match(r'^(\s*)([-*+])\s+(.*)', line)
        num_list_match = re.match(r'^(\s*)(\d+)\.\s+(.*)', line)
        
        if list_match or num_list_match:
            if list_match:
                spaces, marker, text = list_match.groups()
            else:
                spaces, marker, text = num_list_match.groups()
                marker = marker + "."
                
            indent = len(spaces)
            level = get_list_level(indent, list_indent_stack)
            item = {"level": level, "marker": marker, "text": clean_markdown(text)}
            
            if blocks and blocks[-1]["type"] == "list":
                blocks[-1]["items"].append(item)
            else:
                blocks.append({
                    "type": "list",
                    "items": [item]
                })
            i += 1
            continue
            
        # 4. Regular Paragraph or Subheading (Level 3 Header)
        if stripped.startswith('###'):
            subhead_text = clean_markdown(stripped.lstrip('#'))
            blocks.append({
                "type": "paragraph",
                "text": subhead_text,
                "is_subheading": True
            })
        else:
            blocks.append({
                "type": "paragraph",
                "text": clean_markdown(line),
                "is_subheading": False
            })
        list_indent_stack = []
        i += 1
        
    return blocks

def parse_markdown_to_slides(filepath):
    """
    Reads a Markdown file and groups content by headers into individual slide data.
    """
    with open(filepath, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    slides = []
    current_slide = None
    
    for line in lines:
        stripped = line.strip()
        
        # Level 1 Header -> Section / Title Slide
        if line.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            title = line[2:].strip()
            current_slide = {
                "type": "title",
                "title": clean_markdown(title),
                "lines": []
            }
        # Level 2 Header -> Content Slide
        elif line.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            title = line[3:].strip()
            current_slide = {
                "type": "content",
                "title": clean_markdown(title),
                "lines": []
            }
        else:
            if current_slide:
                current_slide["lines"].append(line)
            else:
                if stripped:
                    # Prepend an initial title slide if content starts without header
                    current_slide = {
                        "type": "title",
                        "title": clean_markdown(stripped),
                        "lines": []
                    }
                    
    if current_slide:
        slides.append(current_slide)
        
    # Parse lines of each slide into structured blocks
    for slide in slides:
        slide["blocks"] = parse_slide_content(slide["lines"])
        del slide["lines"]
        
    return slides

def get_font_sizes(blocks):
    """
    Adjusts font sizes dynamically to avoid slide text overflow.
    """
    total_len = 0
    for block in blocks:
        if block["type"] == "list":
            for item in block["items"]:
                total_len += len(item["text"])
        elif block["type"] == "paragraph":
            total_len += len(block["text"])
            
    if total_len > 500:
        return {"para": Pt(14), "l0": Pt(14), "l1": Pt(12), "l2": Pt(10), "sub": Pt(18)}
    elif total_len > 300:
        return {"para": Pt(16), "l0": Pt(16), "l1": Pt(14), "l2": Pt(12), "sub": Pt(20)}
    else:
        return {"para": Pt(18), "l0": Pt(18), "l1": Pt(16), "l2": Pt(14), "sub": Pt(22)}

def calculate_col_widths(headers, rows, total_width_inches=11.7):
    """
    Calculates column widths proportionally based on content length.
    """
    cols = len(headers)
    if cols == 0:
        return []
        
    max_lens = [len(h) for h in headers]
    for row in rows:
        for c_idx, cell in enumerate(row):
            if c_idx < cols:
                max_lens[c_idx] = max(max_lens[c_idx], len(cell))
                
    total_len = sum(max_lens)
    if total_len == 0:
        return [Inches(total_width_inches / cols)] * cols
        
    widths = []
    for l in max_lens:
        w = (l / total_len) * total_width_inches
        w = max(w, 1.2) # Minimum width threshold
        widths.append(Inches(w))
        
    sum_w = sum(w.inches for w in widths)
    scale = total_width_inches / sum_w
    return [Inches(w.inches * scale) for w in widths]

def add_table_to_slide(slide, table_block, top_inches):
    """
    Builds a beautifully styled native PowerPoint table inside a slide.
    """
    headers = table_block["headers"]
    rows = table_block["rows"]
    
    cols = len(headers)
    num_rows = len(rows) + 1
    
    col_widths = calculate_col_widths(headers, rows)
    total_width = sum(w.inches for w in col_widths)
    
    left = Inches((13.333 - total_width) / 2)
    top = Inches(top_inches)
    height = Inches(0.4 * num_rows)
    
    table_shape = slide.shapes.add_table(num_rows, cols, left, top, Inches(total_width), height)
    table = table_shape.table
    
    for col_idx, width in enumerate(col_widths):
        table.columns[col_idx].width = width
        
    # Format table header row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.08)
        cell.margin_bottom = Inches(0.08)
        
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = header_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
    # Format table body rows
    for row_idx, row in enumerate(rows):
        bg_color = RGBColor(241, 245, 249) if row_idx % 2 == 1 else RGBColor(255, 255, 255)
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            cell_text = row[col_idx] if col_idx < len(row) else ""
            p.text = cell_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT_NAME
            p.font.size = Pt(13)
            p.font.color.rgb = LIGHT_BODY

def add_image_to_slide(slide, image_block, box_left, box_top, box_w, box_h):
    """
    Fits and places a picture inside a bounding box on the slide while maintaining its aspect ratio.
    """
    img_path = image_block["path"]
    if not os.path.exists(img_path):
        print(f"Warning: Image file not found: {img_path}")
        return
        
    try:
        from PIL import Image
        img = Image.open(img_path)
        img_w, img_h = img.size
        img_ratio = img_w / img_h
        box_ratio = box_w / box_h
        
        if img_ratio > box_ratio:
            fit_w = box_w
            fit_h = box_w / img_ratio
        else:
            fit_h = box_h
            fit_w = box_h * img_ratio
            
        fit_left = box_left + (box_w - fit_w) / 2
        fit_top = box_top + (box_h - fit_h) / 2
        
        slide.shapes.add_picture(
            img_path, Inches(fit_left), Inches(fit_top), Inches(fit_w), Inches(fit_h)
        )
    except Exception as e:
        print(f"Warning: Could not add image {img_path}: {e}")

def create_presentation(slides, output_path, template_path=None):
    """
    Compiles parsed slide structured data into a PowerPoint presentation file.
    """
    prs = Presentation(template_path) if template_path else Presentation()
    
    # If using a template, remove all existing slides and drop XML relationships
    if template_path:
        for slide in list(prs.slides):
            # Find in slide id list and remove
            for sldId in prs.slides._sldIdLst:
                if sldId.id == slide.slide_id:
                    prs.slides._sldIdLst.remove(sldId)
                    break
            # Drop relationship to clean up XML parts in ppt container
            for rId, rel in prs.part.rels.items():
                if rel.target_part == slide.part:
                    prs.part.drop_rel(rId)
                    break
            
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(slides):
        is_title_slide = (idx == 0 and slide_data["type"] == "title")
        is_section_divider = (idx > 0 and slide_data["type"] == "title")
        
        if template_path:
            # Load template slide layouts
            if is_title_slide:
                layout = prs.slide_layouts[0]
            elif is_section_divider:
                layout = prs.slide_layouts[2]
            else:
                layout = prs.slide_layouts[1]
                
            slide = prs.slides.add_slide(layout)
            
            # Setup title/subtitle with layout placeholders
            if is_title_slide:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 0:
                        shape.text = slide_data["title"]
                        for p in shape.text_frame.paragraphs:
                            p.font.name = FONT_NAME
                    elif shape.placeholder_format.idx == 1:
                        shape.text = "簡報大綱與核心內容整理"
                        for p in shape.text_frame.paragraphs:
                            p.font.name = FONT_NAME
                
                # Delete unused placeholders
                for shape in list(slide.shapes):
                    if shape.is_placeholder and shape.placeholder_format.idx not in (0, 1):
                        sp = shape.element
                        sp.getparent().remove(sp)
                        
            elif is_section_divider:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 0:
                        shape.text = slide_data["title"]
                        for p in shape.text_frame.paragraphs:
                            p.font.name = FONT_NAME
                    elif shape.placeholder_format.idx == 1:
                        shape.text = f"SECTION 0{idx}"
                        for p in shape.text_frame.paragraphs:
                            p.font.name = FONT_NAME
                
                for shape in list(slide.shapes):
                    if shape.is_placeholder and shape.placeholder_format.idx not in (0, 1):
                        sp = shape.element
                        sp.getparent().remove(sp)
            else:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_data["title"]
                    # Position title box to avoid overlapping the Eshine template logo/line
                    title_shape.left = Inches(0.8)
                    title_shape.top = Inches(0.28)
                    title_shape.width = Inches(8.0)
                    title_shape.height = Inches(0.65)
                    
                    tf = title_shape.text_frame
                    tf.margin_left = Inches(0)
                    tf.margin_right = Inches(0)
                    tf.margin_top = Inches(0)
                    tf.margin_bottom = Inches(0)
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    
                    for p in tf.paragraphs:
                        p.font.name = FONT_NAME
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.color.rgb = LIGHT_TITLE
                
                # Delete all other placeholders
                for shape in list(slide.shapes):
                    if shape.is_placeholder and shape != title_shape:
                        sp = shape.element
                        sp.getparent().remove(sp)
        else:
            # Clean layout creation
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            if is_title_slide or is_section_divider:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = DARK_BG
                
                if is_title_slide:
                    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.2))
                    tf = title_box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
                    
                    p_title = tf.paragraphs[0]
                    p_title.text = slide_data["title"]
                    p_title.font.name = FONT_NAME
                    p_title.font.size = Pt(44)
                    p_title.font.bold = True
                    p_title.font.color.rgb = DARK_TITLE
                    p_title.space_after = Pt(12)
                    
                    p_sub = tf.add_paragraph()
                    p_sub.text = "簡報大綱與核心內容整理"
                    p_sub.font.name = FONT_NAME
                    p_sub.font.size = Pt(20)
                    p_sub.font.color.rgb = DARK_ACCENT
                    
                    line = slide.shapes.add_shape(
                        1, Inches(1.0), Inches(4.8), Inches(2.5), Inches(0.06)
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = DARK_ACCENT
                    line.line.color.rgb = DARK_ACCENT
                else:
                    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.2))
                    tf = title_box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
                    
                    p_sec = tf.paragraphs[0]
                    p_sec.text = f"SECTION 0{idx}"
                    p_sec.font.name = FONT_NAME
                    p_sec.font.size = Pt(14)
                    p_sec.font.bold = True
                    p_sec.font.color.rgb = DARK_ACCENT
                    p_sec.space_after = Pt(8)
                    
                    p_title = tf.add_paragraph()
                    p_title.text = slide_data["title"]
                    p_title.font.name = FONT_NAME
                    p_title.font.size = Pt(36)
                    p_title.font.bold = True
                    p_title.font.color.rgb = DARK_TITLE
                    
                    line = slide.shapes.add_shape(
                        1, Inches(1.0), Inches(4.5), Inches(2.0), Inches(0.05)
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = DARK_ACCENT
                    line.line.color.rgb = DARK_ACCENT
            else:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = LIGHT_BG
                
                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
                tf_title = title_box.text_frame
                tf_title.word_wrap = True
                tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = Inches(0)
                
                p_title = tf_title.paragraphs[0]
                p_title.text = slide_data["title"]
                p_title.font.name = FONT_NAME
                p_title.font.size = Pt(30)
                p_title.font.bold = True
                p_title.font.color.rgb = LIGHT_TITLE
                
                rule = slide.shapes.add_shape(
                    1, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02)
                )
                rule.fill.solid()
                rule.fill.fore_color.rgb = LIGHT_ACCENT
                rule.line.color.rgb = LIGHT_ACCENT

        if not is_title_slide and not is_section_divider:
            has_table = any(b["type"] == "table" for b in slide_data["blocks"])
            has_image = any(b["type"] == "image" for b in slide_data["blocks"])
            
            text_blocks = [b for b in slide_data["blocks"] if b["type"] in ("paragraph", "list")]
            table_blocks = [b for b in slide_data["blocks"] if b["type"] == "table"]
            image_blocks = [b for b in slide_data["blocks"] if b["type"] == "image"]
            
            is_card_layout = (
                not has_table and not has_image and
                len(text_blocks) == 1 and text_blocks[0]["type"] == "list" and
                2 <= len(text_blocks[0]["items"]) <= 4 and
                all(item["level"] == 0 for item in text_blocks[0]["items"]) and
                sum(len(item["text"]) for item in text_blocks[0]["items"]) < 200
            )
            
            content_top_y = 1.6 if template_path else 1.7
            
            if is_card_layout:
                items = text_blocks[0]["items"]
                n_cards = len(items)
                card_gap = 0.4
                total_w = 11.7
                card_w = (total_w - (n_cards - 1) * card_gap) / n_cards
                card_h = 3.8
                top_y = content_top_y + 0.4
                
                for c_idx, item in enumerate(items):
                    left_x = 0.8 + c_idx * (card_w + card_gap)
                    
                    card_shape = slide.shapes.add_shape(
                        1, Inches(left_x), Inches(top_y), Inches(card_w), Inches(card_h)
                    )
                    card_shape.fill.solid()
                    card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    card_shape.line.color.rgb = RGBColor(226, 232, 240)
                    card_shape.line.width = Pt(1.5)
                    
                    strip = slide.shapes.add_shape(
                        1, Inches(left_x), Inches(top_y), Inches(card_w), Inches(0.12)
                    )
                    strip.fill.solid()
                    strip.fill.fore_color.rgb = LIGHT_ACCENT
                    strip.line.color.rgb = LIGHT_ACCENT
                    
                    tb = slide.shapes.add_textbox(
                        Inches(left_x + 0.25), Inches(top_y + 0.4), 
                        Inches(card_w - 0.5), Inches(card_h - 0.6)
                    )
                    tf_card = tb.text_frame
                    tf_card.word_wrap = True
                    tf_card.margin_left = tf_card.margin_right = tf_card.margin_top = tf_card.margin_bottom = Inches(0)
                    
                    p = tf_card.paragraphs[0]
                    p.text = f"0{c_idx + 1}"
                    p.font.name = FONT_NAME
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.font.color.rgb = LIGHT_ACCENT
                    p.space_after = Pt(12)
                    
                    p2 = tf_card.add_paragraph()
                    p2.text = item["text"]
                    p2.font.name = FONT_NAME
                    p2.font.size = Pt(16)
                    p2.font.color.rgb = LIGHT_BODY
                    p2.space_after = Pt(0)
            
            elif has_image:
                if text_blocks:
                    font_sizes = get_font_sizes(text_blocks)
                    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(content_top_y), Inches(5.5), Inches(4.8))
                    tf_content = text_box.text_frame
                    tf_content.word_wrap = True
                    tf_content.margin_left = tf_content.margin_right = tf_content.margin_top = tf_content.margin_bottom = Inches(0)
                    render_text_blocks(tf_content, text_blocks, font_sizes)
                    
                    if image_blocks:
                        add_image_to_slide(slide, image_blocks[0], 6.8, content_top_y, 5.7, 4.8)
                else:
                    if image_blocks:
                        add_image_to_slide(slide, image_blocks[0], 0.8, content_top_y, 11.7, 4.8)
                        
            elif has_table:
                if text_blocks:
                    font_sizes = get_font_sizes(text_blocks)
                    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(content_top_y - 0.1), Inches(11.733), Inches(2.2))
                    tf_content = text_box.text_frame
                    tf_content.word_wrap = True
                    tf_content.margin_left = tf_content.margin_right = tf_content.margin_top = tf_content.margin_bottom = Inches(0)
                    render_text_blocks(tf_content, text_blocks, font_sizes)
                    
                    for tbl_b in table_blocks:
                        add_table_to_slide(slide, tbl_b, content_top_y + 2.3)
                else:
                    for tbl_b in table_blocks:
                        add_table_to_slide(slide, tbl_b, content_top_y + 0.3)
            else:
                if not template_path:
                    # Draw standard side vertical accent line when not using template
                    side_bar = slide.shapes.add_shape(
                        1, Inches(0.5), Inches(1.7), Inches(0.06), Inches(4.8)
                    )
                    side_bar.fill.solid()
                    side_bar.fill.fore_color.rgb = LIGHT_ACCENT
                    side_bar.line.color.rgb = LIGHT_ACCENT
                
                font_sizes = get_font_sizes(text_blocks)
                text_box = slide.shapes.add_textbox(Inches(0.8), Inches(content_top_y), Inches(11.733), Inches(4.8))
                tf_content = text_box.text_frame
                tf_content.word_wrap = True
                tf_content.margin_left = tf_content.margin_right = tf_content.margin_top = tf_content.margin_bottom = Inches(0)
                render_text_blocks(tf_content, text_blocks, font_sizes)
                
            # Footer styling
            if not template_path:
                footer_line = slide.shapes.add_shape(
                    1, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.01)
                )
                footer_line.fill.solid()
                footer_line.fill.fore_color.rgb = LIGHT_ACCENT_LINE
                footer_line.line.color.rgb = LIGHT_ACCENT_LINE
            
            footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.733), Inches(0.4))
            tf_foot = footer_box.text_frame
            tf_foot.word_wrap = True
            tf_foot.margin_left = tf_foot.margin_right = tf_foot.margin_top = tf_foot.margin_bottom = Inches(0)
            
            p_foot = tf_foot.paragraphs[0]
            p_foot.text = f"倉儲系統方案比較簡報 | Page {idx + 1}"
            p_foot.font.name = FONT_NAME
            p_foot.font.size = Pt(10)
            p_foot.font.color.rgb = MUTED_TEXT
            p_foot.alignment = PP_ALIGN.LEFT
            
    prs.save(output_path)

def render_text_blocks(text_frame, blocks, font_sizes):
    """
    Renders structured paragraphs and nested lists inside a text frame.
    """
    p_idx = 0
    for block in blocks:
        if block["type"] == "paragraph":
            p = text_frame.paragraphs[0] if p_idx == 0 else text_frame.add_paragraph()
            p_idx += 1
            
            p.text = block["text"]
            p.font.name = FONT_NAME
            p.font.bold = block.get("is_subheading", False)
            
            if block.get("is_subheading", False):
                p.font.size = font_sizes["sub"]
                p.font.color.rgb = LIGHT_ACCENT
                p.space_before = Pt(12)
                p.space_after = Pt(6)
            else:
                p.font.size = font_sizes["para"]
                p.font.color.rgb = LIGHT_BODY
                p.space_after = Pt(8)
                
        elif block["type"] == "list":
            for item in block["items"]:
                p = text_frame.paragraphs[0] if p_idx == 0 else text_frame.add_paragraph()
                p_idx += 1
                
                level = item["level"]
                bullet = "•" if level == 0 else "◦" if level == 1 else "▪"
                
                p.text = f"{bullet}\t{item['text']}"
                p.font.name = FONT_NAME
                p.font.bold = False
                
                if level == 0:
                    p.font.size = font_sizes["l0"]
                    p.font.color.rgb = LIGHT_BODY
                elif level == 1:
                    p.font.size = font_sizes["l1"]
                    p.font.color.rgb = LIGHT_BODY
                else:
                    p.font.size = font_sizes["l2"]
                    p.font.color.rgb = LIGHT_BODY
                    
                p.left_indent = Inches(0.45 * (level + 1))
                p.first_line_indent = Inches(-0.25)
                p.space_after = Pt(6)

def main():
    parser = argparse.ArgumentParser(description="Convert Markdown to PowerPoint Presentation (.pptx)")
    parser.add_argument("input_md", help="Path to the input Markdown file")
    parser.add_argument("output_pptx", nargs="?", default=None, help="Path to save the output PowerPoint presentation (optional)")
    parser.add_argument("--template", default=None, help="Path to the template PowerPoint presentation (optional)")
    
    args = parser.parse_args()
    
    input_file = args.input_md
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found.")
        sys.exit(1)
        
    output_file = args.output_pptx
    if not output_file:
        base_name, _ = os.path.splitext(input_file)
        output_file = base_name + ".pptx"
        
    template_file = args.template
    # Auto-detect Hongsheng template if not specified and exists in the current directory
    if not template_file:
        default_template = "鴻勝彰濱二廠申請保稅工廠評估報告0707.pptx"
        if os.path.exists(default_template):
            template_file = default_template
            print(f"Auto-detected template file: {default_template}")
            
    print(f"Parsing Markdown file: {input_file}...")
    slides = parse_markdown_to_slides(input_file)
    print(f"Found {len(slides)} slides.")
    
    if template_file:
        print(f"Compiling PowerPoint presentation using template '{template_file}' to: {output_file}...")
    else:
        print(f"Compiling PowerPoint presentation using clean layout to: {output_file}...")
        
    create_presentation(slides, output_file, template_file)
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    main()
