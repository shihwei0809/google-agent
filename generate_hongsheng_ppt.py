import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create the presentation object
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Color definitions
BG_COLOR = RGBColor(11, 14, 20)      # #0B0E14 (Premium Dark)
TEXT_WHITE = RGBColor(255, 255, 255) # #FFFFFF
TEXT_GRAY = RGBColor(187, 187, 187)  # #BBBBBB (Subtext)
COLOR_GREEN = RGBColor(39, 174, 96)   # #27AE60 (Neon Success Green)
COLOR_RED = RGBColor(211, 84, 0)      # #D35400 (Alert/Warning Orange-Red)
COLOR_GOLD = RGBColor(212, 175, 55)   # #D4AF37 (Traditional Gold)

# Image paths (Local images in workspace)
workspace_dir = r"g:\我的雲端硬碟\GOOGLE ANGET"
workspace_img_dir = os.path.join(workspace_dir, "images")

images = {
    "cover": os.path.join(workspace_img_dir, "hs_cover.png"),
    "roles": os.path.join(workspace_img_dir, "hs_roles.png"),
    "green": os.path.join(workspace_img_dir, "hs_scan_green.png"),
    "red": os.path.join(workspace_img_dir, "hs_scan_red.png"),
    "qc": os.path.join(workspace_img_dir, "hs_qc_dashboard.png")
}

def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_clean_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    return tf

# ----------------- SLIDE 1: COVER -----------------
slide_layout = prs.slide_layouts[6] # Blank slide
slide1 = prs.slides.add_slide(slide_layout)
set_dark_background(slide1)

if os.path.exists(images["cover"]):
    slide1.shapes.add_picture(images["cover"], Inches(5.5), Inches(0), Inches(7.833), Inches(7.5))

tf1 = add_clean_textbox(slide1, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.0))
p_top = tf1.paragraphs[0]
p_top.text = "企業級教育訓練教材"
p_top.font.name = 'Microsoft JhengHei'
p_top.font.size = Pt(14)
p_top.font.bold = True
p_top.font.color.rgb = COLOR_GOLD
p_top.space_after = Pt(15)

p_title = tf1.add_paragraph()
p_title.text = "軟管對刷稽核系統"
p_title.font.name = 'Microsoft JhengHei'
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE
p_title.space_after = Pt(10)

p_sub = tf1.add_paragraph()
p_sub.text = "鴻勝化學｜現場作業 × QC 授權放行"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(18)
p_sub.font.bold = True
p_sub.font.color.rgb = COLOR_GREEN
p_sub.space_after = Pt(35)

p_desc = tf1.add_paragraph()
p_desc.text = "訓練目標：\n✓ 降低管線誤接、防範人為疏失\n✓ 建立掃碼稽核紀錄與系統存證\n✓ 確保放行前完成 QC 嚴格授權放行"
p_desc.font.name = 'Microsoft JhengHei'
p_desc.font.size = Pt(12)
p_desc.font.color.rgb = TEXT_GRAY

# ----------------- SLIDE 2: PURPOSE -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_dark_background(slide2)

tf2 = add_clean_textbox(slide2, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
p = tf2.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t2 = tf2.add_paragraph()
p_t2.text = "01 目的：訓練定位與管控原則"
p_t2.font.name = 'Microsoft JhengHei'
p_t2.font.size = Pt(28)
p_t2.font.bold = True
p_t2.font.color.rgb = TEXT_WHITE
p_t2.space_after = Pt(20)

p_body = tf2.add_paragraph()
p_body.text = "本教材用於新人上線前訓練與現場複訓，聚焦「儲槽」與「軟管」兩組識別碼的一致性稽核。"
p_body.font.name = 'Microsoft JhengHei'
p_body.font.size = Pt(14)
p_body.font.color.rgb = TEXT_GRAY
p_body.space_after = Pt(30)

# Add Columns manually for sleek layout
tf2_col1 = add_clean_textbox(slide2, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.0))
p_col1_title = tf2_col1.paragraphs[0]
p_col1_title.text = "核心管控點"
p_col1_title.font.name = 'Microsoft JhengHei'
p_col1_title.font.size = Pt(18)
p_col1_title.font.bold = True
p_col1_title.font.color.rgb = COLOR_GREEN
p_col1_title.space_after = Pt(15)

p_col1_desc = tf2_col1.add_paragraph()
p_col1_desc.text = "• 人員安全確認：先確認操作人員，排除無照作業\n• 櫃車任務索引：建立唯一櫃號任務作為關聯\n• 即時對刷掃描：雙端條碼匹配判定\n• QC 最終審查：無口頭替代，必須系統放行"
p_col1_desc.font.name = 'Microsoft JhengHei'
p_col1_desc.font.size = Pt(13)
p_col1_desc.font.color.rgb = TEXT_WHITE
p_col1_desc.space_after = Pt(10)

tf2_col2 = add_clean_textbox(slide2, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.0))
p_col2_title = tf2_col2.paragraphs[0]
p_col2_title.text = "管控輸出與存證"
p_col2_title.font.name = 'Microsoft JhengHei'
p_col2_title.font.size = Pt(18)
p_col2_title.font.bold = True
p_col2_title.font.color.rgb = COLOR_GREEN
p_col2_title.space_after = Pt(15)

p_col2_desc = tf2_col2.add_paragraph()
p_col2_desc.text = "• 暫存紀錄：現場操作之儲槽與軟管掃碼結果\n• 待稽核任務清單：同步派發至 QC 端看板\n• 系統授權軌跡：保存 QC 識別證刷卡時間與 ID\n• 紅色異常阻斷紀錄：自動存證便於後續追蹤"
p_col2_desc.font.name = 'Microsoft JhengHei'
p_col2_desc.font.size = Pt(13)
p_col2_desc.font.color.rgb = TEXT_WHITE

# ----------------- SLIDE 3: ROLES -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_dark_background(slide3)

if os.path.exists(images["roles"]):
    slide3.shapes.add_picture(images["roles"], Inches(7.0), Inches(1.2), Inches(5.533), Inches(5.5))

tf3 = add_clean_textbox(slide3, Inches(0.8), Inches(0.6), Inches(5.8), Inches(6.0))
p = tf3.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t3 = tf3.add_paragraph()
p_t3.text = "02 角色：系統角色與責任分工"
p_t3.font.name = 'Microsoft JhengHei'
p_t3.font.size = Pt(28)
p_t3.font.bold = True
p_t3.font.color.rgb = TEXT_WHITE
p_t3.space_after = Pt(25)

roles = [
    ("👷 現場作業人員 (Operator)", "選擇人員帳號 ➔ 輸入櫃號 ➔ 掃描儲槽碼與軟管對刷 ➔ 匹配後提交 QC 審查"),
    ("🛡️ 現場 QC 稽核員 (QC Auditor)", "查詢待稽核任務看板 ➔ 現場覆核管線狀態 ➔ 掃描 QC 識別證完成授權放行"),
    ("👑 班組主管 (Supervisor)", "每日抽查稽核紀錄 ➔ 異常紅色阻斷處置 ➔ 追蹤新人教育訓練狀態與成效")
]

for role_title, role_desc in roles:
    pr_t = tf3.add_paragraph()
    pr_t.text = role_title
    pr_t.font.name = 'Microsoft JhengHei'
    pr_t.font.size = Pt(15)
    pr_t.font.bold = True
    pr_t.font.color.rgb = COLOR_GOLD
    pr_t.space_after = Pt(5)
    
    pr_d = tf3.add_paragraph()
    pr_d.text = role_desc
    pr_d.font.name = 'Microsoft JhengHei'
    pr_d.font.size = Pt(11)
    pr_d.font.color.rgb = TEXT_GRAY
    pr_d.space_after = Pt(20)

# ----------------- SLIDE 4: SOP -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_dark_background(slide4)

tf4 = add_clean_textbox(slide4, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
p = tf4.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t4 = tf4.add_paragraph()
p_t4.text = "03 SOP：標準作業流程總覽"
p_t4.font.name = 'Microsoft JhengHei'
p_t4.font.size = Pt(28)
p_t4.font.bold = True
p_t4.font.color.rgb = TEXT_WHITE
p_t4.space_after = Pt(30)

steps = [
    ("1. 選擇人員", "未選人員時，對刷掃描功能自動停用鎖定"),
    ("2. 輸入櫃號", "建立唯一櫃車任務編號作為後續稽核索引"),
    ("3. 對刷掃碼", "依序掃描儲槽識別碼與軟管識別碼條碼"),
    ("4. 系統判定", "綠色代表匹配正確准予提交；紅色代表誤接強制阻斷"),
    ("5. QC 放行", "QC 檢視待稽核任務，確認無誤後刷識別證放行")
]

# Render flow horizontally
for idx, (s_title, s_desc) in enumerate(steps):
    left_pos = Inches(0.8 + idx * 2.3)
    tf_flow = add_clean_textbox(slide4, left_pos, Inches(2.2), Inches(2.1), Inches(2.2))
    p_s_t = tf_flow.paragraphs[0]
    p_s_t.text = s_title
    p_s_t.font.name = 'Microsoft JhengHei'
    p_s_t.font.size = Pt(14)
    p_s_t.font.bold = True
    p_s_t.font.color.rgb = COLOR_GREEN
    p_s_t.space_after = Pt(8)
    
    p_s_d = tf_flow.add_paragraph()
    p_s_d.text = s_desc
    p_s_d.font.name = 'Microsoft JhengHei'
    p_s_d.font.size = Pt(11)
    p_s_d.font.color.rgb = TEXT_GRAY

# Add Important Alert at the bottom
tf_alert = add_clean_textbox(slide4, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0))
p_a_t = tf_alert.paragraphs[0]
p_a_t.text = "⚠️ 重要判定與阻斷機制"
p_a_t.font.name = 'Microsoft JhengHei'
p_a_t.font.size = Pt(14)
p_a_t.font.bold = True
p_a_t.font.color.rgb = COLOR_RED
p_a_t.space_after = Pt(5)

p_a_d = tf_alert.add_paragraph()
p_a_d.text = "• 判定原則：儲槽識別碼與軟管識別碼必須在系統註冊中屬於同一條允收管線，才判定為匹配。\n• 嚴防人工口頭替代：若系統出現紅色錯誤阻斷，現場不得口頭妥協，必須停止提交並通報現場複查。"
p_a_d.font.name = 'Microsoft JhengHei'
p_a_d.font.size = Pt(12)
p_a_d.font.color.rgb = TEXT_GRAY

# ----------------- SLIDE 5: FIELD SETUP -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_dark_background(slide5)

if os.path.exists(images["green"]):
    slide5.shapes.add_picture(images["green"], Inches(0.8), Inches(1.2), Inches(5.533), Inches(5.5))

tf5 = add_clean_textbox(slide5, Inches(7.0), Inches(0.6), Inches(5.533), Inches(6.0))
p = tf5.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t5 = tf5.add_paragraph()
p_t5.text = "04 現場：人員與櫃號建立"
p_t5.font.name = 'Microsoft JhengHei'
p_t5.font.size = Pt(28)
p_t5.font.bold = True
p_t5.font.color.rgb = TEXT_WHITE
p_t5.space_after = Pt(30)

items = [
    ("🛡️ 人員防呆限制", "在未完成操作人員選擇前，對刷掃描按鈕將呈現灰色停用狀態，防止無登錄人員作業。"),
    ("🚚 建立櫃車任務", "輸入櫃車之櫃號（作為本次卸料任務之唯一索引識別），方可開啟對刷掃描頁面。"),
    ("⏳ 資料暫存機制", "對刷判定正確後，系統會建立暫存紀錄，並透過網路即時向現場 QC 看板發送待稽核通知。")
]

for title, desc in items:
    pt = tf5.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(15)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GREEN
    pt.space_after = Pt(5)
    
    pd = tf5.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(12)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# ----------------- SLIDE 6: SCAN & MATCH -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_dark_background(slide6)

tf6 = add_clean_textbox(slide6, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
p = tf6.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t6 = tf6.add_paragraph()
p_t6.text = "05 對刷：現場對刷與即時判定"
p_t6.font.name = 'Microsoft JhengHei'
p_t6.font.size = Pt(28)
p_t6.font.bold = True
p_t6.font.color.rgb = TEXT_WHITE
p_t6.space_after = Pt(25)

# Render side-by-side split screen comparisons
# Left half: Mismatch (Red)
tf6_left = add_clean_textbox(slide6, Inches(0.8), Inches(2.0), Inches(5.4), Inches(4.5))
p_l_t = tf6_left.paragraphs[0]
p_l_t.text = "🔴 紅色畫面：管線誤接 (Mismatch)"
p_l_t.font.name = 'Microsoft JhengHei'
p_l_t.font.size = Pt(18)
p_l_t.font.bold = True
p_l_t.font.color.rgb = COLOR_RED
p_l_t.space_after = Pt(10)

p_l_d = tf6_left.add_paragraph()
p_l_d.text = "• 系統狀態：條碼判定不匹配，警示蜂鳴器響起。\n• 強制阻斷：畫面顯示紅色警告阻斷，無法按下提交按鈕。\n• 處置規範：立刻暫停對接，返回現場核對儲槽與軟管編號。"
p_l_d.font.name = 'Microsoft JhengHei'
p_l_d.font.size = Pt(12)
p_l_d.font.color.rgb = TEXT_GRAY

# Right half: Match (Green)
tf6_right = add_clean_textbox(slide6, Inches(7.0), Inches(2.0), Inches(5.4), Inches(4.5))
p_r_t = tf6_right.paragraphs[0]
p_r_t.text = "🟢 綠色畫面：匹配正確 (Match)"
p_r_t.font.name = 'Microsoft JhengHei'
p_r_t.font.size = Pt(18)
p_r_t.font.bold = True
p_r_t.font.color.rgb = COLOR_GREEN
p_r_t.space_after = Pt(10)

p_r_d = tf6_right.add_paragraph()
p_r_d.text = "• 系統狀態：判定儲槽與軟管碼一致，警示綠燈亮起。\n• 允許提交：畫面顯示綠色匹配成功，啟用資料提交按鈕。\n• 下步操作：點選提交資料，通知 QC 稽核員前來授權。"
p_r_d.font.name = 'Microsoft JhengHei'
p_r_d.font.size = Pt(12)
p_r_d.font.color.rgb = TEXT_GRAY

# ----------------- SLIDE 7: EXCEPTION HANDLING -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_dark_background(slide7)

if os.path.exists(images["red"]):
    slide7.shapes.add_picture(images["red"], Inches(0.8), Inches(1.2), Inches(5.533), Inches(5.5))

tf7 = add_clean_textbox(slide7, Inches(7.0), Inches(0.6), Inches(5.533), Inches(6.0))
p = tf7.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t7 = tf7.add_paragraph()
p_t7.text = "06 異常：錯誤情境與現場處置"
p_t7.font.name = 'Microsoft JhengHei'
p_t7.font.size = Pt(28)
p_t7.font.bold = True
p_t7.font.color.rgb = TEXT_WHITE
p_t7.space_after = Pt(30)

items = [
    ("⚠️ 現場強效阻斷機制", "系統拋出 Mismatch 阻斷警告時，現場作業人員絕對禁止強行對接閥門，嚴防誤接造成重大化學事故。"),
    ("📞 主動通報現場主管", "現場作業人員應立即保持現狀，通報現場 QC 稽核員或班組主管至現場進行實體管線查核。"),
    ("✏️ 保留異常紀錄軌跡", "系統會自動留存紅色誤接記錄，由主管審查是否有條碼汙損或管線規劃錯誤，修正後重新對刷。")
]

for title, desc in items:
    pt = tf7.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(15)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_RED
    pt.space_after = Pt(5)
    
    pd = tf7.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(12)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# ----------------- SLIDE 8: SUCCESS SCENARIO -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_dark_background(slide8)

tf8 = add_clean_textbox(slide8, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
p = tf8.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t8 = tf8.add_paragraph()
p_t8.text = "07 成功：匹配正確與資料提交"
p_t8.font.name = 'Microsoft JhengHei'
p_t8.font.size = Pt(28)
p_t8.font.bold = True
p_t8.font.color.rgb = TEXT_WHITE
p_t8.space_after = Pt(30)

success_stages = [
    ("對刷成功", "當系統辨識條碼無誤，且驗證儲槽與軟管碼一致時，畫面轉為綠色 MATCH OK 狀態。"),
    ("提交暫存", "點擊提交，系統將儲槽號、軟管號、作業人員及建立時間加密暫存。"),
    ("派發看板", "現場看板即時出現該筆待稽核任務，提醒 QC 前往現場覆核。")
]

for idx, (stage_title, stage_desc) in enumerate(success_stages):
    left_pos = Inches(0.8 + idx * 3.9)
    tf_stage = add_clean_textbox(slide8, left_pos, Inches(2.2), Inches(3.6), Inches(4.0))
    p_s_t = tf_stage.paragraphs[0]
    p_s_t.text = stage_title
    p_s_t.font.name = 'Microsoft JhengHei'
    p_s_t.font.size = Pt(18)
    p_s_t.font.bold = True
    p_s_t.font.color.rgb = COLOR_GREEN
    p_s_t.space_after = Pt(15)
    
    p_s_d = tf_stage.add_paragraph()
    p_s_d.text = stage_desc
    p_s_d.font.name = 'Microsoft JhengHei'
    p_s_d.font.size = Pt(13)
    p_s_d.font.color.rgb = TEXT_GRAY

# ----------------- SLIDE 9: QC AUDIT -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_dark_background(slide9)

if os.path.exists(images["qc"]):
    slide9.shapes.add_picture(images["qc"], Inches(7.0), Inches(1.2), Inches(5.533), Inches(5.5))

tf9 = add_clean_textbox(slide9, Inches(0.8), Inches(0.6), Inches(5.8), Inches(6.0))
p = tf9.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t9 = tf9.add_paragraph()
p_t9.text = "08 QC：QC 授權放行端操作"
p_t9.font.name = 'Microsoft JhengHei'
p_t9.font.size = Pt(28)
p_t9.font.bold = True
p_t9.font.color.rgb = TEXT_WHITE
p_t9.space_after = Pt(25)

qc_actions = [
    ("📊 檢視待稽核清單", "QC 主控面板登入後，檢視是否有現場作業提交的待放行筆數與詳細欄位。"),
    ("🔍 二次覆核", "QC 核對進出貨單據，確認是否有此槽號，並確認儲槽與管路實體掃描是否正確。"),
    ("💳 識別證授權放行", "確認無誤後，QC 掃描個人識別證條碼，系統紀錄授權軌跡完成安全放行。")
]

for title, desc in qc_actions:
    pt = tf9.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(15)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GOLD
    pt.space_after = Pt(5)
    
    pd = tf9.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(12)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# ----------------- SLIDE 10: TEST & COMPLIANCE -----------------
slide10 = prs.slides.add_slide(slide_layout)
set_dark_background(slide10)

tf10 = add_clean_textbox(slide10, Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.3))
p = tf10.paragraphs[0]
p.text = "鴻勝化學｜軟管對刷稽核系統"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_GOLD
p.space_after = Pt(5)

p_t10 = tf10.add_paragraph()
p_t10.text = "09 測驗：新人上線檢核表與規範"
p_t10.font.name = 'Microsoft JhengHei'
p_t10.font.size = Pt(28)
p_t10.font.bold = True
p_t10.font.color.rgb = TEXT_WHITE
p_t10.space_after = Pt(20)

tf10_left = add_clean_textbox(slide10, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
p_chk_t = tf10_left.paragraphs[0]
p_chk_t.text = "上線前實務檢核項"
p_chk_t.font.name = 'Microsoft JhengHei'
p_chk_t.font.size = Pt(16)
p_chk_t.font.bold = True
p_chk_t.font.color.rgb = COLOR_GREEN
p_chk_t.space_after = Pt(10)

p_chk_d = tf10_left.add_paragraph()
p_chk_d.text = "✓ 能完整說明儲槽與軟管對刷稽核之安控目的\n✓ 能依循 SOP 完成選人員、輸櫃號、掃描操作\n✓ 能正確分辨綠色 Match 與紅色 Mismatch 畫面\n✓ 知道出現紅色阻斷時絕對禁止提交且主動回報\n✓ 能熟練操作 QC 授權放行端的確認與刷卡步驟"
p_chk_d.font.name = 'Microsoft JhengHei'
p_chk_d.font.size = Pt(12)
p_chk_d.font.color.rgb = TEXT_WHITE
p_chk_d.space_after = Pt(10)

tf10_right = add_clean_textbox(slide10, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
p_std_t = tf10_right.paragraphs[0]
p_std_t.text = "評核標準與主管抽查重點"
p_std_t.font.name = 'Microsoft JhengHei'
p_std_t.font.size = Pt(16)
p_std_t.font.bold = True
p_std_t.font.color.rgb = COLOR_GREEN
p_std_t.space_after = Pt(10)

p_std_d = tf10_right.add_paragraph()
p_std_d.text = "• 考核要求：學科測驗 80 分以上，且在主管觀察下獨立完成現場實際情境演練一次。\n• 抽查稽核要點：主管應定期抽查後台，確認現場是否有「未刷卡直接對接」或口頭核准等不合規行為。\n• 異常持續追蹤：對於阻斷紀錄進行二次原因查證，並存檔於專案歷史庫。"
p_std_d.font.name = 'Microsoft JhengHei'
p_std_d.font.size = Pt(12)
p_std_d.font.color.rgb = TEXT_WHITE

# Save the presentation
output_path = os.path.join(workspace_dir, "鴻勝化學_軟管對刷稽核系統_正式版_10頁_圖文可編輯版.pptx")
prs.save(output_path)
print(f"Presentation successfully created and saved to {output_path}!")
