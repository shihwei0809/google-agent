import os
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import google.generativeai as genai
from dotenv import load_dotenv
import shutil
import io

load_dotenv()

app = FastAPI(title="AI 訓練平台 API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 掛載 materials 目錄為靜態資料夾，用來提供圖片等檔案
MATERIALS_DIR = "materials"
os.makedirs(MATERIALS_DIR, exist_ok=True)
app.mount("/materials_static", StaticFiles(directory=MATERIALS_DIR), name="materials_static")

# 讀取多個 API Key (用逗號分隔)
GEMINI_API_KEYS_STR = os.getenv("GEMINI_API_KEY", "")
API_KEYS = [k.strip() for k in GEMINI_API_KEYS_STR.split(",") if k.strip()]

if not API_KEYS:
    print("[警告] 未設定任何 GEMINI_API_KEY")
else:
    genai.configure(api_key=API_KEYS[0])

# 定義模型優先順序清單 (2026 最新版本，以免費/高速的 Flash 為主)
MODEL_FALLBACKS = [
    'gemini-3.8-flash',
    'gemini-3.7-flash',
    'gemini-3.1-pro'
]

def call_gemini_with_fallback(prompt_or_list):
    """共用的 Gemini API 呼叫函數，支援多組 API Key 與多模型自動降級 (Fallback) 機制"""
    for key in API_KEYS:
        genai.configure(api_key=key)
        for model_name in MODEL_FALLBACKS:
            try:
                model = genai.GenerativeModel(model_name)
                res = model.generate_content(prompt_or_list)
                return res
            except Exception as e:
                print(f"[Fallback] 模型 {model_name} (Key: {key[:4]}...) 呼叫失敗: {e}")
                continue
    raise Exception("所有模型與 API Key 皆無法順利回應，請檢查配額與網路連線")

class ChatRequest(BaseModel):
    message: str
    context: str = ""
    material_name: str = "未知教材"

@app.get("/")
def read_root():
    return {"status": "ok", "message": "AI 訓練平台 API 運作中"}

# --- 教材管理 API ---

@app.get("/materials")
def list_materials():
    files = [f for f in os.listdir(MATERIALS_DIR) if f.endswith('.md') or f.endswith('.txt')]
    return {"materials": files}

@app.get("/materials/{filename}")
def get_material(filename: str):
    filepath = os.path.join(MATERIALS_DIR, filename)
    if not os.path.exists(filepath):
        return {"error": "找不到該教材"}
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()
    return {"content": content}

import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

@app.post("/materials")
async def upload_material(file: UploadFile = File(...)):
    filename = file.filename
    ext = filename.lower().split('.')[-1]
    
    # 讀取檔案內容至記憶體
    content_bytes = await file.read()
    
    # 若是 md 或 txt，直接存檔
    if ext in ['md', 'txt']:
        filepath = os.path.join(MATERIALS_DIR, filename)
        with open(filepath, "wb") as f:
            f.write(content_bytes)
        return {"message": "上傳成功", "filename": filename}
    
    # 若是其他格式，進行文字萃取並轉為 md
    extracted_text = f"# {filename} (系統自動轉換)\n\n"
    new_filename = filename.rsplit('.', 1)[0] + ".md"
    filepath = os.path.join(MATERIALS_DIR, new_filename)
    
    try:
        if ext == 'pdf':
            doc = fitz.open(stream=content_bytes, filetype="pdf")
            for page in doc:
                extracted_text += page.get_text() + "\n\n"
                
        elif ext == 'docx':
            doc_file = io.BytesIO(content_bytes)
            doc = docx.Document(doc_file)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n\n"
                
        elif ext == 'xlsx':
            excel_file = io.BytesIO(content_bytes)
            # 讀取所有 sheet
            xl = pd.ExcelFile(excel_file)
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(xl, sheet_name=sheet_name)
                extracted_text += f"## {sheet_name}\n\n"
                extracted_text += df.to_markdown(index=False) + "\n\n"
                
        elif ext == 'pptx':
            ppt_file = io.BytesIO(content_bytes)
            prs = Presentation(ppt_file)
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for i, slide in enumerate(prs.slides):
                extracted_text += f"## 第 {i+1} 頁\n\n"
                for j, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        extracted_text += shape.text + "\n"
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_bytes = shape.image.blob
                        img_ext = shape.image.ext
                        img_filename = f"{filename.rsplit('.', 1)[0]}_p{i+1}_{j+1}.{img_ext}"
                        img_filepath = os.path.join(MATERIALS_DIR, img_filename)
                        with open(img_filepath, "wb") as img_f:
                            img_f.write(img_bytes)
                        extracted_text += f"\n![圖片]({img_filename})\n\n"
                extracted_text += "\n"
        elif ext in ['mp4', 'mov', 'avi', 'webm']:
            import tempfile
            import time
            
            # 暫存影片檔案供 Gemini 讀取
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                tmp.write(content_bytes)
                tmp_path = tmp.name
                
            try:
                if not API_KEYS:
                    raise Exception("尚未設定 GEMINI_API_KEY，無法執行影片解析")
                genai.configure(api_key=API_KEYS[0])
                
                print(f"正在上傳影片至 Gemini: {filename}...")
                video_file = genai.upload_file(path=tmp_path)
                
                print("等待影片處理中...")
                while video_file.state.name == 'PROCESSING':
                    time.sleep(2)
                    video_file = genai.get_file(video_file.name)
                
                if video_file.state.name == 'FAILED':
                    raise Exception("Gemini 影片處理失敗")
                    
                print("影片處理完成，開始產生 SOP...")
                prompt = """你是一個專業的教育訓練教材撰寫專家。
請仔細觀看這段系統操作影片（無論是否有聲音），將人員的操作流程轉化為一份高品質的 Markdown 圖文教學教材。
請必須嚴格包含以下區塊：
1. 💡 教材核心重點 (Key Takeaways)：提煉出這個影片中最核心的 3 個作業重點或防呆注意事項。
2. 🖼️ 系統流程圖：請根據影片的操作邏輯，使用 Mermaid 語法繪製一段精簡的流程圖 (graph TD)。
3. 📖 步驟解析：詳細記錄每個重要的點擊位置與欄位輸入，並使用要點式 (bullet points) 條列說明。在每個重要步驟下方，請加入「![畫面截圖](圖示建議)」作為圖片佔位符，確保基層員工能圖文對照學習。
"""
                
                res = call_gemini_with_fallback([video_file, prompt])
                extracted_text += res.text
                
                try:
                    genai.delete_file(video_file.name)  # 清理雲端空間
                except:
                    pass
            finally:
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
        else:
            return {"error": "不支援的檔案格式"}
            
        # ==========================================
        # AI 重寫與精煉 (適用於 PDF, DOCX, PPTX 等靜態文件)
        # ==========================================
        if ext not in ['mp4', 'mov', 'avi', 'webm', 'md', 'txt']:
            print(f"原始文件萃取完成，開始使用 AI 提煉精華與重寫版面 ({ext})...")
            rewrite_prompt = f"""你是一個專業的教育訓練教材撰寫專家。
以下是從原始文件中萃取出來的文字（以及保留的圖片標籤）。請仔細閱讀並理解內容，重新排版並提煉精華，產出一份給基層員工閱讀的高品質 Markdown 圖文教學教材。
請嚴格包含以下區塊：
1. 💡 教材核心重點 (Key Takeaways)：提煉出這份教材最核心的 3 個重點。
2. 🖼️ 系統流程圖：請根據內容邏輯，使用 Mermaid 語法繪製精簡的流程圖 (graph TD)。
3. 📖 步驟解析：將原本的內容有邏輯地分章節列出，文字敘述必須簡單易懂。**重要：請務必保留原文中所有的圖片標籤 `![圖片](...)` 不可刪除，將它們安插在適合的步驟段落中，以便員工圖文對照學習**。

原始內容如下：
{extracted_text}
"""
            try:
                res = call_gemini_with_fallback(rewrite_prompt)
                extracted_text = res.text
            except Exception as e:
                print(f"AI 提煉失敗，退回原始萃取文字: {e}")
                # 若 AI 提煉失敗，保留原本原始萃取的 extracted_text

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(extracted_text)
            
        return {"message": "轉換並上傳成功", "filename": new_filename}
        
    except Exception as e:
        print(f"轉換失敗: {e}")
        return {"error": f"檔案解析失敗: {str(e)}"}

@app.delete("/materials/{filename}")
def delete_material(filename: str):
    filepath = os.path.join(MATERIALS_DIR, filename)
    if os.path.exists(filepath):
        os.remove(filepath)
        return {"message": "刪除成功"}
    return {"error": "檔案不存在"}

from pydantic import BaseModel

class UpdateMaterialRequest(BaseModel):
    content: str

@app.put("/materials/{filename}")
def update_material(filename: str, req: UpdateMaterialRequest):
    filepath = os.path.join(MATERIALS_DIR, filename)
    if not os.path.exists(filepath):
        return {"error": "找不到該教材"}
    try:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(req.content)
        return {"message": "更新成功"}
    except Exception as e:
        return {"error": f"更新失敗: {str(e)}"}

class GenerateImageRequest(BaseModel):
    prompt: str

import urllib.request
import urllib.parse
import time

@app.post("/generate_image")
def generate_image(req: GenerateImageRequest):
    if not API_KEYS:
        return {"error": "未設定 API Key"}
    try:
        # 1. 將使用者的中文提示翻譯成精準的英文繪圖提示詞 (使用 Gemini 自動降級機制)
        prompt = f"Translate this image generation prompt to English. Just output the English text, add details if needed to make it look professional, no extra words: {req.prompt}"
        trans_res = call_gemini_with_fallback(prompt)
        eng_prompt = trans_res.text.strip()
        
        # 2. 呼叫外部免金鑰 AI 繪圖 API (Pollinations) 進行繪圖
        safe_prompt = urllib.parse.quote(eng_prompt)
        url = f"https://image.pollinations.ai/prompt/{safe_prompt}?nologo=true&width=800&height=600"
        
        # 3. 將生成的圖片下載並存入系統的教材圖片庫
        timestamp = int(time.time())
        filename = f"ai_img_{timestamp}.jpg"
        filepath = os.path.join(MATERIALS_DIR, filename)
        
        req_img = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'})
        with urllib.request.urlopen(req_img) as response, open(filepath, 'wb') as out_file:
            out_file.write(response.read())
        
        return {"filename": filename}
    except Exception as e:
        return {"error": str(e)}

@app.post("/upload_image")
async def upload_image(file: UploadFile = File(...)):
    """專門處理前端貼上 (Paste) 截圖的 API"""
    try:
        ext = file.filename.split('.')[-1].lower()
        if ext not in ['png', 'jpg', 'jpeg', 'gif']:
            ext = 'png' # 預設副檔名
            
        content_bytes = await file.read()
        timestamp = int(time.time())
        filename = f"screenshot_{timestamp}.{ext}"
        filepath = os.path.join(MATERIALS_DIR, filename)
        
        with open(filepath, "wb") as f:
            f.write(content_bytes)
            
        return {"url": filename}
    except Exception as e:
        return {"error": str(e)}

# --- AI 問答 API ---

# 定義模型優先順序清單 (2026 最新版本)
MODEL_FALLBACKS = [
    'gemini-3.8-flash',
    'gemini-3.7-flash',
    'gemini-3.1-pro'
]

import csv
import datetime

from fastapi.responses import StreamingResponse

@app.post("/chat")
async def chat_with_ai(req: ChatRequest):
    if not API_KEYS:
        return {"response": "系統尚未設定任何 GEMINI_API_KEY，無法提供 AI 服務。"}
    
    # 自動讀取所有教材 (全知模式)
    all_materials_content = ""
    if os.path.exists(MATERIALS_DIR):
        for filename in os.listdir(MATERIALS_DIR):
            if filename.endswith(".md"):
                try:
                    with open(os.path.join(MATERIALS_DIR, filename), "r", encoding="utf-8") as f:
                        all_materials_content += f"\n\n--- 教材: {filename} ---\n" + f.read()
                except:
                    pass

    prompt = f"你是一個專業的企業內訓 AI 助教。請根據以下【所有教材內容】，親切且專業地回答學員的問題。\n如果學員的問題跨越了多份教材，請幫忙統整答案。\n如果問題與教材完全無關，請委婉告知。\n\n【所有教材內容】\n{all_materials_content}\n\n【學員問題】\n{req.message}"
    
    # 雙重備援機制：先輪替 API Keys，再輪替模型
    for key_idx, current_key in enumerate(API_KEYS):
        # 切換當前使用的 API Key
        genai.configure(api_key=current_key)
        
        for model_name in MODEL_FALLBACKS:
            try:
                model = genai.GenerativeModel(model_name)
                # 開啟 stream=True 模式
                response = model.generate_content(prompt, stream=True)
                
                async def generate():
                    full_text = ""
                    # 逐字回傳給前端
                    for chunk in response:
                        if chunk.text:
                            full_text += chunk.text
                            yield chunk.text
                    
                    # 提示當下使用的模型與 Key (僅顯示前幾碼以利辨識)
                    masked_key = f"{current_key[:4]}...{current_key[-4:]}"
                    footer = f"\n\n*(Powered by {model_name} / Key: {masked_key})*"
                    full_text += footer
                    yield footer
                    
                    # --- 在回傳結束後寫入 Excel (以月份分頁) ---
                    try:
                        import openpyxl
                        log_file = "chat_logs.xlsx"
                        now = datetime.datetime.now()
                        month_str = now.strftime("%Y-%m")
                        timestamp = now.strftime("%Y-%m-%d %H:%M:%S")
                        
                        if os.path.exists(log_file):
                            wb = openpyxl.load_workbook(log_file)
                        else:
                            wb = openpyxl.Workbook()
                            # 移除預設的空 Sheet
                            if "Sheet" in wb.sheetnames:
                                del wb["Sheet"]
                                
                        # 若當月的分頁不存在，則建立並寫入標題列
                        if month_str not in wb.sheetnames:
                            ws = wb.create_sheet(title=month_str)
                            ws.append(["時間", "當前檢視教材", "學員提問", "AI回覆"])
                        else:
                            ws = wb[month_str]
                            
                        # 寫入提問紀錄
                        ws.append([timestamp, req.material_name, req.message, full_text])
                        wb.save(log_file)
                    except Exception as log_e:
                        print(f"寫入 Excel 日誌失敗: {log_e}")
                    # -----------------------------
                    
                return StreamingResponse(generate(), media_type="text/plain")
            except Exception as e:
                error_msg = str(e)
                print(f"[警告] Key({key_idx+1}/{len(API_KEYS)}) 模型 {model_name} 呼叫失敗 ({error_msg})，嘗試切換...")
                continue
                
    return {"response": "系統內建的所有 API Key 以及 Gemini 模型備援方案皆已用盡或發生異常，請聯絡系統管理員！"}
