import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_optimized_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # ------------------ DESIGN SYSTEM COLORS ------------------
    BG_COLOR = RGBColor(11, 14, 20)        # #0B0E14 Premium Dark
    TEXT_WHITE = RGBColor(255, 255, 255)  # #FFFFFF High contrast text
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Cool gray subtext
    
    COLOR_RAW = RGBColor(56, 189, 248)     # #38BDF8 Sky Blue for Raw Materials
    COLOR_PROCESS = RGBColor(251, 191, 36) # #FBBF24 Amber for Process / Intermediate
    COLOR_FINISH = RGBColor(52, 211, 153)  # #34D399 Emerald for Finished Products
    COLOR_OFFGRADE = RGBColor(244, 63, 94) # #F43F5E Coral/Rose for Off-grade / Checks
    COLOR_UTILITY = RGBColor(203, 213, 225) # #CBD5E1 Off-white/Slate for N2 / Utility
    
    CARD_BG = RGBColor(22, 28, 38)         # #161C26 Dark slate for container boxes
    
    # Fonts
    FONT_FAMILY = "Microsoft JhengHei"
    
    # Helper: Set Slide Dark Background
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Helper: Create Header
    def add_slide_header(slide, title_text, total_capacity=""):
        # Decorative top border
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_PROCESS
        top_bar.line.fill.background()
        
        # Header text frame (shortened to 6.5 inches to prevent overlapping legend)
        header_tf = add_clean_textbox(slide, Inches(0.8), Inches(0.3), Inches(6.5), Inches(0.9))
        p = header_tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        if total_capacity:
            p_cap = header_tf.add_paragraph()
            p_cap.text = f"系統總容量: {total_capacity} | 流程配置: 原料區 ➔ 製程 ➔ 待驗槽 ➔ 成品區"
            p_cap.font.name = FONT_FAMILY
            p_cap.font.size = Pt(9.5)
            p_cap.font.bold = True
            p_cap.font.color.rgb = COLOR_PROCESS
            
    # Helper: Add Slide Legend (compacted to prevent header overlapping)
    def add_slide_legend(slide):
        legend_items = [
            ("原料", COLOR_RAW),
            ("製程 / 調合", COLOR_PROCESS),
            ("成品槽", COLOR_FINISH),
            ("格外品 / 下腳料", COLOR_OFFGRADE),
            ("公用 / N2", COLOR_UTILITY)
        ]
        
        for idx, (label, color) in enumerate(legend_items):
            # Compact spacing to 1.3 inches to shift legend group to the right
            left = Inches(13.333 - 1.3 * (5 - idx) - 0.4)
            top = Inches(0.35)
            width = Inches(1.2)
            height = Inches(0.25)
            
            # Indicator box
            indicator = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.12), Inches(0.12))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = color
            indicator.line.fill.background()
            
            # Label
            tf = add_clean_textbox(slide, left + Inches(0.18), top - Inches(0.04), width - Inches(0.18), height)
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = FONT_FAMILY
            p.font.size = Pt(8.5)
            p.font.color.rgb = TEXT_MUTED
            
    # Helper: Clean Textbox
    def add_clean_textbox(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        return tf
        
    # Helper: Draw Tank Node Card
    def draw_tank_card(slide, left, top, width, height, title, capacity, tank_type="process", special_note=""):
        # Select color based on type
        if tank_type == "raw":
            theme_color = COLOR_RAW
        elif tank_type == "finish":
            theme_color = COLOR_FINISH
        elif tank_type == "offgrade":
            theme_color = COLOR_OFFGRADE
        elif tank_type == "utility":
            theme_color = COLOR_UTILITY
        else:
            theme_color = COLOR_PROCESS
            
        # Draw background shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = theme_color
        card.line.width = Pt(1.5)
        
        # Text Frame
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        
        # Title (Tank ID)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(1)
        
        # Capacity
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.text = capacity
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.bold = True
        p2.font.color.rgb = theme_color
        
        # Special note or subtext
        if special_note:
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.LEFT
            p3.text = special_note
            p3.font.name = FONT_FAMILY
            p3.font.size = Pt(7.5)
            p3.font.color.rgb = TEXT_MUTED
            p3.space_before = Pt(2)

    # Helper: Draw Group/Stage Bounding Box
    def draw_group_container(slide, left, top, width, height, label, total_cap, group_color):
        # Draw dotted bounding box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = BG_COLOR
        box.line.color.rgb = group_color
        box.line.width = Pt(1.5)
        box.line.dash_style = 2 # Dashed line
        
        # Label container at top left of box
        lbl_w = Inches(len(label) * 0.11 + 0.5)
        lbl_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top - Inches(0.12), lbl_w, Inches(0.24))
        lbl_box.fill.solid()
        lbl_box.fill.fore_color.rgb = BG_COLOR
        lbl_box.line.fill.background()
        
        tf = lbl_box.text_frame
        tf.word_wrap = False
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = group_color
        
        # Total Capacity label at top right of box
        if total_cap:
            cap_box = add_clean_textbox(slide, left + width - Inches(2.2), top + Inches(0.1), Inches(2.0), Inches(0.25))
            p_cap = cap_box.paragraphs[0]
            p_cap.alignment = PP_ALIGN.RIGHT
            p_cap.text = total_cap
            p_cap.font.name = FONT_FAMILY
            p_cap.font.size = Pt(9.5)
            p_cap.font.bold = True
            p_cap.font.color.rgb = TEXT_MUTED

    # Helper: Draw Thick Flow Arrow connecting Stages
    def draw_thick_flow_arrow(slide, start_x, start_y, end_x, end_y, label=""):
        dx = end_x - start_x
        dy = end_y - start_y
        abs_dx = abs(dx)
        abs_dy = abs(dy)
        
        # Thicker conduit line
        if abs_dx >= abs_dy:
            conduit = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, start_y - Inches(0.04), abs_dx, Inches(0.08))
            conduit.fill.solid()
            conduit.fill.fore_color.rgb = TEXT_MUTED
            conduit.line.fill.background()
            
            # Larger arrowhead pointing horizontally
            arrow_w = Inches(0.20)
            arrow_h = Inches(0.24)
            if dx >= 0:
                arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, end_x - arrow_w, start_y - arrow_h/2, arrow_w, arrow_h)
                arrow.rotation = 90
            else:
                arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, end_x, start_y - arrow_h/2, arrow_w, arrow_h)
                arrow.rotation = 270
        else:
            conduit = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x - Inches(0.04), start_y, Inches(0.08), abs_dy)
            conduit.fill.solid()
            conduit.fill.fore_color.rgb = TEXT_MUTED
            conduit.line.fill.background()
            
            # Larger arrowhead pointing vertically
            arrow_w = Inches(0.24)
            arrow_h = Inches(0.20)
            if dy >= 0:
                arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, start_x - arrow_w/2, end_y - arrow_h, arrow_w, arrow_h)
                arrow.rotation = 180
            else:
                arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, start_x - arrow_w/2, end_y, arrow_w, arrow_h)
                arrow.rotation = 0
                
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_MUTED
        arrow.line.fill.background()
        
        # Flow text label
        if label:
            lbl_w = Inches(2.0)
            if abs_dy > abs_dx:
                # Vertical flow line: place text to the right of the line and left-align
                lbl_box = add_clean_textbox(slide, start_x + Inches(0.15), start_y + dy/2 - Inches(0.15), lbl_w, Inches(0.3))
                p = lbl_box.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
            else:
                # Horizontal flow line: place text centered above the line and center-align
                lbl_box = add_clean_textbox(slide, start_x + dx/2 - lbl_w/2, start_y + dy/2 - Inches(0.35), lbl_w, Inches(0.3))
                p = lbl_box.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                
            p.text = label
            p.font.name = FONT_FAMILY
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = TEXT_MUTED

    slide_layout = prs.slide_layouts[6] # Blank Slide
    
    # ------------------ SLIDE 1: IPA 流程 ------------------
    slide1 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide1)
    add_slide_header(slide1, "01. I P A 流程圖", "原料共 2,850 KL / 成品共 770 KL")
    add_slide_legend(slide1)
    
    # 1. Bounding Boxes for Stages
    draw_group_container(slide1, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 2,850 KL", COLOR_RAW)
    draw_group_container(slide1, Inches(4.5), Inches(3.4), Inches(2.4), Inches(1.3), "IPAHW下腳料區", "共 50 KL", COLOR_OFFGRADE)
    draw_group_container(slide1, Inches(4.5), Inches(4.9), Inches(2.4), Inches(1.8), "IPA下腳料區", "共 150 KL", COLOR_OFFGRADE)
    draw_group_container(slide1, Inches(8.5), Inches(1.5), Inches(2.0), Inches(2.5), "Check Tank 待驗 (IPAUPS)", "共 25 KL", COLOR_PROCESS)
    draw_group_container(slide1, Inches(10.8), Inches(1.5), Inches(2.0), Inches(2.5), "IPAUPS成品區", "共 570 KL", COLOR_FINISH)
    draw_group_container(slide1, Inches(8.5), Inches(4.2), Inches(2.0), Inches(2.5), "Check Tank 待驗 (IPAHQ)", "共 60 KL", COLOR_PROCESS)
    draw_group_container(slide1, Inches(10.8), Inches(4.2), Inches(2.0), Inches(2.5), "IPAHQ成品區", "共 200 KL", COLOR_FINISH)
    
    # S1 Central Process card
    proc_s1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.4), Inches(2.0), Inches(0.7))
    proc_s1.fill.solid()
    proc_s1.fill.fore_color.rgb = CARD_BG
    proc_s1.line.color.rgb = COLOR_PROCESS
    proc_s1.line.width = Pt(2.5)
    p_proc_s1 = proc_s1.text_frame.paragraphs[0]
    p_proc_s1.alignment = PP_ALIGN.CENTER
    p_proc_s1.text = "S1製程生產"
    p_proc_s1.font.name = FONT_FAMILY
    p_proc_s1.font.size = Pt(12)
    p_proc_s1.font.bold = True
    p_proc_s1.font.color.rgb = TEXT_WHITE
    
    # S3 Central Process card
    proc_s3 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.35), Inches(2.0), Inches(0.7))
    proc_s3.fill.solid()
    proc_s3.fill.fore_color.rgb = CARD_BG
    proc_s3.line.color.rgb = COLOR_PROCESS
    proc_s3.line.width = Pt(2.5)
    p_proc_s3 = proc_s3.text_frame.paragraphs[0]
    p_proc_s3.alignment = PP_ALIGN.CENTER
    p_proc_s3.text = "S3製程生產"
    p_proc_s3.font.name = FONT_FAMILY
    p_proc_s3.font.size = Pt(12)
    p_proc_s3.font.bold = True
    p_proc_s3.font.color.rgb = TEXT_WHITE
    
    # High-level thick flow arrows
    # Raw to S1 & S3
    draw_thick_flow_arrow(slide1, Inches(3.3), Inches(1.75), Inches(4.5), Inches(1.75))
    draw_thick_flow_arrow(slide1, Inches(3.3), Inches(2.7), Inches(4.5), Inches(2.7))
    
    # S1 Waste to TK-652 (Left bypass to avoid crossing S3)
    pipe_h1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(1.75) - Inches(0.04), Inches(0.6), Inches(0.08))
    pipe_h1.fill.solid()
    pipe_h1.fill.fore_color.rgb = TEXT_MUTED
    pipe_h1.line.fill.background()
    
    pipe_v1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9) - Inches(0.04), Inches(1.75), Inches(0.08), Inches(2.25))
    pipe_v1.fill.solid()
    pipe_v1.fill.fore_color.rgb = TEXT_MUTED
    pipe_v1.line.fill.background()
    
    draw_thick_flow_arrow(slide1, Inches(3.9), Inches(4.0), Inches(4.5), Inches(4.0), "格外品排料")
    
    # S3 Waste to TK-690 (Left bypass)
    pipe_h2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.7), Inches(2.7) - Inches(0.04), Inches(0.8), Inches(0.08))
    pipe_h2.fill.solid()
    pipe_h2.fill.fore_color.rgb = TEXT_MUTED
    pipe_h2.line.fill.background()
    
    pipe_v2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.7) - Inches(0.04), Inches(2.7), Inches(0.08), Inches(2.8))
    pipe_v2.fill.solid()
    pipe_v2.fill.fore_color.rgb = TEXT_MUTED
    pipe_v2.line.fill.background()
    
    draw_thick_flow_arrow(slide1, Inches(3.7), Inches(5.5), Inches(4.5), Inches(5.5), "格外品排料")

    
    # Flow to Check Tank IPA (S1 to top branch)
    draw_thick_flow_arrow(slide1, Inches(6.5), Inches(1.75), Inches(8.5), Inches(2.45), "送待驗")
    
    # Flow to Check Tank IPAHQ (S3 to bottom branch: down then right)
    vert_pipe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5) - Inches(0.04), Inches(2.7), Inches(0.08), Inches(2.55))
    vert_pipe.fill.solid()
    vert_pipe.fill.fore_color.rgb = TEXT_MUTED
    vert_pipe.line.fill.background()
    draw_thick_flow_arrow(slide1, Inches(6.5), Inches(5.25), Inches(8.5), Inches(5.25), "送待驗")
    
    # Flow from Check Tank to Finish Tank
    draw_thick_flow_arrow(slide1, Inches(10.5), Inches(2.45), Inches(10.8), Inches(2.45), "N2 吹掃放行")
    draw_thick_flow_arrow(slide1, Inches(10.5), Inches(5.25), Inches(10.8), Inches(5.25), "N2 吹掃放行")
    
    # Tanks in Raw (Centered single column, x = 1.4, spacing = 0.8)
    draw_tank_card(slide1, Inches(1.4), Inches(1.8), Inches(1.4), Inches(0.65), "TK-602", "500 KL", "raw", "原料槽")
    draw_tank_card(slide1, Inches(1.4), Inches(2.6), Inches(1.4), Inches(0.65), "TK-696", "550 KL", "raw", "原料槽")
    draw_tank_card(slide1, Inches(1.4), Inches(3.4), Inches(1.4), Inches(0.65), "TK-697", "550 KL", "raw", "原料槽")
    draw_tank_card(slide1, Inches(1.4), Inches(4.2), Inches(1.4), Inches(0.65), "TK-693", "250 KL", "raw", "一廠專供原料槽")
    draw_tank_card(slide1, Inches(1.4), Inches(5.0), Inches(1.4), Inches(0.65), "TK-604A", "500 KL", "raw", "精餾塔A")
    draw_tank_card(slide1, Inches(1.4), Inches(5.8), Inches(1.4), Inches(0.65), "TK-604B", "500 KL", "raw", "精餾塔B")
    
    # Tanks in Waste (IPAHW at top, IPA at bottom)
    draw_tank_card(slide1, Inches(5.0), Inches(3.7), Inches(1.4), Inches(0.65), "TK-652", "50 KL", "offgrade", "IPAHW下腳料")
    draw_tank_card(slide1, Inches(5.0), Inches(5.2), Inches(1.4), Inches(0.65), "TK-690", "50 KL", "offgrade", "格外品回收")
    draw_tank_card(slide1, Inches(5.0), Inches(5.95), Inches(1.4), Inches(0.65), "TK-691", "100 KL", "offgrade", "不合格回收")
    
    # Tanks in Check Tank (Centered single column, Standard IPA at top, IPAHQ at bottom)
    draw_tank_card(slide1, Inches(8.8), Inches(2.4), Inches(1.4), Inches(0.7), "TK-614", "25 KL", "process", "待驗槽A")
    draw_tank_card(slide1, Inches(8.8), Inches(4.55), Inches(1.4), Inches(0.7), "TK-675", "30 KL", "process", "待驗槽B")
    draw_tank_card(slide1, Inches(8.8), Inches(5.55), Inches(1.4), Inches(0.7), "TK-676", "30 KL", "process", "待驗槽C")
    
    # Tanks in Finish (Centered single column, Standard IPA at top, IPAHQ at bottom)
    draw_tank_card(slide1, Inches(11.1), Inches(1.8), Inches(1.4), Inches(0.65), "TK-624", "500 KL", "finish", "成品槽")
    draw_tank_card(slide1, Inches(11.1), Inches(2.8), Inches(1.4), Inches(0.65), "TK-672", "70 KL", "finish", "工業級成品槽")
    draw_tank_card(slide1, Inches(11.1), Inches(4.55), Inches(1.4), Inches(0.65), "TK-681", "100 KL", "finish", "電子級成品槽A")
    draw_tank_card(slide1, Inches(11.1), Inches(5.55), Inches(1.4), Inches(0.65), "TK-682", "100 KL", "finish", "電子級成品槽B")
    
    # ------------------ SLIDE 2: EG 流程 ------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "02. E G 流程圖", "原料共 350 KL / 成品共 375 KL")
    add_slide_legend(slide2)
    
    draw_group_container(slide2, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 350 KL", COLOR_RAW)
    draw_group_container(slide2, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 50 KL", COLOR_OFFGRADE)
    draw_group_container(slide2, Inches(8.5), Inches(1.5), Inches(2.0), Inches(5.2), "Check Tank 待驗", "共 55 KL", COLOR_PROCESS)
    draw_group_container(slide2, Inches(10.8), Inches(1.5), Inches(2.0), Inches(5.2), "成品區", "共 375 KL", COLOR_FINISH)
    
    proc2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc2.fill.solid()
    proc2.fill.fore_color.rgb = CARD_BG
    proc2.line.color.rgb = COLOR_PROCESS
    proc2.line.width = Pt(2.5)
    p_proc2 = proc2.text_frame.paragraphs[0]
    p_proc2.alignment = PP_ALIGN.CENTER
    p_proc2.text = "製程生產"
    p_proc2.font.name = FONT_FAMILY
    p_proc2.font.size = Pt(14)
    p_proc2.font.bold = True
    p_proc2.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide2, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide2, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品排料")
    draw_thick_flow_arrow(slide2, Inches(6.5), Inches(2.25), Inches(8.5), Inches(2.25), "調合配料")
    draw_thick_flow_arrow(slide2, Inches(10.5), Inches(2.25), Inches(10.8), Inches(2.25), "放行成品")
    
    # TK-603 and TK-689 are the raw tanks for EG
    draw_tank_card(slide2, Inches(1.4), Inches(2.0), Inches(1.4), Inches(0.8), "TK-603", "250 KL", "raw", "原料槽")
    draw_tank_card(slide2, Inches(1.4), Inches(3.0), Inches(1.4), Inches(0.8), "TK-689", "100 KL", "raw", "原料槽")
    
    # Check tanks: TK-613 & TK-678
    draw_tank_card(slide2, Inches(8.8), Inches(2.0), Inches(1.4), Inches(0.8), "TK-613", "25 KL", "process", "待驗配料槽A")
    draw_tank_card(slide2, Inches(8.8), Inches(3.2), Inches(1.4), Inches(0.8), "TK-678", "30 KL", "process", "待驗配料槽B")
    
    # Finished products
    draw_tank_card(slide2, Inches(11.1), Inches(2.3), Inches(1.4), Inches(0.8), "TK-623", "125 KL", "finish", "成品大槽A")
    draw_tank_card(slide2, Inches(11.1), Inches(3.3), Inches(1.4), Inches(0.8), "TK-692", "250 KL", "finish", "成品大槽B")
    
    # Waste: TK-656
    draw_tank_card(slide2, Inches(5.0), Inches(4.2), Inches(1.4), Inches(0.8), "TK-656", "50 KL", "offgrade", "格外品回收槽")
    
    # ------------------ SLIDE 3: NMP 流程 ------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "03. N M P 流程圖", "原料共 445 KL / 成品共 250 KL")
    add_slide_legend(slide3)
    
    draw_group_container(slide3, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 445 KL", COLOR_RAW)
    draw_group_container(slide3, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "格外品", COLOR_OFFGRADE)
    draw_group_container(slide3, Inches(8.5), Inches(1.5), Inches(2.0), Inches(5.2), "Check Tank 待驗", "共 25 KL", COLOR_PROCESS)
    draw_group_container(slide3, Inches(10.8), Inches(1.5), Inches(2.0), Inches(5.2), "成品區", "共 250 KL", COLOR_FINISH)
    
    proc3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc3.fill.solid()
    proc3.fill.fore_color.rgb = CARD_BG
    proc3.line.color.rgb = COLOR_PROCESS
    proc3.line.width = Pt(2.5)
    p_proc3 = proc3.text_frame.paragraphs[0]
    p_proc3.alignment = PP_ALIGN.CENTER
    p_proc3.text = "製程生產"
    p_proc3.font.name = FONT_FAMILY
    p_proc3.font.size = Pt(14)
    p_proc3.font.bold = True
    p_proc3.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide3, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide3, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品排料")
    draw_thick_flow_arrow(slide3, Inches(6.5), Inches(2.25), Inches(8.5), Inches(2.25), "送待驗")
    draw_thick_flow_arrow(slide3, Inches(10.5), Inches(2.25), Inches(10.8), Inches(2.25), "N2 吹掃")
    
    # Raw tanks stacked vertically
    draw_tank_card(slide3, Inches(1.4), Inches(1.8), Inches(1.4), Inches(0.65), "TK-632", "250 KL", "raw", "原料槽A")
    draw_tank_card(slide3, Inches(1.4), Inches(2.6), Inches(1.4), Inches(0.65), "TK-633", "125 KL", "raw", "原料槽B")
    draw_tank_card(slide3, Inches(1.4), Inches(3.4), Inches(1.4), Inches(0.65), "TK-671", "70 KL", "raw", "原料槽C")
    
    # Check Tank (TK-611)
    draw_tank_card(slide3, Inches(8.8), Inches(2.5), Inches(1.4), Inches(0.8), "TK-611", "25 KL", "process", "製程待驗槽")
    
    # Finished Product (TK-621)
    draw_tank_card(slide3, Inches(11.1), Inches(2.5), Inches(1.4), Inches(0.8), "TK-621", "250 KL", "finish", "成品大槽")
    
    # Waste (IBC桶)
    draw_tank_card(slide3, Inches(5.0), Inches(3.8), Inches(1.4), Inches(0.65), "IBC桶", "1T", "offgrade", "格外品收集")
    
    # ------------------ SLIDE 4: CPNE4 流程 ------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "04. C P N E 4 流程圖", "原料共 500 KL / 成品共 50 KL")
    add_slide_legend(slide4)
    
    draw_group_container(slide4, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 500 KL", COLOR_RAW)
    draw_group_container(slide4, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 50 KL", COLOR_OFFGRADE)
    draw_group_container(slide4, Inches(8.5), Inches(1.5), Inches(2.0), Inches(5.2), "Check Tank 待驗", "共 30 KL", COLOR_PROCESS)
    draw_group_container(slide4, Inches(10.8), Inches(1.5), Inches(2.0), Inches(5.2), "成品區", "共 50 KL", COLOR_FINISH)
    
    proc4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc4.fill.solid()
    proc4.fill.fore_color.rgb = CARD_BG
    proc4.line.color.rgb = COLOR_PROCESS
    proc4.line.width = Pt(2.5)
    p_proc4 = proc4.text_frame.paragraphs[0]
    p_proc4.alignment = PP_ALIGN.CENTER
    p_proc4.text = "製程生產"
    p_proc4.font.name = FONT_FAMILY
    p_proc4.font.size = Pt(14)
    p_proc4.font.bold = True
    p_proc4.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide4, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide4, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品排料")
    draw_thick_flow_arrow(slide4, Inches(6.5), Inches(2.25), Inches(8.5), Inches(2.25), "送待驗")
    draw_thick_flow_arrow(slide4, Inches(10.5), Inches(2.25), Inches(10.8), Inches(2.25), "N2 吹掃")
    
    # Draw vertical split pipeline down from between Process and Check-Tank (starts at x=7.5, y=2.25)
    conduit_split = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5) - Inches(0.04), Inches(2.25), Inches(0.08), Inches(1.55))
    conduit_split.fill.solid()
    conduit_split.fill.fore_color.rgb = TEXT_MUTED
    conduit_split.line.fill.background()
    
    # Draw horizontal split arrow into Finished bottom TK-667
    draw_thick_flow_arrow(slide4, Inches(7.5), Inches(3.8), Inches(10.8), Inches(3.8), "直接放行")
    
    # Raw tanks: TK-684, TK-685, TK-686, TK-687, TK-688 (100 KL each)
    draw_tank_card(slide4, Inches(1.4), Inches(1.8), Inches(1.4), Inches(0.65), "TK-684", "100 KL", "raw", "原料槽A")
    draw_tank_card(slide4, Inches(1.4), Inches(2.6), Inches(1.4), Inches(0.65), "TK-685", "100 KL", "raw", "原料槽B")
    draw_tank_card(slide4, Inches(1.4), Inches(3.4), Inches(1.4), Inches(0.65), "TK-686", "100 KL", "raw", "原料槽C")
    draw_tank_card(slide4, Inches(1.4), Inches(4.2), Inches(1.4), Inches(0.65), "TK-687", "100 KL", "raw", "原料槽D")
    draw_tank_card(slide4, Inches(1.4), Inches(5.0), Inches(1.4), Inches(0.65), "TK-688", "100 KL", "raw", "原料槽E")
    
    # Waste: TK-655 (50 KL)
    draw_tank_card(slide4, Inches(5.0), Inches(4.5), Inches(1.4), Inches(0.8), "TK-655", "50 KL", "offgrade", "格外品大槽")
    
    # Check Tank: TK-677 (30 KL)
    draw_tank_card(slide4, Inches(8.8), Inches(2.5), Inches(1.4), Inches(0.8), "TK-677", "30 KL", "process", "待驗槽")
    
    # Finished: TK-680 (50 KL) centered vertically
    draw_tank_card(slide4, Inches(11.1), Inches(2.5), Inches(1.4), Inches(0.8), "TK-680", "50 KL", "finish", "成品大槽A")
    
    # ------------------ SLIDE 5: CPNE3R / CPNE3 / 2CPN-P1 ------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide5)
    add_slide_header(slide5, "05. CPNE3R & CPNE3 & CPNE3T & 2CPN-P1 流程圖", "原料共 80 KL / CPNE3成品共 295 KL / CPNE3T成品共 100 KL")
    add_slide_legend(slide5)
    
    # 5-Column Containers
    draw_group_container(slide5, Inches(0.9), Inches(1.5), Inches(2.0), Inches(5.2), "原料區", "共 80 KL", COLOR_RAW)
    draw_group_container(slide5, Inches(3.3), Inches(1.5), Inches(2.0), Inches(5.2), "半成品 CPN-P1", "共 500 KL", COLOR_PROCESS)
    draw_group_container(slide5, Inches(5.7), Inches(1.5), Inches(2.0), Inches(5.2), "半成品 CPNE3R", "共 250 KL", COLOR_PROCESS)
    draw_group_container(slide5, Inches(8.1), Inches(1.5), Inches(2.0), Inches(5.2), "成品 CPNE3", "共 295 KL", COLOR_FINISH)
    draw_group_container(slide5, Inches(10.5), Inches(1.5), Inches(2.0), Inches(5.2), "成品 CPNE3T", "共 100 KL", COLOR_FINISH)
    
    # Flow Arrows between columns
    # Raw to 2CPN-P1 (Left 2.9 to 3.3)
    draw_thick_flow_arrow(slide5, Inches(2.9), Inches(2.95), Inches(3.3), Inches(2.95), "前處理")
    # 2CPN-P1 to CPNE3R (Left 5.3 to 5.7)
    draw_thick_flow_arrow(slide5, Inches(5.3), Inches(2.95), Inches(5.7), Inches(2.95), "製程")
    # CPNE3R to CPNE3 (Left 7.7 to 8.1)
    draw_thick_flow_arrow(slide5, Inches(7.7), Inches(2.95), Inches(8.1), Inches(2.95), "製程")
    # CPNE3 to CPNE3T (Left 10.1 to 10.5)
    draw_thick_flow_arrow(slide5, Inches(10.1), Inches(2.95), Inches(10.5), Inches(2.95), "製程")
    
    # Extra products (drawing vertical waste arrows and cards below the conduits)
    # Gap 1: x = 3.1. Vertical arrow down to IBC桶
    draw_thick_flow_arrow(slide5, Inches(3.1), Inches(3.1), Inches(3.1), Inches(4.3), "格外品")
    draw_tank_card(slide5, Inches(2.5), Inches(4.4), Inches(1.2), Inches(0.65), "IBC桶", "1T", "offgrade", "格外品前處理")
    
    # Gap 2: x = 5.5. Vertical arrow down to CPNR (TK-660 70KL)
    draw_thick_flow_arrow(slide5, Inches(5.5), Inches(3.1), Inches(5.5), Inches(4.3), "格外品")
    draw_tank_card(slide5, Inches(4.9), Inches(4.4), Inches(1.2), Inches(0.65), "TK-660", "70 KL", "offgrade", "格外品回收槽")
    
    # Gap 3: x = 7.9. Vertical arrow down to CPNE3R1 (TK-655 50KL)
    draw_thick_flow_arrow(slide5, Inches(7.9), Inches(3.1), Inches(7.9), Inches(4.3), "格外品")
    draw_tank_card(slide5, Inches(7.3), Inches(4.4), Inches(1.2), Inches(0.65), "TK-655", "50 KL", "offgrade", "格外品大槽")
    
    # Gap 4: x = 10.3. Vertical arrow down to CPNE3T offgrade (TK-655 50KL)
    draw_thick_flow_arrow(slide5, Inches(10.3), Inches(3.1), Inches(10.3), Inches(4.3), "格外品")
    draw_tank_card(slide5, Inches(9.7), Inches(4.4), Inches(1.2), Inches(0.65), "TK-655", "50 KL", "offgrade", "格外品大槽")
    
    # Raw Tank Cards
    draw_tank_card(slide5, Inches(1.2), Inches(2.0), Inches(1.4), Inches(0.8), "TK-643", "40 KL", "raw", "攪拌調合罐A")
    draw_tank_card(slide5, Inches(1.2), Inches(3.0), Inches(1.4), Inches(0.8), "TK-645", "40 KL", "raw", "攪拌調合罐B")
    
    # 2CPN-P1 Card
    draw_tank_card(slide5, Inches(3.6), Inches(2.55), Inches(1.4), Inches(0.9), "TK-634", "500 KL", "raw", "半成品大槽")
    
    # CPNE3R Card
    draw_tank_card(slide5, Inches(6.0), Inches(2.55), Inches(1.4), Inches(0.9), "TK-601", "250 KL", "process", "半成品反應槽")
    
    # CPNE3 Cards (grouped as stacked cards to fit perfectly)
    draw_tank_card(slide5, Inches(8.4), Inches(2.0), Inches(1.4), Inches(0.8), "TK-661/662/664", "65 KL / 槽", "finish", "成品大槽組 (3槽)")
    draw_tank_card(slide5, Inches(8.4), Inches(3.2), Inches(1.4), Inches(0.8), "TK-667/668", "50 KL / 槽", "finish", "成品大槽組 (2槽)")
    
    # CPNE3T Card
    draw_tank_card(slide5, Inches(10.8), Inches(2.55), Inches(1.4), Inches(0.9), "TK-683", "100 KL", "finish", "成品放行大槽")
    
    # ------------------ SLIDE 6: ACT 流程 ------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide6)
    add_slide_header(slide6, "06. A C T 流程圖", "原料共 70 KL / 成品共 250 KL")
    add_slide_legend(slide6)
    
    draw_group_container(slide6, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 70 KL", COLOR_RAW)
    draw_group_container(slide6, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 25 KL", COLOR_OFFGRADE)
    draw_group_container(slide6, Inches(8.5), Inches(1.5), Inches(2.0), Inches(5.2), "Check Tank 待驗", "共 25 KL", COLOR_PROCESS)
    draw_group_container(slide6, Inches(10.8), Inches(1.5), Inches(2.0), Inches(5.2), "成品區", "共 250 KL", COLOR_FINISH)
    
    proc6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc6.fill.solid()
    proc6.fill.fore_color.rgb = CARD_BG
    proc6.line.color.rgb = COLOR_PROCESS
    proc6.line.width = Pt(2.5)
    p_proc6 = proc6.text_frame.paragraphs[0]
    p_proc6.alignment = PP_ALIGN.CENTER
    p_proc6.text = "製程生產"
    p_proc6.font.name = FONT_FAMILY
    p_proc6.font.size = Pt(14)
    p_proc6.font.bold = True
    p_proc6.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide6, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide6, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品回流")
    draw_thick_flow_arrow(slide6, Inches(6.5), Inches(2.25), Inches(8.5), Inches(2.25), "送精製")
    draw_thick_flow_arrow(slide6, Inches(10.5), Inches(2.25), Inches(10.8), Inches(2.25), "放行成品")
    
    draw_tank_card(slide6, Inches(1.4), Inches(2.5), Inches(1.4), Inches(0.8), "TK-657", "70 KL", "raw", "ACT 原料槽")
    draw_tank_card(slide6, Inches(8.8), Inches(2.5), Inches(1.4), Inches(0.8), "TK-612", "25 KL", "process", "精餾製程釜 (N2 吹掃)")
    draw_tank_card(slide6, Inches(11.1), Inches(2.5), Inches(1.4), Inches(0.8), "TK-622", "250 KL", "finish", "成品大槽")
    draw_tank_card(slide6, Inches(5.0), Inches(4.5), Inches(1.4), Inches(0.8), "1T桶", "1T", "offgrade", "格外品回流")
    
    # ------------------ SLIDE 7: EBR 流程 ------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide7)
    add_slide_header(slide7, "07. E B R 流程圖", "原料共 450 KL / 成品共 800 KL")
    add_slide_legend(slide7)
    
    draw_group_container(slide7, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 450 KL", COLOR_RAW)
    draw_group_container(slide7, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 200 KL", COLOR_OFFGRADE)
    draw_group_container(slide7, Inches(8.5), Inches(1.5), Inches(2.0), Inches(5.2), "成品區", "共 800 KL", COLOR_FINISH)
    
    proc7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc7.fill.solid()
    proc7.fill.fore_color.rgb = CARD_BG
    proc7.line.color.rgb = COLOR_PROCESS
    proc7.line.width = Pt(2.5)
    p_proc7 = proc7.text_frame.paragraphs[0]
    p_proc7.alignment = PP_ALIGN.CENTER
    p_proc7.text = "製程生產"
    p_proc7.font.name = FONT_FAMILY
    p_proc7.font.size = Pt(14)
    p_proc7.font.bold = True
    p_proc7.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide7, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide7, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品排料")
    draw_thick_flow_arrow(slide7, Inches(6.5), Inches(2.25), Inches(8.5), Inches(2.25), "成品放行")
    
    draw_tank_card(slide7, Inches(1.4), Inches(2.5), Inches(1.4), Inches(0.8), "TKC01-09", "50 KL / 槽", "raw", "平行原料組槽")
    
    draw_tank_card(slide7, Inches(5.0), Inches(3.8), Inches(1.4), Inches(0.65), "TK-646", "50 KL", "offgrade", "格外品槽A")
    draw_tank_card(slide7, Inches(5.0), Inches(4.6), Inches(1.4), Inches(0.65), "TK-647", "50 KL", "offgrade", "格外品槽B")
    draw_tank_card(slide7, Inches(5.0), Inches(5.4), Inches(1.4), Inches(0.65), "TK-648", "100 KL", "offgrade", "格外品槽C")
    
    # Finished tanks (TK-603 removed, only TK-699 & TK-631)
    draw_tank_card(slide7, Inches(8.8), Inches(2.3), Inches(1.4), Inches(0.8), "TK-699", "550 KL", "finish", "成品大槽A")
    draw_tank_card(slide7, Inches(8.8), Inches(3.5), Inches(1.4), Inches(0.8), "TK-631", "250 KL", "finish", "成品大槽B")
    
    # ------------------ SLIDE 8: HEA-R 流程 ------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide8)
    add_slide_header(slide8, "08. HEA-R 流程圖", "原料共 50 KL / 成品A共 140 KL / 成品B共 125 KL")
    add_slide_legend(slide8)
    
    draw_group_container(slide8, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 50 KL", COLOR_RAW)
    draw_group_container(slide8, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 1T", COLOR_OFFGRADE)
    draw_group_container(slide8, Inches(10.8), Inches(1.5), Inches(2.0), Inches(2.5), "成品區 B (DPM-B1)", "共 125 KL", COLOR_FINISH)
    
    # Finished A is inside a separate large bounding box at the bottom right for MIXED ETHER
    draw_group_container(slide8, Inches(8.5), Inches(4.2), Inches(4.3), Inches(2.5), "成品區 A (MIXED ETHER)", "共 140 KL", COLOR_FINISH)
    
    proc8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(2.0), Inches(0.9))
    proc8.fill.solid()
    proc8.fill.fore_color.rgb = CARD_BG
    proc8.line.color.rgb = COLOR_PROCESS
    proc8.line.width = Pt(2.5)
    p_proc8 = proc8.text_frame.paragraphs[0]
    p_proc8.alignment = PP_ALIGN.CENTER
    p_proc8.text = "製程生產"
    p_proc8.font.name = FONT_FAMILY
    p_proc8.font.size = Pt(14)
    p_proc8.font.bold = True
    p_proc8.font.color.rgb = TEXT_WHITE
    
    draw_thick_flow_arrow(slide8, Inches(3.3), Inches(2.25), Inches(4.5), Inches(2.25))
    draw_thick_flow_arrow(slide8, Inches(5.5), Inches(2.7), Inches(5.5), Inches(3.4), "格外品排料")
    draw_thick_flow_arrow(slide8, Inches(6.5), Inches(2.25), Inches(10.8), Inches(2.25), "成品放行")
    
    # Draw vertical split pipeline down from between Process and Finished B (starts at x=7.5, y=2.25)
    conduit_split = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5) - Inches(0.04), Inches(2.25), Inches(0.08), Inches(3.15))
    conduit_split.fill.solid()
    conduit_split.fill.fore_color.rgb = TEXT_MUTED
    conduit_split.line.fill.background()
    
    # Draw horizontal split arrow into Finished A
    draw_thick_flow_arrow(slide8, Inches(7.5), Inches(5.4), Inches(8.5), Inches(5.4), "成品放行")
    
    # Raw tanks: Only TK654 (50 KL)
    draw_tank_card(slide8, Inches(1.4), Inches(2.5), Inches(1.4), Inches(0.8), "TK654", "50 KL", "raw", "原料槽")
    
    # Waste: IBC Tank (1T)
    draw_tank_card(slide8, Inches(5.0), Inches(4.5), Inches(1.4), Inches(0.8), "IBC桶", "1T", "offgrade", "格外品收集")
    
    # Finished B: TK-641 (125 KL)
    draw_tank_card(slide8, Inches(11.1), Inches(2.2), Inches(1.4), Inches(0.8), "TK-641", "125 KL", "finish", "成品 DPM-B1")
    
    # Finished A: TK-673 (70 KL), TK-674 (70 KL) side-by-side
    draw_tank_card(slide8, Inches(9.2), Inches(5.0), Inches(1.4), Inches(0.8), "TK-673", "70 KL", "finish", "MIXED ETHER 成品B")
    # ------------------ SLIDE 9: 崙尾 IPA HQ 流程 ------------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide9)
    add_slide_header(slide9, "09. 崙尾 I P A HQ 流程圖", "原料共 2,000 KL / 成品共 400 KL")
    add_slide_legend(slide9)
    
    draw_group_container(slide9, Inches(0.9), Inches(1.5), Inches(2.4), Inches(5.2), "原料區", "共 2,000 KL", COLOR_RAW)
    draw_group_container(slide9, Inches(4.5), Inches(3.4), Inches(2.4), Inches(3.3), "下腳料區", "共 300 KL", COLOR_OFFGRADE)
    draw_group_container(slide9, Inches(8.5), Inches(1.5), Inches(2.0), Inches(2.5), "Check Tank 待驗 (S4)", "共 60 KL", COLOR_PROCESS)
    draw_group_container(slide9, Inches(10.8), Inches(1.5), Inches(2.0), Inches(2.5), "成品區 (S4)", "共 200 KL", COLOR_FINISH)
    draw_group_container(slide9, Inches(8.5), Inches(4.2), Inches(2.0), Inches(2.5), "Check Tank 待驗 (S5)", "共 60 KL", COLOR_PROCESS)
    draw_group_container(slide9, Inches(10.8), Inches(4.2), Inches(2.0), Inches(2.5), "成品區 (S5)", "共 200 KL", COLOR_FINISH)
    
    # S4 Central Process card
    proc9_s4 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.4), Inches(2.0), Inches(0.7))
    proc9_s4.fill.solid()
    proc9_s4.fill.fore_color.rgb = CARD_BG
    proc9_s4.line.color.rgb = COLOR_PROCESS
    proc9_s4.line.width = Pt(2.5)
    p_proc9_s4 = proc9_s4.text_frame.paragraphs[0]
    p_proc9_s4.alignment = PP_ALIGN.CENTER
    p_proc9_s4.text = "S4製程生產"
    p_proc9_s4.font.name = FONT_FAMILY
    p_proc9_s4.font.size = Pt(12)
    p_proc9_s4.font.bold = True
    p_proc9_s4.font.color.rgb = TEXT_WHITE
    
    # S5 Central Process card
    proc9_s5 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.35), Inches(2.0), Inches(0.7))
    proc9_s5.fill.solid()
    proc9_s5.fill.fore_color.rgb = CARD_BG
    proc9_s5.line.color.rgb = COLOR_PROCESS
    proc9_s5.line.width = Pt(2.5)
    p_proc9_s5 = proc9_s5.text_frame.paragraphs[0]
    p_proc9_s5.alignment = PP_ALIGN.CENTER
    p_proc9_s5.text = "S5製程生產"
    p_proc9_s5.font.name = FONT_FAMILY
    p_proc9_s5.font.size = Pt(12)
    p_proc9_s5.font.bold = True
    p_proc9_s5.font.color.rgb = TEXT_WHITE
    
    # Raw to S4 & S5
    draw_thick_flow_arrow(slide9, Inches(3.3), Inches(1.75), Inches(4.5), Inches(1.75))
    draw_thick_flow_arrow(slide9, Inches(3.3), Inches(2.7), Inches(4.5), Inches(2.7))
    
    # S4 Waste to Waste area (Left bypass to avoid crossing S5)
    pipe_h9_s4 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(1.75) - Inches(0.04), Inches(0.6), Inches(0.08))
    pipe_h9_s4.fill.solid()
    pipe_h9_s4.fill.fore_color.rgb = TEXT_MUTED
    pipe_h9_s4.line.fill.background()
    
    pipe_v9_s4 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9) - Inches(0.04), Inches(1.75), Inches(0.08), Inches(2.5))
    pipe_v9_s4.fill.solid()
    pipe_v9_s4.fill.fore_color.rgb = TEXT_MUTED
    pipe_v9_s4.line.fill.background()
    
    draw_thick_flow_arrow(slide9, Inches(3.9), Inches(4.25), Inches(4.5), Inches(4.25))
    
    # S5 Waste to Waste area (Left bypass)
    pipe_h9_s5 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.7), Inches(2.7) - Inches(0.04), Inches(0.8), Inches(0.08))
    pipe_h9_s5.fill.solid()
    pipe_h9_s5.fill.fore_color.rgb = TEXT_MUTED
    pipe_h9_s5.line.fill.background()
    
    pipe_v9_s5 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.7) - Inches(0.04), Inches(2.7), Inches(0.08), Inches(2.55))
    pipe_v9_s5.fill.solid()
    pipe_v9_s5.fill.fore_color.rgb = TEXT_MUTED
    pipe_v9_s5.line.fill.background()
    
    draw_thick_flow_arrow(slide9, Inches(3.7), Inches(5.25), Inches(4.5), Inches(5.25))

    # Add a single label in the gap between S5 and waste group
    lbl_box = add_clean_textbox(slide9, Inches(2.2), Inches(3.12), Inches(1.0), Inches(0.3))
    p = lbl_box.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = "格外品排料"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED

    # Draw a thin red connector arrow pointing from the label to the pipeline
    conn_arrow = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.35), Inches(3.22), Inches(0.28), Inches(0.1))
    conn_arrow.fill.solid()
    conn_arrow.fill.fore_color.rgb = RGBColor(244, 63, 94) # var(--color-offgrade)
    conn_arrow.line.fill.background()

    
    # S4 output to S4 Check Tank
    draw_thick_flow_arrow(slide9, Inches(6.5), Inches(1.75), Inches(8.5), Inches(2.45), "待檢驗")
    
    # S5 output to S5 Check Tank (branch: down then right)
    vert_pipe9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5) - Inches(0.04), Inches(2.7), Inches(0.08), Inches(2.55))
    vert_pipe9.fill.solid()
    vert_pipe9.fill.fore_color.rgb = TEXT_MUTED
    vert_pipe9.line.fill.background()
    draw_thick_flow_arrow(slide9, Inches(6.5), Inches(5.25), Inches(8.5), Inches(5.25), "待檢驗")
    
    # Check Tank to Finish Tank
    draw_thick_flow_arrow(slide9, Inches(10.5), Inches(2.45), Inches(10.8), Inches(2.45), "成品放行")
    draw_thick_flow_arrow(slide9, Inches(10.5), Inches(5.25), Inches(10.8), Inches(5.25), "成品放行")
    
    # Raw Tanks
    draw_tank_card(slide9, Inches(1.4), Inches(2.2), Inches(1.4), Inches(0.8), "TK-617", "1000 KL", "raw", "原料大槽A")
    draw_tank_card(slide9, Inches(1.4), Inches(3.2), Inches(1.4), Inches(0.8), "TK-618", "1000 KL", "raw", "原料大槽B")
    
    # Waste Tanks
    draw_tank_card(slide9, Inches(5.0), Inches(3.9), Inches(1.4), Inches(0.7), "TK-611", "200 KL", "offgrade", "格外品槽A")
    draw_tank_card(slide9, Inches(5.0), Inches(4.9), Inches(1.4), Inches(0.7), "TK-613", "100 KL", "offgrade", "格外品槽B")
    
    # Check Tanks (S4 at top, S5 at bottom)
    draw_tank_card(slide9, Inches(8.8), Inches(1.75), Inches(1.4), Inches(0.65), "TK-601", "30 KL", "process", "待驗槽A")
    draw_tank_card(slide9, Inches(8.8), Inches(2.55), Inches(1.4), Inches(0.65), "TK-602", "30 KL", "process", "待驗槽B")
    draw_tank_card(slide9, Inches(8.8), Inches(4.55), Inches(1.4), Inches(0.65), "TK-603", "30 KL", "process", "待驗槽C")
    draw_tank_card(slide9, Inches(8.8), Inches(5.55), Inches(1.4), Inches(0.65), "TK-604", "30 KL", "process", "待驗槽D")
    
    # Finished Tanks (S4 at top, S5 at bottom)
    draw_tank_card(slide9, Inches(11.1), Inches(1.75), Inches(1.4), Inches(0.65), "TK-605", "100 KL", "finish", "成品大槽A")
    draw_tank_card(slide9, Inches(11.1), Inches(2.55), Inches(1.4), Inches(0.65), "TK-606", "100 KL", "finish", "成品大槽B")
    draw_tank_card(slide9, Inches(11.1), Inches(4.55), Inches(1.4), Inches(0.65), "TK-607", "1000 KL", "finish", "成品大槽C")
    draw_tank_card(slide9, Inches(11.1), Inches(5.55), Inches(1.4), Inches(0.65), "TK-608", "100 KL", "finish", "成品大槽D")
    
    # Save Presentation (Saved as _v2 to bypass file lock if the user has _產品流程圖 open in PowerPoint)
    output_path = os.path.join(r"g:\我的雲端硬碟\GOOGLE ANGET", "20260511產品流程圖_優化版_v2.pptx")
    prs.save(output_path)
    print(f"Presentation successfully created and saved to {output_path}!")

if __name__ == "__main__":
    create_optimized_ppt()
