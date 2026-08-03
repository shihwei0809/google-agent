import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_optimized_sop_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # ------------------ DESIGN SYSTEM COLORS ------------------
    BG_COLOR = RGBColor(9, 13, 22)         # #090D16 Premium Dark Space Gray
    TEXT_WHITE = RGBColor(248, 250, 252)   # #F8FAFC Pure white text
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 Cool gray subtext
    
    COLOR_APPLY = RGBColor(56, 189, 248)   # #38BDF8 Sky Blue for Apply/Monitor
    COLOR_QA = RGBColor(168, 85, 247)      # #A855F7 Purple for QA Appraisal
    COLOR_PROD = RGBColor(16, 185, 129)    # #10B981 Emerald Green for Production
    COLOR_SALES = RGBColor(245, 158, 11)   # #F59E0B Amber for Sales
    COLOR_WASTE = RGBColor(239, 68, 68)    # #EF4444 Rose Red for GM/Waste
    
    CARD_BG = RGBColor(22, 28, 45)         # #161C2D Dark slate container boxes
    CARD_BG_HOVER = RGBColor(30, 41, 59)   # #1E293B High contrast slate
    
    # Fonts
    FONT_FAMILY = "Microsoft JhengHei"
    
    # Helper: Set Slide Dark Background
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Helper: Create Header
    def add_slide_header(slide, title_text, subtitle_text=""):
        # Decorative top border
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_PROD
        top_bar.line.fill.background()
        
        # Header text frame
        header_tf = add_clean_textbox(slide, Inches(0.6), Inches(0.3), Inches(7.5), Inches(0.9))
        p = header_tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        if subtitle_text:
          p_sub = header_tf.add_paragraph()
          p_sub.text = subtitle_text
          p_sub.font.name = FONT_FAMILY
          p_sub.font.size = Pt(10)
          p_sub.font.bold = True
          p_sub.font.color.rgb = COLOR_APPLY
            
    # Helper: Add Slide Legend
    def add_slide_legend(slide):
        legend_items = [
            ("申請/廠務", COLOR_APPLY),
            ("品保鑑定", COLOR_QA),
            ("生產再製", COLOR_PROD),
            ("業務銷售", COLOR_SALES),
            ("核准報廢", COLOR_WASTE)
        ]
        
        for idx, (label, color) in enumerate(legend_items):
            left = Inches(13.333 - 1.25 * (5 - idx) - 0.6)
            top = Inches(0.4)
            width = Inches(1.15)
            height = Inches(0.25)
            
            # Indicator box
            indicator = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.12), Inches(0.12))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = color
            indicator.line.fill.background()
            
            # Label
            tf = add_clean_textbox(slide, left + Inches(0.18), top - Inches(0.04), width - Inches(0.18), height)
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = FONT_FAMILY
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_MUTED
            
    # Helper: Clean Textbox
    def add_clean_textbox(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        return tf
        
    # Helper: Draw Card
    def draw_sop_card(slide, left, top, width, height, title, subtitle="", detail="", type_color=COLOR_APPLY, is_diamond=False, is_doc=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        if is_diamond:
            shape_type = MSO_SHAPE.DIAMOND
        elif is_doc:
            shape_type = MSO_SHAPE.RECTANGLE # Document style border
            
        card = slide.shapes.add_shape(shape_type, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = type_color
        card.line.width = Pt(1.5)
        
        # Text Frame
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        
        # Title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Subtitle
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = subtitle
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(8)
            p2.font.bold = True
            p2.font.color.rgb = type_color
            p2.space_before = Pt(1)
            
        # Detail / Code
        if detail:
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            p3.text = detail
            p3.font.name = FONT_FAMILY
            p3.font.size = Pt(7)
            p3.font.color.rgb = TEXT_MUTED
            p3.space_before = Pt(2)
            
        return card

    # Helper: Draw Group Container (Swimlane Column)
    def draw_swimlane_col(slide, left, top, width, height, title, subtitle, group_color):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(17, 24, 39)
        box.line.color.rgb = BG_COLOR
        box.line.width = Pt(1.0)
        
        # Header Rectangle inside
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.45))
        header.fill.solid()
        header.fill.fore_color.rgb = CARD_BG
        header.line.fill.background()
        
        tf = header.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p_sub = tf.add_paragraph()
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.text = subtitle
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(7.5)
        p_sub.font.color.rgb = group_color

    # Helper: Draw Conduit Line (Thick Pipe)
    def draw_conduit_pipe(slide, start_x, start_y, end_x, end_y, pipe_color=RGBColor(100, 116, 139), dash_style=1, label=""):
        dx = end_x - start_x
        dy = end_y - start_y
        abs_dx = abs(dx)
        abs_dy = abs(dy)
        
        # Draw base conduit
        if abs_dx >= abs_dy: # horizontal
            x = min(start_x, end_x)
            conduit = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y - Inches(0.02), abs_dx, Inches(0.04))
        else: # vertical
            y = min(start_y, end_y)
            conduit = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x - Inches(0.02), y, Inches(0.04), abs_dy)
            
        conduit.fill.solid()
        conduit.fill.fore_color.rgb = pipe_color
        conduit.line.fill.background()
        
        # Label placement
        if label:
            lbl_w = Inches(2.0)
            if abs_dy > abs_dx:
                lbl_box = add_clean_textbox(slide, start_x + Inches(0.1), start_y + dy/2 - Inches(0.12), lbl_w, Inches(0.25))
                p = lbl_box.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
            else:
                lbl_box = add_clean_textbox(slide, start_x + dx/2 - lbl_w/2, start_y - Inches(0.28), lbl_w, Inches(0.25))
                p = lbl_box.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                
            p.text = label
            p.font.name = FONT_FAMILY
            p.font.size = Pt(7.5)
            p.font.bold = True
            p.font.color.rgb = pipe_color

    slide_layout = prs.slide_layouts[6] # Blank Slide
    
    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide1)
    
    # Bottom subtle geometric graphic block
    dec_block = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7))
    dec_block.fill.solid()
    dec_block.fill.fore_color.rgb = RGBColor(17, 24, 39)
    dec_block.line.fill.background()
    
    # Title decorative accent line
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.3), Inches(0.15), Inches(2.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_PROD
    accent.line.fill.background()
    
    # Main text block
    main_tf = add_clean_textbox(slide1, Inches(1.6), Inches(2.2), Inches(10.5), Inches(2.5))
    p1 = main_tf.paragraphs[0]
    p1.text = "格外品、久滯品處理流程優化"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(8)
    
    p2 = main_tf.add_paragraph()
    p2.text = "S.O.P. 標準泳道管制與教育訓練指引"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_APPLY
    p2.space_after = Pt(20)
    
    # Subtext on footer block
    foot_tf = add_clean_textbox(slide1, Inches(1.6), Inches(6.1), Inches(10), Inches(1.0))
    pf1 = foot_tf.paragraphs[0]
    pf1.text = "鴻勝化學工業股份有限公司 · 廠務暨品保管理部"
    pf1.font.name = FONT_FAMILY
    pf1.font.size = Pt(11)
    pf1.font.bold = True
    pf1.font.color.rgb = TEXT_WHITE
    
    pf2 = foot_tf.add_paragraph()
    pf2.text = "本流程基於公司《格外品處理、儲存管理辦法》優化編制，提供 🟢回流再製 / 🔵折價銷售 / 🔴廢棄報廢 三大標準路徑。"
    pf2.font.name = FONT_FAMILY
    pf2.font.size = Pt(9.5)
    pf2.font.color.rgb = TEXT_MUTED
    pf2.space_before = Pt(4)
    
    
    # ==========================================
    # SLIDE 2: The Swimlane Flowchart Overview
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "01. 格外品、久滯品流程圖", "全體流程單向向下流動，清晰分工無折返路徑")
    add_slide_legend(slide2)
    
    # 1. Bounding Bins for Swimlanes
    col_w = Inches(2.25)
    col_h = Inches(5.3)
    col_top = Inches(1.5)
    
    draw_swimlane_col(slide2, Inches(0.6), col_top, col_w, col_h, "申請與監控單位", "申請部門 & 廠務管制", COLOR_APPLY)
    draw_swimlane_col(slide2, Inches(3.0), col_top, col_w, col_h, "品質鑑定單位", "品保部 (QA)", COLOR_QA)
    draw_swimlane_col(slide2, Inches(5.4), col_top, col_w, col_h, "再製執行單位", "生產現場主管/各課室", COLOR_PROD)
    draw_swimlane_col(slide2, Inches(7.8), col_top, col_w, col_h, "價值評估單位", "業務部 (Sales)", COLOR_SALES)
    draw_swimlane_col(slide2, Inches(10.2), col_top, col_w, col_h, "核准與報廢處置", "(副)總經理 / 環安單位", COLOR_WASTE)
    
    # 2. Draw Pipes (Flow conduits)
    # Col 1 down
    draw_conduit_pipe(slide2, Inches(1.725), Inches(2.8), Inches(1.725), Inches(3.0), COLOR_APPLY)
    draw_conduit_pipe(slide2, Inches(1.725), Inches(3.6), Inches(1.725), Inches(3.8), COLOR_APPLY)
    draw_conduit_pipe(slide2, Inches(1.725), Inches(4.4), Inches(1.725), Inches(5.95), COLOR_APPLY) # table tracing line
    
    # Col 1 -> Col 2
    draw_conduit_pipe(slide2, Inches(2.625), Inches(4.1), Inches(3.225), Inches(4.1), COLOR_APPLY)
    
    # Col 2 -> Col 3 (Reclaimable path)
    draw_conduit_pipe(slide2, Inches(5.025), Inches(4.1), Inches(5.625), Inches(4.1), COLOR_PROD, label="🟢可再製")
    
    # Col 2 -> Col 4 (Non-reclaimable path)
    draw_conduit_pipe(slide2, Inches(4.125), Inches(4.4), Inches(4.125), Inches(4.8), COLOR_SALES)
    draw_conduit_pipe(slide2, Inches(4.125), Inches(4.8), Inches(8.025), Inches(4.8), COLOR_SALES, label="🔵不可再製")
    
    # Col 3 down
    draw_conduit_pipe(slide2, Inches(6.525), Inches(4.4), Inches(6.525), Inches(5.0), COLOR_PROD)
    # Col 3 INV-05 path to close
    draw_conduit_pipe(slide2, Inches(6.525), Inches(5.65), Inches(6.525), Inches(5.9), COLOR_PROD)
    draw_conduit_pipe(slide2, Inches(6.525), Inches(5.9), Inches(1.725), Inches(5.9), COLOR_PROD)
    
    # Col 3 Alert Backtrack path (Dotted conduit from Col 3 to Col 4 Sales)
    draw_conduit_pipe(slide2, Inches(7.425), Inches(4.1), Inches(7.725), Inches(4.1), COLOR_SALES, label="⚠️再製失敗")
    draw_conduit_pipe(slide2, Inches(7.725), Inches(4.1), Inches(7.725), Inches(4.8), COLOR_SALES)
    draw_conduit_pipe(slide2, Inches(7.725), Inches(4.8), Inches(8.025), Inches(4.8), COLOR_SALES)
    
    # Col 4 down
    draw_conduit_pipe(slide2, Inches(8.925), Inches(5.1), Inches(8.925), Inches(5.4), COLOR_SALES)
    # Col 4 SOP to close
    draw_conduit_pipe(slide2, Inches(8.925), Inches(6.05), Inches(8.925), Inches(6.2), COLOR_SALES)
    draw_conduit_pipe(slide2, Inches(8.925), Inches(6.2), Inches(1.725), Inches(6.2), COLOR_SALES)
    
    # Col 4 -> Col 5 (Waste path)
    draw_conduit_pipe(slide2, Inches(9.825), Inches(4.8), Inches(10.425), Inches(4.8), COLOR_WASTE, label="🔴無殘值")
    
    # Col 5 down
    draw_conduit_pipe(slide2, Inches(11.325), Inches(5.1), Inches(11.325), Inches(5.4), COLOR_WASTE)
    # Col 5 SOP to close
    draw_conduit_pipe(slide2, Inches(11.325), Inches(6.05), Inches(11.325), Inches(6.35), COLOR_WASTE)
    draw_conduit_pipe(slide2, Inches(11.325), Inches(6.35), Inches(1.725), Inches(6.35), COLOR_WASTE)
    
    # Close -> End
    draw_conduit_pipe(slide2, Inches(1.725), Inches(6.55), Inches(1.725), Inches(6.7), COLOR_APPLY)
    
    # 3. Draw Node Cards
    card_w = Inches(1.8)
    card_h = Inches(0.6)
    c1_x = Inches(0.6 + 0.225)
    c2_x = Inches(3.0 + 0.225)
    c3_x = Inches(5.4 + 0.225)
    c4_x = Inches(7.8 + 0.225)
    c5_x = Inches(10.2 + 0.225)
    
    # Column 1 Cards (申請與監控)
    draw_sop_card(slide2, c1_x, Inches(2.2), card_w, card_h, "無法自行回收料品", "", "", COLOR_APPLY)
    draw_sop_card(slide2, c1_x, Inches(3.0), card_w, card_h, "填寫處理單", "表HS-QA-F-082", "", COLOR_APPLY, is_doc=True)
    draw_sop_card(slide2, c1_x, Inches(3.8), card_w, card_h, "登錄管制表列管", "表HS-QA-T-015", "廠務專責列管", COLOR_APPLY)
    draw_sop_card(slide2, c1_x, Inches(5.95), card_w, card_h, "列管至處理完成", "處理單/管制表銷案", "", COLOR_APPLY)
    
    # Column 2 Card (品質鑑定)
    draw_sop_card(slide2, c2_x, Inches(3.8), card_w, card_h, "品保鑑定品質", "判定是否可回收", "品質鑑定準則", COLOR_QA, is_diamond=True)
    
    # Column 3 Cards (再製執行)
    draw_sop_card(slide2, c3_x, Inches(3.8), card_w, card_h, "生產單位再製處理", "回流精餾 / 大槽調配", "現場排程操作", COLOR_PROD)
    draw_sop_card(slide2, c3_x, Inches(5.0), card_w, Inches(0.65), "格外品處理、儲存辦法", "C50110-INV-05", "再製回收及儲存規範", COLOR_PROD, is_doc=True)
    
    # Column 4 Cards (價值評估)
    draw_sop_card(slide2, c4_x, Inches(4.5), card_w, card_h, "業務部評估出售", "評估理化殘值", "折價销售準則", COLOR_SALES, is_diamond=True)
    draw_sop_card(slide2, c4_x, Inches(5.4), card_w, Inches(0.65), "安排降級出貨", "C20110-SAL-02", "次級品合約及地磅過重", COLOR_SALES, is_doc=True)
    
    # Column 5 Cards (核准與報廢)
    draw_sop_card(slide2, c5_x, Inches(4.5), card_w, card_h, "(副)總經理審核", "核准報廢資產", "授權授信矩陣(LOA)", COLOR_WASTE, is_diamond=True)
    draw_sop_card(slide2, c5_x, Inches(5.4), card_w, Inches(0.65), "事業廢棄物處置", "C10500-EPO-02", "委託甲級清除商/GPS軌跡", COLOR_WASTE, is_doc=True)


    # ==========================================
    # SLIDE 3: Responsibilities & Form Lifecycle
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "02. 權責分工與表單生命週期", "明確判定與執行角色，確保格外品全程有跡可循")
    
    # Left Box: Roles & Responsibilities (R&R)
    rr_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2))
    rr_box.fill.solid()
    rr_box.fill.fore_color.rgb = CARD_BG
    rr_box.line.color.rgb = COLOR_APPLY
    rr_box.line.width = Pt(1.5)
    
    tf_rr = rr_box.text_frame
    tf_rr.word_wrap = True
    tf_rr.margin_left = Inches(0.2)
    tf_rr.margin_right = Inches(0.2)
    tf_rr.margin_top = Inches(0.2)
    
    p = tf_rr.paragraphs[0]
    p.text = "各部門核心權責分工 (R&R)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_APPLY
    p.space_after = Pt(12)
    
    roles = [
        ("申請部門", "負責在現場隔離異常品、掛牌標記並於24小時內填寫格外品處理單。", COLOR_APPLY),
        ("品保部", "負責重新採樣化驗，並依理化指標做「🟢可回收」或「🔴不可回收」的關鍵分流鑑定。", COLOR_QA),
        ("生產課室", "負責將核准之可回收品進行再製（精餾回流/大槽按比例調混），若失敗則通報業務。", COLOR_PROD),
        ("業務部", "負責不可回收產品的商業價值殘值判定，尋找次級品買家折價出售。", COLOR_SALES),
        ("總經理/環安", "核准無殘值產品的資產報廢，環安嚴格遵照環保法規委外清運處置。", COLOR_WASTE)
    ]
    
    for title, desc, color in roles:
        p_role = tf_rr.add_paragraph()
        p_role.text = f"■ {title}："
        p_role.font.name = FONT_FAMILY
        p_role.font.size = Pt(11)
        p_role.font.bold = True
        p_role.font.color.rgb = color
        p_role.space_before = Pt(8)
        
        p_desc = tf_rr.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_after = Pt(4)
        p_desc.margin_left = Inches(0.2)

    # Right Box: Document/Form Lifecycle
    doc_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2))
    doc_box.fill.solid()
    doc_box.fill.fore_color.rgb = CARD_BG
    doc_box.line.color.rgb = COLOR_QA
    doc_box.line.width = Pt(1.5)
    
    tf_doc = doc_box.text_frame
    tf_doc.word_wrap = True
    tf_doc.margin_left = Inches(0.2)
    tf_doc.margin_right = Inches(0.2)
    tf_doc.margin_top = Inches(0.2)
    
    p = tf_doc.paragraphs[0]
    p.text = "雙表管理生命週期 (Form Lifecycle)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_QA
    p.space_after = Pt(12)
    
    lifecycles = [
        ("【表1】格外品、久滯品處理單 (流轉專用表單)", "● 誕生：現場產出異常料品時，由申請單位發起，檢附COA原始化驗單。\n● 流轉：品保填寫品質鑑定結論 ➔ 現場填寫再製紀錄 ➔ 或業務填寫降級合約 ➔ 或總經理報廢簽核。\n● 銷案：於最後處置完成並上傳憑證後，廠務管制窗口確認結案銷案。"),
        ("【表2】格外品、久滯品管制表 (追蹤列管台帳)", "● 誕生：品保部收到處理單後立案，發放唯一系統列管序號（HS-OP-XXXX）。\n● 追蹤：列為廠務部「每週生產調度會議」重點追蹤台帳，無結案憑證者不得擅自銷案。\n● 存查：結案後，憑證與處理單一同歸檔，系統保存5年以上以備稽核調閱。")
    ]
    
    for title, text in lifecycles:
        p_title = tf_doc.add_paragraph()
        p_title.text = title
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(11.5)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.space_before = Pt(8)
        p_title.space_after = Pt(4)
        
        p_body = tf_doc.add_paragraph()
        p_body.text = text
        p_body.font.name = FONT_FAMILY
        p_body.font.size = Pt(9.5)
        p_body.font.color.rgb = TEXT_MUTED
        p_body.space_after = Pt(10)


    # ==========================================
    # SLIDE 4: Core SOP & Audit Checks
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "03. 核心管理辦法與稽核防錯機制", "三大 SOP 互鎖，確保財務帳實一致性與環保合規性")
    
    # 3 SOP Cards
    sop_w = Inches(3.7)
    sop_h = Inches(5.2)
    sop_top = Inches(1.5)
    
    # SOP 1: INV-05
    sop1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), sop_top, sop_w, sop_h)
    sop1.fill.solid()
    sop1.fill.fore_color.rgb = CARD_BG
    sop1.line.color.rgb = COLOR_PROD
    sop1.line.width = Pt(1.5)
    tf1 = sop1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.15)
    tf1.margin_right = Inches(0.15)
    tf1.margin_top = Inches(0.15)
    
    p = tf1.paragraphs[0]
    p.text = "格外品處理、儲存管理辦法"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PROD
    p.space_after = Pt(2)
    p_code1 = tf1.add_paragraph()
    p_code1.text = "辦法編號：C50110-INV-05 (3.0版)"
    p_code1.font.name = FONT_FAMILY
    p_code1.font.size = Pt(8.5)
    p_code1.font.bold = True
    p_code1.font.color.rgb = TEXT_MUTED
    p_code1.space_after = Pt(12)
    
    inv_bullets = [
        "● 久滯品期限判定：依最後異動起算，電子級儲槽超60天、工業級儲槽超120天即列入久滯品台帳追蹤去化。",
        "● 儲存容器規範：1M3塑膠新/回收桶需在2年內製造；格外品存放超半年的料品，必須以ISO-TANK或白鐵桶存放。",
        "● 外觀與BPM控制：桶裝存放區須黏貼A4黃色「格外品專用外觀標示」註明啟用日期與年限，並於BPM開立處理單。"
    ]
    for b in inv_bullets:
        p_b = tf1.add_paragraph()
        p_b.text = b
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_after = Pt(8)
        
    # SOP 2: SAL-02
    sop2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), sop_top, sop_w, sop_h)
    sop2.fill.solid()
    sop2.fill.fore_color.rgb = CARD_BG
    sop2.line.color.rgb = COLOR_SALES
    sop2.line.width = Pt(1.5)
    tf2 = sop2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.15)
    tf2.margin_right = Inches(0.15)
    tf2.margin_top = Inches(0.15)
    
    p = tf2.paragraphs[0]
    p.text = "出貨作業辦法"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_SALES
    p.space_after = Pt(2)
    p_code2 = tf2.add_paragraph()
    p_code2.text = "辦法編號：C20110-SAL-02"
    p_code2.font.name = FONT_FAMILY
    p_code2.font.size = Pt(8.5)
    p_code2.font.bold = True
    p_code2.font.color.rgb = TEXT_MUTED
    p_code2.space_after = Pt(12)
    
    sal_bullets = [
        "● 次級品標籤化：銷售出貨單、合約與發票上必須明文備註「格外品/降級品，概不退換」，出廠檢驗單(COA)必須以實測異常值開具。",
        "● 二次地磅檢驗：槽車裝載出廠前，必須經過工廠大門過重地磅，確認出廠淨重與申報無誤後方可放行。",
        "● ⚠️稽核紅線：嚴禁使用合格品包裝或COA來混淆降級品出貨！一旦查獲，涉案同仁將直接予以開除處分並送法辦。"
    ]
    for b in sal_bullets:
        p_b = tf2.add_paragraph()
        p_b.text = b
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_after = Pt(8)

    # SOP 3: EPO-02
    sop3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), sop_top, sop_w, sop_h)
    sop3.fill.solid()
    sop3.fill.fore_color.rgb = CARD_BG
    sop3.line.color.rgb = COLOR_WASTE
    sop3.line.width = Pt(1.5)
    tf3 = sop3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = Inches(0.15)
    tf3.margin_right = Inches(0.15)
    tf3.margin_top = Inches(0.15)
    
    p = tf3.paragraphs[0]
    p.text = "廢棄物管理辦法"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WASTE
    p.space_after = Pt(2)
    p_code3 = tf3.add_paragraph()
    p_code3.text = "辦法編號：C10500-EPO-02"
    p_code3.font.name = FONT_FAMILY
    p_code3.font.size = Pt(8.5)
    p_code3.font.bold = True
    p_code3.font.color.rgb = TEXT_MUTED
    p_code3.space_after = Pt(12)
    
    epo_bullets = [
        "● 甲級清除資質互鎖：廢棄物處置必須委託經環保署核可、具資質之甲級事業廢棄物清除商，簽署正式清運合約。",
        "● 三聯單與GPS追蹤：清運時線上填報「三聯單」並封存備查，追蹤清運槽車之GPS軌跡，確保清運至合規焚化爐。",
        "● ⚠️稽核紅線：嚴禁將化學廢棄物交予地下無證回收商或非法傾倒！違者除面臨數百萬罰鍰外，環安主管將負刑事責任。"
    ]
    for b in epo_bullets:
        p_b = tf3.add_paragraph()
        p_b.text = b
        p_b.font.name = FONT_FAMILY
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_after = Pt(8)

    # Save Presentation
    output_path = os.path.join(r"g:\我的雲端硬碟\GOOGLE ANGET", "20260602格外品久滯品流程圖_優化版.pptx")
    prs.save(output_path)
    print(f"PowerPoint Presentation successfully created and saved to {output_path}!")

if __name__ == "__main__":
    create_optimized_sop_ppt()
