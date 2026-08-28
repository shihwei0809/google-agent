import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_edu_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ------------------ DESIGN SYSTEM COLORS ------------------
    BG_COLOR = RGBColor(9, 13, 22)
    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    COLOR_PRIMARY = RGBColor(56, 189, 248) # Sky Blue
    
    BG_LIGHT = RGBColor(248, 250, 252)
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_DARK_MUTED = RGBColor(71, 85, 105)
    CARD_LIGHT = RGBColor(241, 245, 249)
    CARD_DARK = RGBColor(15, 23, 42)
    
    FONT_FAMILY = "Microsoft JhengHei"
    
    def set_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, x, y, w, h, text, font_size, color, bold=False):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        return txBox

    # SLIDE 1: Title Slide (Dark)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide1, BG_COLOR)
    
    # Left accent bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    
    add_textbox(slide1, Inches(2), Inches(2.2), Inches(8), Inches(1), "教育訓練作業管理辦法", 48, TEXT_WHITE, True)
    add_textbox(slide1, Inches(2), Inches(3.2), Inches(8), Inches(0.8), "4.0版重點宣導", 32, COLOR_PRIMARY, True)
    
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.2), Inches(6), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = TEXT_MUTED
    line.line.fill.background()
    
    add_textbox(slide1, Inches(2), Inches(4.5), Inches(8), Inches(0.5), "把訓練做成可追蹤、可驗證、可持續改善的管理流程", 20, TEXT_WHITE, False)
    
    add_textbox(slide1, Inches(2), Inches(6.5), Inches(8), Inches(0.5), "適用：全公司從業人員之職前、在職教育訓練及相關作業", 14, TEXT_MUTED, False)

    # SLIDE 2: Purpose (Light)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide2, BG_LIGHT)
    
    add_textbox(slide2, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4), "01 | 先理解目的", 14, COLOR_PRIMARY, True)
    add_textbox(slide2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8), "這份辦法要解決的，是「能力有沒有真的到位」", 36, TEXT_DARK, True)
    add_textbox(slide2, Inches(0.6), Inches(1.6), Inches(10), Inches(0.4), "訓練不只是在上課，而是從需求、執行到資格與成效的完整閉環。", 16, TEXT_DARK_MUTED, False)
    
    # Left Box
    left_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(5.5), Inches(4))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = CARD_LIGHT
    left_box.line.fill.background()
    
    tf = left_box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.4)
    p = tf.paragraphs[0]
    p.text = "管理目的"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "讓員工了解公司理念、目標與作業\n流程，具備本職所需的專業知識與\n技能。"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_DARK
    p2.font.bold = True
    p2.space_before = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "最終目標"
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_DARK_MUTED
    p3.space_before = Pt(40)
    
    p4 = tf.add_paragraph()
    p4.text = "品質 | 效率 | 環境 | 安全衛生"
    p4.font.name = FONT_FAMILY
    p4.font.size = Pt(20)
    p4.font.color.rgb = COLOR_PRIMARY
    
    # Right items
    rights = [
        ("誰適用？", "全公司所屬從業人員"),
        ("涵蓋什麼？", "職前、在職訓練與相關作業"),
        ("如何落地？", "內訓、外訓、證照、資格驗證與系統留存")
    ]
    for i, (rt, rb) in enumerate(rights):
        y = 2.8 + (i * 1.2)
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(y), Inches(0.1), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY
        bar.line.fill.background()
        add_textbox(slide2, Inches(7.1), Inches(y-0.1), Inches(5), Inches(0.4), rt, 18, TEXT_DARK, True)
        add_textbox(slide2, Inches(7.1), Inches(y+0.25), Inches(5), Inches(0.4), rb, 14, TEXT_DARK_MUTED, False)

    # SLIDE 3: Flowchart (Light)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide3, BG_LIGHT)
    
    add_textbox(slide3, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4), "02 | 管理主線", 14, COLOR_PRIMARY, True)
    add_textbox(slide3, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8), "每一堂訓練，都要走完這條管理鏈", 36, TEXT_DARK, True)
    add_textbox(slide3, Inches(0.6), Inches(1.6), Inches(10), Inches(0.4), "從職能差距出發，最後回到成效與稽核。", 16, TEXT_DARK_MUTED, False)
    
    steps = [
        ("01", "找需求", "策略 / KPI / 職能差距", CARD_LIGHT, TEXT_DARK),
        ("02", "排計畫", "需求調查→年度計畫", CARD_LIGHT, TEXT_DARK),
        ("03", "做訓練", "內訓、外訓、職前/變更工作前", CARD_LIGHT, TEXT_DARK),
        ("04", "驗資格", "測驗、實作、證照與名冊", CARD_LIGHT, TEXT_DARK),
        ("05", "留證據", "ERP履歷、簽到、心得、成效", CARD_DARK, TEXT_WHITE)
    ]
    
    start_x = 0.6
    w = 2.2
    for i, (num, title, desc, bg_c, text_c) in enumerate(steps):
        x = start_x + (i * 2.5)
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(w), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_c
        box.line.fill.background()
        
        num_c = COLOR_PRIMARY if bg_c == CARD_LIGHT else COLOR_PRIMARY
        add_textbox(slide3, Inches(x+0.2), Inches(2.8), Inches(1.8), Inches(0.4), num, 24, num_c, True)
        add_textbox(slide3, Inches(x+0.2), Inches(3.4), Inches(1.8), Inches(0.4), title, 20, text_c, True)
        
        desc_box = add_textbox(slide3, Inches(x+0.2), Inches(4.0), Inches(1.8), Inches(0.8), desc, 12, TEXT_DARK_MUTED if bg_c == CARD_LIGHT else TEXT_WHITE, False)
        desc_box.text_frame.word_wrap = True
        
        if i < 4:
            arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+w+0.05), Inches(3.6), Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_PRIMARY
            arrow.line.fill.background()
            
    add_textbox(slide3, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5), "判斷標準：訓練完成 ≠ 上完課；必須能證明「人員受訓、能力驗證、紀錄可查」。", 18, COLOR_PRIMARY, True)

    output_path = r"D:\GOOGLE ANGET\教育訓練-C0588-教材\C10200-EDU-01 教育訓練作業管理辦法(Codex樣式).pptx"
    prs.save(output_path)

create_edu_ppt()

