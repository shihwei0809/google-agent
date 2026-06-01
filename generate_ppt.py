import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create the presentation object
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Color definitions
BG_COLOR = RGBColor(11, 14, 20)      # #0B0E14 (Premium Dark)
TEXT_WHITE = RGBColor(255, 255, 255) # #FFFFFF
TEXT_GRAY = RGBColor(187, 187, 187)  # #BBBBBB (Subtext)
COLOR_PINK = RGBColor(255, 42, 122)   # #FF2A7A (Modern Neon Pink)
COLOR_GOLD = RGBColor(212, 175, 55)   # #D4AF37 (Traditional Gold)

# Image paths (Absolute paths in the brain directory)
brain_dir = r"C:\Users\C606\.gemini\antigravity\brain\62d1b10e-92db-4d8f-b075-b9c75d607001"
images = {
    "cover": os.path.join(brain_dir, "osaka_cover_1780231751100.png"),
    "map": os.path.join(brain_dir, "osaka_map_1780231767135.png"),
    "day1": os.path.join(brain_dir, "dotonbori_night_1780231782470.png"),
    "day2": os.path.join(brain_dir, "osaka_castle_1780231802361.png"),
    "day3": os.path.join(brain_dir, "usj_fantasy_1780231818761.png"),
    "day4": os.path.join(brain_dir, "kyoto_torii_1780231835224.png"),
    "day5": os.path.join(brain_dir, "kuromon_food_1780231856935.png")
}

# Create a local images folder in the workspace and copy images there for user access
workspace_dir = r"G:\我的雲端硬碟\GOOGLE ANGET"
workspace_img_dir = os.path.join(workspace_dir, "images")
os.makedirs(workspace_img_dir, exist_ok=True)

local_images = {}
for key, filepath in images.items():
    if os.path.exists(filepath):
        dest = os.path.join(workspace_img_dir, f"{key}.png")
        shutil.copy(filepath, dest)
        local_images[key] = dest
        print(f"Copied {key} image to {dest}")
    else:
        print(f"Warning: Image {filepath} not found.")
        local_images[key] = filepath # fallback to absolute path if copying failed

def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_clean_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    return tf

# ----------------- SLIDE 1: COVER -----------------
slide_layout = prs.slide_layouts[6] # Blank slide
slide1 = prs.slides.add_slide(slide_layout)
set_dark_background(slide1)

# Add Cover Image on the right (60% width)
if os.path.exists(local_images["cover"]):
    slide1.shapes.add_picture(local_images["cover"], Inches(5.5), Inches(0), Inches(7.833), Inches(7.5))

# Text Frame on the left panel
tf1 = add_clean_textbox(slide1, Inches(0.8), Inches(1.8), Inches(4.2), Inches(4.5))

# Title
p = tf1.paragraphs[0]
p.text = "大阪五天四夜"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(10)

# Subtitle (Accent color)
p2 = tf1.add_paragraph()
p2.text = "璀璨霓虹與經典古蹟"
p2.font.name = 'Microsoft JhengHei'
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = COLOR_PINK
p2.space_after = Pt(30)

p3 = tf1.add_paragraph()
p3.text = "關西自由行規劃指南"
p3.font.name = 'Microsoft JhengHei'
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT_GRAY
p3.space_after = Pt(50)

# Footer info
p4 = tf1.add_paragraph()
p4.text = "✈️ 深度關西旅程\n🗺️ 現代與歷史的完美交織"
p4.font.name = 'Microsoft JhengHei'
p4.font.size = Pt(12)
p4.font.color.rgb = TEXT_GRAY


# ----------------- SLIDE 2: OVERVIEW -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_dark_background(slide2)

# Left Map Image
if os.path.exists(local_images["map"]):
    slide2.shapes.add_picture(local_images["map"], Inches(0.8), Inches(1.0), Inches(4.8), Inches(5.5))

# Right Content Box
tf2 = add_clean_textbox(slide2, Inches(6.2), Inches(0.8), Inches(6.333), Inches(6.0))

# Title
p = tf2.paragraphs[0]
p.text = "五天四夜 玩轉關西"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

# Subtitle
p_sub = tf2.add_paragraph()
p_sub.text = "行程概覽與地理動線"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_GOLD
p_sub.space_after = Pt(25)

# Day by day list
days_info = [
    ("Day 1", "降落關西 ➔ 道頓堀璀璨霓虹與美食巡禮", "體驗大阪活力夜生活，打卡固力果跑跑人"),
    ("Day 2", "穿越古今 ➔ 大阪城天守閣與梅田空中花園", "欣賞宏偉天守閣與 360 度無死角百萬都市夜景"),
    ("Day 3", "奇幻冒險 ➔ 日本環球影城 USJ 一整天狂歡", "沉浸於瑪利歐賽車與哈利波特魔法世界"),
    ("Day 4", "古都慢活 ➔ 京都伏見稻荷與清水寺一日遊", "漫步千本鳥居與清水舞台的朱紅歷史記憶"),
    ("Day 5", "最後衝刺 ➔ 黑門市場海鮮與心齋橋採購 ➔ 返台", "享受烤和牛與巨無霸干貝，臨空城 Outlets 血拼")
]

for day, desc, details in days_info:
    p_day = tf2.add_paragraph()
    p_day.text = f"📍 {day} ： {desc}"
    p_day.font.name = 'Microsoft JhengHei'
    p_day.font.size = Pt(14)
    p_day.font.bold = True
    p_day.font.color.rgb = COLOR_PINK
    
    p_det = tf2.add_paragraph()
    p_det.text = f"     {details}"
    p_det.font.name = 'Microsoft JhengHei'
    p_det.font.size = Pt(11)
    p_det.font.color.rgb = TEXT_GRAY
    p_det.space_after = Pt(12)


# ----------------- SLIDE 3: DAY 1 -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_dark_background(slide3)

# Left Content Box
tf3 = add_clean_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.8), Inches(6.0))

p = tf3.paragraphs[0]
p.text = "Day 1 | 降落關西，擁抱道頓堀"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

p_sub = tf3.add_paragraph()
p_sub.text = "感受大阪最奔放的熱情與美食之夜"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_PINK
p_sub.space_after = Pt(30)

items = [
    ("⚡ 交通直達：南海電鐵 Rapi:t", "從關西機場搭乘特急 Rapi:t 藍色鋼鐵車身列車，只需 38 分鐘即可直達難波站，舒適又快速。"),
    ("🛍️ 潮流前線：心齋橋筋商店街", "日本最具代表性的拱廊商店街！聚集各大家藥妝店、日系潮流品牌與國際精品，滿足一切購物慾望。"),
    ("🐙 美食天堂：道頓堀川與巨型招牌", "在固力果跑跑人霓虹招牌前打卡！沿河川品嚐本家章魚燒、金龍拉麵、美津的大阪燒，感受關西庶民美食實力。")
]

for title, desc in items:
    pt = tf3.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GOLD
    
    pd = tf3.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(11)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# Right Image
if os.path.exists(local_images["day1"]):
    slide3.shapes.add_picture(local_images["day1"], Inches(7.0), Inches(1.0), Inches(5.533), Inches(5.5))


# ----------------- SLIDE 4: DAY 2 -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_dark_background(slide4)

# Left Image
if os.path.exists(local_images["day2"]):
    slide4.shapes.add_picture(local_images["day2"], Inches(0.8), Inches(1.0), Inches(5.533), Inches(5.5))

# Right Content Box
tf4 = add_clean_textbox(slide4, Inches(7.0), Inches(0.8), Inches(5.533), Inches(6.0))

p = tf4.paragraphs[0]
p.text = "Day 2 | 歷史與現代的交織"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

p_sub = tf4.add_paragraph()
p_sub.text = "從豐臣秀吉的城堡到摩天大樓展望台"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_GOLD
p_sub.space_after = Pt(30)

items = [
    ("🏯 歷史巡禮：大阪城天守閣", "參觀白色城牆配上亮眼金箔的雄偉天守閣，登上 8 樓展望台俯瞰大阪城公園的美麗綠意與護城河景致。"),
    ("🍜 庶民漫步：天神橋筋商店街", "全長 2.6 公里為日本最長商店街！充滿懷舊氣息的在地小吃、便宜藥妝與生活雜貨，感受最道地的大阪生活。"),
    ("🌃 百萬夜景：梅田藍天大廈", "登上 173 公尺高的空中庭園展望台，享受 360 度無死角的室外露天微風，將大阪繁華的都市霓虹盡收眼底。")
]

for title, desc in items:
    pt = tf4.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_PINK
    
    pd = tf4.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(11)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)


# ----------------- SLIDE 5: DAY 3 -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_dark_background(slide5)

# Left Content Box
tf5 = add_clean_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.8), Inches(6.0))

p = tf5.paragraphs[0]
p.text = "Day 3 | 環球影城的奇幻冒險"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

p_sub = tf5.add_paragraph()
p_sub.text = "進入電影與遊戲實體化的夢幻國度"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_PINK
p_sub.space_after = Pt(30)

items = [
    ("🍄 瑪利歐：超級任天堂世界", "利用能量手環在園區內敲擊磚頭收集金幣！必玩「瑪利歐賽車～庫巴的挑戰書～」，並到奇諾比奧咖啡店享受美味料理。"),
    ("🧙 魔法世界：哈利波特的禁忌之旅", "走入壯麗的霍格華茲城堡，手持互動魔杖施展魔法。在活米村喝杯冰涼奶油啤酒，體驗極致魔法魅力。"),
    ("💡 樂園攻略：提早入園與快速通關", "務必下載 USJ 官方 App 搶整理券，建議提早 1-1.5 小時至大門口排隊，或提前購買快速通關 (Express Pass) 省去排隊時間。")
]

for title, desc in items:
    pt = tf5.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GOLD
    
    pd = tf5.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(11)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# Right Image
if os.path.exists(local_images["day3"]):
    slide5.shapes.add_picture(local_images["day3"], Inches(7.0), Inches(1.0), Inches(5.533), Inches(5.5))


# ----------------- SLIDE 6: DAY 4 -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_dark_background(slide6)

# Left Image
if os.path.exists(local_images["day4"]):
    slide6.shapes.add_picture(local_images["day4"], Inches(0.8), Inches(1.0), Inches(5.533), Inches(5.5))

# Right Content Box
tf6 = add_clean_textbox(slide6, Inches(7.0), Inches(0.8), Inches(5.533), Inches(6.0))

p = tf6.paragraphs[0]
p.text = "Day 4 | 京都古都一日巡禮"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

p_sub = tf6.add_paragraph()
p_sub.text = "漫步千本鳥居與朱紅舞台的關西記憶"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_GOLD
p_sub.space_after = Pt(30)

items = [
    ("⛩️ 視覺震撼：伏見稻荷大社", "朝聖綿延整座稻荷山的「千本鳥居」！極具視覺張力的鮮紅鳥居隧道，是京都最經典的明信片級畫面。"),
    ("🍂 歷史懸空：清水寺與二年坂", "站在巨大的木造「清水舞台」上眺望京都街景。下山沿著二年坂、三年坂石板路漫步，體驗傳統抹茶與日式紙傘風情。"),
    ("🚆 快捷交通：京阪電車一日券", "使用京阪電車一日乘車券，從大阪難波/梅田輕鬆直達京都，沿途暢遊宇治、伏見稻荷與祇園，便利又划算。")
]

for title, desc in items:
    pt = tf6.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_PINK
    
    pd = tf6.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(11)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)


# ----------------- SLIDE 7: DAY 5 -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_dark_background(slide7)

# Left Content Box
tf7 = add_clean_textbox(slide7, Inches(0.8), Inches(0.8), Inches(5.8), Inches(6.0))

p = tf7.paragraphs[0]
p.text = "Day 5 | 美食廚房與最後血拼"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.space_after = Pt(5)

p_sub = tf7.add_paragraph()
p_sub.text = "帶著滿滿回憶與超值戰利品凱旋歸國"
p_sub.font.name = 'Microsoft JhengHei'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = COLOR_PINK
p_sub.space_after = Pt(30)

items = [
    ("🦞 胃袋覺醒：黑門市場海鮮", "大阪的庶民廚房！現場品嚐現烤神戶牛、巨無霸大干貝、現開生蠔、海膽與極致新鮮的黑鮪魚生魚片，大飽口福。"),
    ("🛍️ 機場前哨：臨空城 Rinku Outlets", "搭乘接駁車僅需 20 分鐘即可抵達關西機場。在此瘋狂採購折價國際名牌、運動服飾，為行李箱做最後衝刺。"),
    ("✈️ 機場補貨與登機小提醒", "提早 2.5 至 3 小時抵達關西機場辦理托運。入關後在免稅店採購薯條三兄弟、白色戀人等超人氣伴手禮，完美收尾！")
]

for title, desc in items:
    pt = tf7.add_paragraph()
    pt.text = title
    pt.font.name = 'Microsoft JhengHei'
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GOLD
    
    pd = tf7.add_paragraph()
    pd.text = desc
    pd.font.name = 'Microsoft JhengHei'
    pd.font.size = Pt(11)
    pd.font.color.rgb = TEXT_GRAY
    pd.space_after = Pt(20)

# Right Image
if os.path.exists(local_images["day5"]):
    slide7.shapes.add_picture(local_images["day5"], Inches(7.0), Inches(1.0), Inches(5.533), Inches(5.5))

# Save the presentation
output_path = os.path.join(workspace_dir, "大阪五天四夜自由行.pptx")
prs.save(output_path)
print(f"Presentation successfully created and saved to {output_path}!")
