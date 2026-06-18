#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
簡報生成工具：讀取大綱文字，利用 python-pptx 生成 16:9 孟菲斯波普風格投影片。
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

OUTLINE_PATH = "input/slide_outline.txt"
OUTPUT_PATH = "output/memphis_presentation.pptx"

# 孟菲斯波普調色盤 (RGB)
BG_COLOR = RGBColor(0xFF, 0xFD, 0xF0)       # 淺奶油色 #FFFDF0
TEXT_COLOR = RGBColor(0x1A, 0x1A, 0x1A)     # 炭黑色 #1A1A1A
ACCENT_YELLOW = RGBColor(0xFF, 0xE0, 0x00)  # 電力黃 #FFE000
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)     # 珊瑚紅 #FF4D4D
ACCENT_BLUE = RGBColor(0x2E, 0x5B, 0xFF)    # 電光藍 #2E5BFF

# 裝飾用顏色池
ACCENT_COLORS = [ACCENT_YELLOW, ACCENT_RED, ACCENT_BLUE]

def load_outline():
    if not os.path.exists(OUTLINE_PATH):
        # 預設簡報大綱
        return [
            {"title": "鴻勝化學：智慧物流與品質控制", "content": ["1. 槽車物流實時調度系統", "2. TSMC N-series 高精準條碼驗證", "3. 現場操作離線優先防呆方案"]},
            {"title": "ISOTANK 物流痛點與解決方案", "content": ["1. LLM 運算評估時間易出錯，存在數學幻覺", "2. 解決方案：導入 Python 物流計算核心", "3. 司機與調度端即時數據同步"]},
            {"title": "未來展望：React & Supabase 升級", "content": ["1. 現代化 Web 前端架構重構", "2. Supabase RLS 資料庫層級安全防護", "3. 全面提升系統擴充性與響應速度"]}
        ]
        
    slides_data = []
    current_slide = None
    
    with open(OUTLINE_PATH, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if line.startswith("TITLE:"):
                if current_slide:
                    slides_data.append(current_slide)
                current_slide = {"title": line[6:].strip(), "content": []}
            elif line.startswith("-") and current_slide:
                current_slide["content"].append(line[1:].strip())
                
    if current_slide:
        slides_data.append(current_slide)
        
    return slides_data

def set_slide_background(slide):
    # 用一個滿版矩形設置奶油色背景
    left = top = Inches(0)
    width = Inches(13.333)
    height = Inches(7.5)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background() # 無邊框

def add_memphis_decorations(slide, slide_index):
    accent_color = ACCENT_COLORS[slide_index % len(ACCENT_COLORS)]
    other_color = ACCENT_COLORS[(slide_index + 1) % len(ACCENT_COLORS)]
    
    # 根據不同投影片索引，套用三種交替的孟菲斯風格排版與幾何拼貼
    if slide_index % 3 == 0:
        # === 版面 A：標題背景色塊與右下幾何拼貼 ===
        # 標題背後的實心偏移色塊
        title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.85), Inches(11.833), Inches(1.3))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = accent_color
        title_bg.line.color.rgb = TEXT_COLOR
        title_bg.line.width = Pt(3.5)
        
        # 標題背景黑底陰影
        title_shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(0.92), Inches(11.833), Inches(1.3))
        title_shadow.fill.solid()
        title_shadow.fill.fore_color.rgb = TEXT_COLOR
        title_shadow.line.fill.background()
        
        # 將陰影與背景送至下層
        slide.shapes._spTree.remove(title_shadow._element)
        slide.shapes._spTree.remove(title_bg._element)
        slide.shapes._spTree.insert(1, title_shadow._element)
        slide.shapes._spTree.insert(2, title_bg._element)
        
        # 右下角圓形與三角拼貼
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(5.2), Inches(1.4), Inches(1.4))
        c.fill.solid()
        c.fill.fore_color.rgb = other_color
        c.line.color.rgb = TEXT_COLOR
        c.line.width = Pt(3)
        
        t = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10.2), Inches(5.8), Inches(1.2), Inches(1.2))
        t.fill.solid()
        t.fill.fore_color.rgb = TEXT_COLOR
        t.line.fill.background()
        t.rotation = 180
        
    elif slide_index % 3 == 1:
        # === 版面 B：左側粗邊框邊條與右上重疊圓形 ===
        # 左側裝飾粗邊條
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(0.3), Inches(6.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.color.rgb = TEXT_COLOR
        bar.line.width = Pt(3)
        
        bar_shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.46), Inches(0.46), Inches(0.3), Inches(6.7))
        bar_shadow.fill.solid()
        bar_shadow.fill.fore_color.rgb = TEXT_COLOR
        bar_shadow.line.fill.background()
        
        slide.shapes._spTree.remove(bar_shadow._element)
        slide.shapes._spTree.remove(bar._element)
        slide.shapes._spTree.insert(1, bar_shadow._element)
        slide.shapes._spTree.insert(2, bar._element)
        
        # 右上角重疊圓形
        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(0.6), Inches(1.1), Inches(1.1))
        c1.fill.solid()
        c1.fill.fore_color.rgb = other_color
        c1.line.color.rgb = TEXT_COLOR
        c1.line.width = Pt(2.5)
        
        c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.6), Inches(1.0), Inches(0.9), Inches(0.9))
        c2.fill.solid()
        c2.fill.fore_color.rgb = ACCENT_COLORS[(slide_index + 2) % len(ACCENT_COLORS)]
        c2.line.color.rgb = TEXT_COLOR
        c2.line.width = Pt(2.5)
        
    else:
        # === 版面 C：底部雙層長橫條與右上立體菱形 ===
        # 底部橫向條
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.35))
        banner.fill.solid()
        banner.fill.fore_color.rgb = accent_color
        banner.line.color.rgb = TEXT_COLOR
        banner.line.width = Pt(3)
        
        banner_shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.86), Inches(6.66), Inches(11.7), Inches(0.35))
        banner_shadow.fill.solid()
        banner_shadow.fill.fore_color.rgb = TEXT_COLOR
        banner_shadow.line.fill.background()
        
        slide.shapes._spTree.remove(banner_shadow._element)
        slide.shapes._spTree.remove(banner._element)
        slide.shapes._spTree.insert(1, banner_shadow._element)
        slide.shapes._spTree.insert(2, banner._element)
        
        # 右上角立體菱形
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(11.3), Inches(1.2), Inches(1.2), Inches(1.2))
        d.fill.solid()
        d.fill.fore_color.rgb = other_color
        d.line.color.rgb = TEXT_COLOR
        d.line.width = Pt(3)

def create_memphis_presentation():
    prs = Presentation()
    
    # 設置投影片比例為 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6] # 6號為空白版面
    
    slides_data = load_outline()
    
    for idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide)
        add_memphis_decorations(slide, idx)
        
        # 1. 建立標題文字方塊
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_title = tf.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.name = "Alibaba PuHuiTi 3.0"
        p_title.font.size = Pt(38)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_COLOR
        
        # 微調版面 A 的標題垂直位置（使其完美置中於背景塊中）
        if idx % 3 == 0:
            title_box.top = Inches(1.1)
            p_title.font.size = Pt(34)
            p_title.alignment = PP_ALIGN.CENTER
            
        # 2. 建立內容文字方塊 (使用霞鶩文楷/楷體樣式)
        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(10.5), Inches(3.5))
        
        # 微調不同版面的內文邊距避開裝飾物
        if idx % 3 == 1:
            content_box.left = Inches(1.5) # 避開左側粗邊條
            
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = tf_content.margin_top = tf_content.margin_right = tf_content.margin_bottom = 0
        
        for bullet_idx, item in enumerate(data["content"]):
            p = tf_content.add_paragraph() if bullet_idx > 0 else tf_content.paragraphs[0]
            p.text = "   " + item
            p.font.name = "LXGW WenKai"
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(20) # 段落間距
            
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✓ 成功生成優化版孟菲斯風格投影片：{OUTPUT_PATH}")
    print(f"  -> 總投影片張數：{len(slides_data)} 張")

if __name__ == "__main__":
    create_memphis_presentation()
