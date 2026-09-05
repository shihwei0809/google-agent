import os
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import google.generativeai as genai
from dotenv import load_dotenv
import shutil
import io

load_dotenv()

app = FastAPI(title="AI 訓練平台 API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 掛載 materials 目錄為靜態資料夾，用來提供圖片等檔案
MATERIALS_DIR = "materials"
os.makedirs(MATERIALS_DIR, exist_ok=True)
app.mount("/materials_static", StaticFiles(directory=MATERIALS_DIR), name="materials_static")

# 讀取多個 API Key (用逗號分隔)
GEMINI_API_KEYS_STR = os.getenv("GEMINI_API_KEY", "")
API_KEYS = [k.strip() for k in GEMINI_API_KEYS_STR.split(",") if k.strip()]

if not API_KEYS:
    print("[警告] 未設定任何 GEMINI_API_KEY！")
else:
    # 預設先用第一組 Key 初始化
    genai.configure(api_key=API_KEYS[0])

class ChatRequest(BaseModel):
    message: str
    context: str = ""
    material_name: str = "未知教材"

@app.get("/")
def read_root():
    return {"status": "ok", "message": "AI 訓練平台 API 運作中"}

# --- 教材管理 API ---

@app.get("/materials")
def list_materials():
    files = [f for f in os.listdir(MATERIALS_DIR) if f.endswith('.md') or f.endswith('.txt')]
    return {"materials": files}

@app.get("/materials/{filename}")
def get_material(filename: str):
    filepath = os.path.join(MATERIALS_DIR, filename)
    if not os.path.exists(filepath):
        return {"error": "找不到該教材"}
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()
    return {"content": content}

import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

@app.post("/materials")
async def upload_material(file: UploadFile = File(...)):
    filename = file.filename
    ext = filename.lower().split('.')[-1]
    
    # 讀取檔案內容至記憶體
    content_bytes = await file.read()
    
    # 若是 md 或 txt，直接存檔
    if ext in ['md', 'txt']:
        filepath = os.path.join(MATERIALS_DIR, filename)
        with open(filepath, "wb") as f:
            f.write(content_bytes)
        return {"message": "上傳成功", "filename": filename}
    
    # 若是其他格式，進行文字萃取並轉為 md
    extracted_text = f"# {filename} (系統自動轉換)\n\n"
    new_filename = filename.rsplit('.', 1)[0] + ".md"
    filepath = os.path.join(MATERIALS_DIR, new_filename)
    
    try:
        if ext == 'pdf':
            doc = fitz.open(stream=content_bytes, filetype="pdf")
            for page in doc:
                extracted_text += page.get_text() + "\n\n"
                
        elif ext == 'docx':
            doc_file = io.BytesIO(content_bytes)
            doc = docx.Document(doc_file)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n\n"
                
        elif ext == 'xlsx':
            excel_file = io.BytesIO(content_bytes)
            # 讀取所有 sheet
            xl = pd.ExcelFile(excel_file)
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(xl, sheet_name=sheet_name)
                extracted_text += f"## {sheet_name}\n\n"
                extracted_text += df.to_markdown(index=False) + "\n\n"
                
        elif ext == 'pptx':
            ppt_file = io.BytesIO(content_bytes)
            prs = Presentation(ppt_file)
            for i, slide in enumerate(prs.slides):
                extracted_text += f"## 第 {i+1} 頁\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
                extracted_text += "\n"
        else:
            return {"error": "不支援的檔案格式"}
            
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(extracted_text)
            
        return {"message": "轉換並上傳成功", "filename": new_filename}
        
    except Exception as e:
        print(f"轉換失敗: {e}")
        return {"error": f"檔案解析失敗: {str(e)}"}

@app.delete("/materials/{filename}")
def delete_material(filename: str):
    filepath = os.path.join(MATERIALS_DIR, filename)
    if os.path.exists(filepath):
        os.remove(filepath)
        return {"message": "刪除成功"}
    return {"error": "檔案不存在"}

# --- AI 問答 API ---

# 定義模型優先順序清單 (2026 最新版本)
MODEL_FALLBACKS = [
    'gemini-3.8-flash',
    'gemini-3.7-flash',
    'gemini-3.1-pro'
]

import csv
import datetime

from fastapi.responses import StreamingResponse

@app.post("/chat")
async def chat_with_ai(req: ChatRequest):
    if not API_KEYS:
        return {"response": "系統尚未設定任何 GEMINI_API_KEY，無法提供 AI 服務。"}
    
    # 自動讀取所有教材 (全知模式)
    all_materials_content = ""
    if os.path.exists(MATERIALS_DIR):
        for filename in os.listdir(MATERIALS_DIR):
            if filename.endswith(".md"):
                try:
                    with open(os.path.join(MATERIALS_DIR, filename), "r", encoding="utf-8") as f:
                        all_materials_content += f"\n\n--- 教材: {filename} ---\n" + f.read()
                except:
                    pass

    prompt = f"你是一個專業的企業內訓 AI 助教。請根據以下【所有教材內容】，親切且專業地回答學員的問題。\n如果學員的問題跨越了多份教材，請幫忙統整答案。\n如果問題與教材完全無關，請委婉告知。\n\n【所有教材內容】\n{all_materials_content}\n\n【學員問題】\n{req.message}"
    
    # 雙重備援機制：先輪替 API Keys，再輪替模型
    for key_idx, current_key in enumerate(API_KEYS):
        # 切換當前使用的 API Key
        genai.configure(api_key=current_key)
        
        for model_name in MODEL_FALLBACKS:
            try:
                model = genai.GenerativeModel(model_name)
                # 開啟 stream=True 模式
                response = model.generate_content(prompt, stream=True)
                
                async def generate():
                    full_text = ""
                    # 逐字回傳給前端
                    for chunk in response:
                        if chunk.text:
                            full_text += chunk.text
                            yield chunk.text
                    
                    # 提示當下使用的模型與 Key (僅顯示前幾碼以利辨識)
                    masked_key = f"{current_key[:4]}...{current_key[-4:]}"
                    footer = f"\n\n*(Powered by {model_name} / Key: {masked_key})*"
                    full_text += footer
                    yield footer
                    
                    # --- 在回傳結束後寫入 Excel (以月份分頁) ---
                    try:
                        import openpyxl
                        log_file = "chat_logs.xlsx"
                        now = datetime.datetime.now()
                        month_str = now.strftime("%Y-%m")
                        timestamp = now.strftime("%Y-%m-%d %H:%M:%S")
                        
                        if os.path.exists(log_file):
                            wb = openpyxl.load_workbook(log_file)
                        else:
                            wb = openpyxl.Workbook()
                            # 移除預設的空 Sheet
                            if "Sheet" in wb.sheetnames:
                                del wb["Sheet"]
                                
                        # 若當月的分頁不存在，則建立並寫入標題列
                        if month_str not in wb.sheetnames:
                            ws = wb.create_sheet(title=month_str)
                            ws.append(["時間", "當前檢視教材", "學員提問", "AI回覆"])
                        else:
                            ws = wb[month_str]
                            
                        # 寫入提問紀錄
                        ws.append([timestamp, req.material_name, req.message, full_text])
                        wb.save(log_file)
                    except Exception as log_e:
                        print(f"寫入 Excel 日誌失敗: {log_e}")
                    # -----------------------------
                    
                return StreamingResponse(generate(), media_type="text/plain")
            except Exception as e:
                error_msg = str(e)
                print(f"[警告] Key({key_idx+1}/{len(API_KEYS)}) 模型 {model_name} 呼叫失敗 ({error_msg})，嘗試切換...")
                continue
                
    return {"response": "系統內建的所有 API Key 以及 Gemini 模型備援方案皆已用盡或發生異常，請聯絡系統管理員！"}
