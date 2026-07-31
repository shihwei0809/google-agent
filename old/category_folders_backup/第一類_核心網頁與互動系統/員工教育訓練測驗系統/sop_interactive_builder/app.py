import asyncio
import base64
import json
import logging
import os
import shutil
import uuid
import wave
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

try:
    from dotenv import load_dotenv
    load_dotenv(override=False)
except ImportError:
    pass

import edge_tts
import fitz  # PyMuPDF
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from PIL import Image

# For PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
import docx

# ─── Config ───────────────────────────────────────────────────────────────────
JOBS_DIR = Path("jobs")
JOBS_DIR.mkdir(exist_ok=True)
TEMPLATES_DIR = Path("templates")
TEMPLATES_DIR.mkdir(exist_ok=True)

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

executor = ThreadPoolExecutor(max_workers=2)

# ─── Voice Configurations ───────────────────────────────────────────────────────
VOICE_GROUPS = {
    "🇹🇼 繁體中文（台灣）": [
        {"id": "zh-TW-HsiaoChenNeural", "label": "曉臻（女聲・自然）"},
        {"id": "zh-TW-HsiaoYuNeural",   "label": "曉雨（女聲・活潑）"},
        {"id": "zh-TW-YunJheNeural",    "label": "雲哲（男聲）"},
    ],
    "🇨🇳 普通話（大陸）": [
        {"id": "zh-CN-XiaoxiaoNeural",  "label": "曉曉（女聲）"},
        {"id": "zh-CN-YunxiNeural",     "label": "雲希（男聲）"},
    ],
    "🇺🇸 英文（美國）": [
        {"id": "en-US-JennyNeural",  "label": "Jenny（Female）"},
        {"id": "en-US-GuyNeural",    "label": "Guy（Male）"},
    ]
}

GEMINI_TTS_VOICES = [
    {"id": "Zephyr", "label": "Zephyr"},
    {"id": "Puck",   "label": "Puck"},
    {"id": "Charon", "label": "Charon"},
    {"id": "Kore",   "label": "Kore"}
]

# ─── Helpers ──────────────────────────────────────────────────────────────────
def parse_api_keys(raw_keys: str) -> list[str]:
    if not raw_keys:
        return []
    keys = []
    for line in raw_keys.replace(",", "\n").splitlines():
        k = line.strip()
        if k:
            keys.append(k)
    return keys

# Gemini Integration
import requests

def call_gemini_vision(api_key: str, model: str, prompt: str, image_path: str = None, text_content: str = None) -> str:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    
    parts = [{"text": prompt}]
    if image_path:
        with open(image_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode("utf-8")
        parts.append({
            "inlineData": {
                "mimeType": "image/png",
                "data": img_b64
            }
        })
    if text_content:
        parts.append({"text": f"\n\n參考內容：\n{text_content}"})
        
    payload = {"contents": [{"parts": parts}]}
    
    response = requests.post(url, json=payload, headers=headers, timeout=60)
    response.raise_for_status()
    data = response.json()
    return data["candidates"][0]["content"]["parts"][0]["text"].strip()

def pcm_to_wav(pcm_data: bytes, sample_rate: int = 24000, channels: int = 1, sample_width: int = 2) -> bytes:
    import io
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(channels)
        wf.setsampwidth(sample_width)
        wf.setframerate(sample_rate)
        wf.writeframes(pcm_data)
    return buf.getvalue()

def call_gemini_tts(api_key: str, model: str, voice: str, text: str, output_path: str) -> None:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    headers = {"Content-Type": "application/json"}
    payload = {
        "contents": [{"parts": [{"text": text}]}],
        "generationConfig": {
            "responseModalities": ["AUDIO"],
            "speechConfig": {
                "voiceConfig": {
                    "prebuiltVoiceConfig": {
                        "voiceName": voice
                    }
                }
            }
        }
    }
    response = requests.post(url, json=payload, headers=headers, timeout=120)
    response.raise_for_status()
    data = response.json()

    audio_b64 = None
    mime_type = ""
    for cand in data.get("candidates", []):
        content = cand.get("content", {})
        if isinstance(content, dict):
            for part in content.get("parts", []):
                inline = part.get("inlineData", {})
                if inline.get("mimeType", "").startswith("audio/"):
                    audio_b64 = inline["data"]
                    mime_type = inline["mimeType"]
                    break
    if not audio_b64:
        raise ValueError(f"Gemini TTS: no audio data in response. Raw: {str(data)[:300]}")

    raw_bytes = base64.b64decode(audio_b64)
    # Gemini 回傳 audio/L16 (raw PCM) 時需要包成 WAV
    if "L16" in mime_type or "pcm" in mime_type.lower():
        raw_bytes = pcm_to_wav(raw_bytes)
    with open(output_path, "wb") as f:
        f.write(raw_bytes)

# ─── FastAPI Application ──────────────────────────────────────────────────────
app = FastAPI(title="SOP Interactive Builder")
templates = Jinja2Templates(directory="templates")
app.mount("/jobs", StaticFiles(directory="jobs"), name="jobs")

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse(request=request, name="index.html")

@app.get("/api/voices")
async def get_voices():
    return {"edge": VOICE_GROUPS, "gemini": GEMINI_TTS_VOICES}

@app.get("/api/jobs")
async def list_jobs():
    """載入歷史專案清單"""
    jobs = []
    for d in JOBS_DIR.iterdir():
        if d.is_dir() and (d / "script.json").exists():
            jobs.append(d.name)
    return {"jobs": sorted(jobs, reverse=True)}

@app.get("/api/jobs/{job_id}")
async def get_job_data(job_id: str):
    job_dir = JOBS_DIR / job_id
    script_path = job_dir / "script.json"
    if not script_path.exists():
        raise HTTPException(status_code=404, detail="找不到該專案")
    with open(script_path, "r", encoding="utf-8") as f:
        return json.load(f)

@app.post("/api/extract")
async def extract_document(
    file: UploadFile = File(...),
    api_keys: str = Form(""),
    model: str = Form("gemini-3.5-flash"),
    extract_mode: str = Form("gemini")
):
    """解析 PDF / DOCX 並切分為投影片腳本"""
    import datetime, re
    safe_name = re.sub(r'[^\w\u4e00-\u9fa5\-]', '_', Path(file.filename).stem)
    job_id = f"{safe_name}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    job_dir = JOBS_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    
    keys = parse_api_keys(api_keys)
    env_keys = parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
    keys = keys or env_keys
    
    if extract_mode == "gemini" and not keys:
        extract_mode = "local"  # 自動退回到本機模式
    
    content_text = ""
    file_bytes = await file.read()
    
    local_slides = []
    
    if file.filename.lower().endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            content_text += f"\n--- 第 {i+1} 頁 ---\n" + text
            
            # 準備本機模式的腳本
            if text:
                local_slides.append({
                    "title": f"投影片 {i+1}",
                    "content": text[:150] + ("..." if len(text)>150 else ""),
                    "narration": text.replace("\n", "，")[:300]
                })
                
            # 同時產出圖片做備用/參考
            pix = page.get_pixmap()
            pix.save(str(job_dir / f"page_{i}.png"))
    elif file.filename.lower().endswith(".docx"):
        import io
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for i, p_text in enumerate(paragraphs):
            content_text += p_text + "\n"
            local_slides.append({
                "title": f"段落 {i+1}",
                "content": p_text[:150] + ("..." if len(p_text)>150 else ""),
                "narration": p_text.replace("\n", "，")[:300]
            })
    else:
        raise HTTPException(status_code=400, detail="僅支援 PDF 或 DOCX")

    if not local_slides:
        local_slides.append({
            "title": "無內容",
            "content": "無法解析出文字",
            "narration": "系統無法從此檔案解析出文字"
        })

    if extract_mode == "local":
        parsed_data = {"slides": local_slides, "quiz": []}
    else:
        # 使用 Gemini 分析並輸出 JSON 格式的大綱與旁白，以及測驗題目
        prompt = (
            "請將以下教育訓練教材內容整理成多頁的『簡報腳本』，並根據內容生成 3 到 5 題『選擇題測驗』。\n"
            "請以嚴格的 JSON 格式回傳，**絕對不能漏掉 quiz 屬性**，包含兩個主要屬性：\n"
            '1. "slides": 陣列，每個元素代表一頁簡報，包含:\n'
            '   - "title": 投影片標題\n'
            '   - "content": 投影片重點內容（條列式）\n'
            '   - "narration": 口語化的朗讀旁白腳本（適合語音播放）\n'
            '2. "quiz": 陣列，每個元素代表一題選擇題（必須提供 3~5 題），包含:\n'
            '   - "question": 題目敘述\n'
            '   - "options": 長度為 4 的字串陣列 (選項 A, B, C, D)\n'
            '   - "answer": 正確選項的索引值 (0 到 3)\n'
            "只回傳 JSON，不要加 markdown 標記。"
        )
    
        # Model fallback 清單：選到的 model 排第一，失敗自動往下試
        FALLBACK_MODELS = [
            "gemini-3.5-flash",
            "gemini-2.5-flash",
            "gemini-2.0-flash",
            "gemini-2.0-flash-lite",
        ]
        if model in FALLBACK_MODELS:
            model_queue = [model] + [m for m in FALLBACK_MODELS if m != model]
        else:
            model_queue = [model] + FALLBACK_MODELS

        success, result_text, last_err = False, "", None
        for try_model in model_queue:
            for key in keys:
                try:
                    result_text = call_gemini_vision(key, try_model, prompt, text_content=content_text)
                    logger.info(f"Gemini 成功：model={try_model}")
                    success = True
                    break
                except Exception as e:
                    last_err = e
                    logger.warning(f"Gemini 失敗：model={try_model} err={e}")
            if success:
                break

        if not success:
            raise HTTPException(status_code=500, detail=f"AI 分析失敗（已試所有 model 與 key）: {last_err}")
        
        try:
            if result_text.startswith("```json"):
                result_text = result_text.strip("`").replace("json\n", "")
            parsed_data = json.loads(result_text)
            
            if isinstance(parsed_data, list):
                parsed_data = {"slides": parsed_data, "quiz": []}
            elif "slides" not in parsed_data:
                parsed_data = {"slides": local_slides, "quiz": []}
                
            if "quiz" not in parsed_data or not isinstance(parsed_data["quiz"], list):
                parsed_data["quiz"] = []
                
            slides = parsed_data.get("slides", [])
            quiz = parsed_data.get("quiz", [])
            
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"解析 AI JSON 失敗: {str(e)}\n\n{result_text}")

    slides = parsed_data.get("slides", [])
    quiz   = parsed_data.get("quiz", [])

    job_data = {
        "job_id": job_id,
        "title": Path(file.filename).stem,   # 保留原始檔名作為標題
        "slides": slides,
        "quiz": quiz,
        "voice": "zh-TW-HsiaoChenNeural",
        "tts_engine": "edge"
    }
    with open(job_dir / "script.json", "w", encoding="utf-8") as f:
        json.dump(job_data, f, ensure_ascii=False, indent=2)
        
    return {"job_id": job_id, "slides": slides, "quiz": quiz}

@app.post("/api/generate-quiz")
async def generate_quiz_endpoint(
    api_keys: str = Form(""),
    model: str = Form("gemini-3.5-flash"),
    content: str = Form(...),
    count: int = Form(5),
):
    """根據腳本內容 AI 生成考題"""
    keys = parse_api_keys(api_keys) or parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
    if not keys:
        raise HTTPException(status_code=400, detail="需要 Gemini API 金鑰才能 AI 生成考題")

    prompt = (
        f"請根據以下教育訓練腳本內容，生成 {count} 題繁體中文選擇題測驗。\n"
        "每題必須包含 4 個選項（A/B/C/D）且只有一個正確答案。\n"
        "以 JSON 陣列格式回傳，每個元素包含：\n"
        '- "question": 題目敘述\n'
        '- "options": 長度為 4 的字串陣列\n'
        '- "answer": 正確選項索引（0=A, 1=B, 2=C, 3=D）\n'
        "只回傳 JSON 陣列，不要加 markdown 標記。\n\n"
        f"腳本內容：\n{content[:6000]}"
    )

    FALLBACK_MODELS = [
        "gemini-3.5-flash", "gemini-2.5-flash",
        "gemini-2.0-flash", "gemini-2.0-flash-lite",
    ]
    model_queue = [model] + [m for m in FALLBACK_MODELS if m != model]

    result_text, last_err, success = "", None, False
    for try_model in model_queue:
        for key in keys:
            try:
                result_text = call_gemini_vision(key, try_model, prompt)
                logger.info(f"generate-quiz 成功：model={try_model}")
                success = True
                break
            except Exception as e:
                last_err = e
                logger.warning(f"generate-quiz 失敗：model={try_model} err={e}")
        if success:
            break

    if not success:
        raise HTTPException(status_code=500, detail=f"AI 生成考題失敗: {last_err}")

    try:
        if result_text.startswith("```"):
            result_text = result_text.strip("`").replace("json\n", "").strip()
        quiz = json.loads(result_text)
        if not isinstance(quiz, list):
            quiz = quiz.get("quiz", [])
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"解析考題 JSON 失敗: {e}\n{result_text[:300]}")

    return {"quiz": quiz}

@app.post("/api/generate-audio")
async def generate_audio(
    job_id: str = Form(...),
    page_index: int = Form(...),
    text: str = Form(...),
    tts_engine: str = Form("edge"),
    voice: str = Form(...),
    api_keys: str = Form("")
):
    """生成單頁語音"""
    job_dir = JOBS_DIR / job_id
    if not job_dir.exists():
        raise HTTPException(status_code=404, detail="找不到專案")
        
    audio_ext = "wav" if tts_engine == "gemini" else "mp3"
    audio_path = job_dir / f"audio_{page_index}.{audio_ext}"
    
    try:
        if tts_engine == "gemini":
            keys = parse_api_keys(api_keys) or parse_api_keys(os.environ.get("GEMINI_API_KEYS", ""))
            if not keys:
                raise HTTPException(status_code=400, detail="Gemini TTS 需要 API 金鑰")
            success = False
            for key in keys:
                try:
                    call_gemini_tts(key, "gemini-2.5-flash-preview-tts", voice, text, str(audio_path))
                    success = True
                    break
                except:
                    continue
            if not success:
                raise ValueError("所有金鑰皆失效")
        else:
            communicate = edge_tts.Communicate(text, voice)
            await communicate.save(str(audio_path))
            
        import time
        return {"url": f"/jobs/{job_id}/audio_{page_index}.{audio_ext}?t={int(time.time())}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/save-script")
async def save_script(request: Request):
    data = await request.json()
    job_id = data.get("job_id")
    job_dir = JOBS_DIR / job_id
    if job_dir.exists():
        with open(job_dir / "script.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    return {"status": "ok"}

# ── 欄位型態正規化：Gemini 有時回傳 list，需轉成 str ──
def coerce_str(val) -> str:
    if val is None:
        return ""
    if isinstance(val, list):
        return "\n".join(str(v) for v in val)
    return str(val)


@app.post("/api/export-pptx")
async def export_pptx(job_id: str = Form(...)):
    import time, re
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt, Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    job_dir = JOBS_DIR / job_id
    with open(job_dir / "script.json", "r", encoding="utf-8") as f:
        job_data = json.load(f)

    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    W = prs.slide_width
    H = prs.slide_height

    C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
    C_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    C_BG     = RGBColor(0xF1, 0xF5, 0xF9)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
    C_MUTED  = RGBColor(0x64, 0x74, 0x8B)

    def hex_rgb(c):
        return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

    def force_fill(shape, color):
        """直接操作 XML 注入 solidFill，繞過 theme 繼承"""
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(shape._element, qn("p:spPr"))
        for tag in [qn("a:noFill"), qn("a:gradFill"), qn("a:solidFill"), qn("a:pattFill")]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        sr = etree.SubElement(sf, qn("a:srgbClr"))
        sr.set("val", hex_rgb(color))

    def add_rect(sl, l, t, w, h, color):
        shp = sl.shapes.add_shape(1, int(l), int(t), int(w), int(h))
        force_fill(shp, color)
        shp.line.fill.background()
        return shp

    def add_text(sl, l, t, w, h, text, size=20, bold=False, color=None,
                 align=PP_ALIGN.LEFT, line_sp=1.4):
        if color is None:
            color = C_TEXT
        tb = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in text.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spc   = etree.SubElement(lnSpc, qn("a:spcPct"))
            spc.set("val", str(int(line_sp * 100000)))
            r = p.add_run()
            r.text = line
            r.font.name  = "\u5fae\u8edf\u6b63\u9ed1\u9ad4"
            r.font.size  = Pt(size)
            r.font.bold  = bold
            r.font.color.rgb = color
        return tb

    blank = prs.slide_layouts[6]  # Blank

    # ─── 封面頁 ───────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = C_DARK
    add_rect(cover, 0, H * 0.60, W, H * 0.40, C_BG)
    add_rect(cover, 0, H * 0.10, Emu(14000), H * 0.45, C_ACCENT)
    title_str = job_data.get("title", job_id)
    add_text(cover, Emu(100000), H * 0.12, W * 0.92, H * 0.42,
             title_str, size=38, bold=True, color=C_WHITE, line_sp=1.4)
    add_text(cover, Emu(100000), H * 0.67, W * 0.60, H * 0.18,
             "\u54e1\u5de5\u6559\u80b2\u8a13\u7df4\u6559\u6750", size=17, color=C_MUTED)

    # ─── 內容頁 ───────────────────────────────────────────
    for slide_data in job_data.get("slides", []):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = C_BG
        add_rect(sl, 0, 0, W, H * 0.19, C_DARK)
        add_rect(sl, 0, H * 0.19, Emu(13000), H * 0.81, C_ACCENT)

        t_text = coerce_str(slide_data.get("title", ""))
        add_text(sl, Emu(70000), H * 0.02, W * 0.94, H * 0.16,
                 t_text, size=30, bold=True, color=C_WHITE, line_sp=1.2)

        raw = coerce_str(slide_data.get("content", ""))
        bullets = [s.strip() for s in re.split(r"[\u3002\uff1b;\n]", raw) if s.strip()]
        body = "\n".join(f"\u2022  {b}" for b in bullets) if bullets else raw
        add_text(sl, Emu(80000), H * 0.23, W * 0.94, H * 0.72,
                 body, size=19, color=C_TEXT, line_sp=1.7)

        sl.notes_slide.notes_text_frame.text = coerce_str(slide_data.get("narration", ""))

    out_path = job_dir / f"{job_id}.pptx"
    prs.save(str(out_path))
    return {"url": f"/jobs/{job_id}/{job_id}.pptx?t={int(time.time())}"}


@app.post("/api/export-web")
async def export_web(job_id: str = Form(...)):
    import zipfile
    job_dir = JOBS_DIR / job_id
    with open(job_dir / "script.json", "r", encoding="utf-8") as f:
        job_data = json.load(f)

    template_path = TEMPLATES_DIR / "player_template.html"
    if not template_path.exists():
        raise HTTPException(status_code=500, detail="Player template missing")

    with open(template_path, "r", encoding="utf-8") as f:
        html = f.read()

    title = job_data.get("title", job_id.rsplit("_", 2)[0].replace("_", " "))

    slides_config = []
    audio_files_to_pack = []
    for i, s in enumerate(job_data.get("slides", [])):
        audio_ext = "wav" if job_data.get("tts_engine") == "gemini" else "mp3"
        audio_file = f"audio_{i}.{audio_ext}"
        has_audio = (job_dir / audio_file).exists()
        slides_config.append({
            "title":     coerce_str(s.get("title", "")),
            "content":   coerce_str(s.get("content", "")),
            "narration": coerce_str(s.get("narration", "")),
            "audio":     audio_file if has_audio else None
        })
        if has_audio:
            audio_files_to_pack.append(audio_file)

    html = html.replace("{{TITLE}}",      title)
    html = html.replace("{{SLIDES_JSON}}", json.dumps(slides_config, ensure_ascii=False))
    html = html.replace("{{QUIZ_JSON}}",   json.dumps(job_data.get("quiz", []), ensure_ascii=False))

    # ── 1. 寫入 index.html
    index_path = job_dir / "index.html"
    with open(index_path, "w", encoding="utf-8") as f:
        f.write(html)

    # ── 2. serve_intranet.ps1 — 用 Python list join 組合，完全避開引號跳脫問題
    ps1_lines = [
        "$port = 18080",
        "",
        "# 取得本機對外內網 IP",
        "$socket = New-Object System.Net.Sockets.UdpClient",
        "$ip = $null",
        "try {",
        "    $socket.Connect('8.8.8.8', 80)",
        "    $ip = $socket.Client.LocalEndPoint.Address.IPAddressToString",
        "} catch {} finally { if ($socket) { $socket.Close() } }",
        "if (-not $ip) {",
        "    $ip = (Get-NetIPAddress | Where-Object {",
        "        $_.AddressFamily -eq 'InterNetwork' -and",
        "        $_.IPAddress -notmatch '^127\\.' -and",
        "        $_.IPAddress -notmatch '^169\\.254\\.'",
        "    } | Select-Object -First 1).IPAddress",
        "}",
        "if (-not $ip) { $ip = '127.0.0.1' }",
        "",
        "$localIP = [System.Net.IPAddress]::Any",
        "$listener = $null",
        "$bound = $false",
        "while (-not $bound -and $port -lt 19000) {",
        "    try {",
        "        $listener = New-Object System.Net.Sockets.TcpListener($localIP, $port)",
        "        $listener.Start()",
        "        $bound = $true",
        "    } catch { $port++ }",
        "}",
        "if (-not $bound) {",
        f"    Write-Host '  找不到可用連接埠 (18080-19000)。' -ForegroundColor Red",
        "    Read-Host '按 Enter 結束'; exit",
        "}",
        "",
        "Write-Host '==================================================' -ForegroundColor Cyan",
        "Write-Host '  員工教育訓練測驗系統 -- 本機內網伺服器' -ForegroundColor Cyan",
        "Write-Host '==================================================' -ForegroundColor Cyan",
        "Write-Host '  請勿關閉此視窗，關閉即結束服務。'",
        "Write-Host '  同仁需連至同一 Wi-Fi 或公司內網。'",
        "Write-Host ''",
        "Write-Host \"  http://$($ip):$($port)/index.html\" -ForegroundColor Green",
        "Write-Host '==================================================' -ForegroundColor Cyan",
        "",
        "$currentDir = $PSScriptRoot",
        "if (-not $currentDir) { $currentDir = (Get-Location).Path }",
        "",
        "function EscCsv($s) { return '\"' + ($s -replace '\"','\"\"') + '\"' }",
        "",
        "while ($true) {",
        "    try {",
        "        if (-not $listener.Pending()) { Start-Sleep -Milliseconds 100; continue }",
        "        $client  = $listener.AcceptTcpClient()",
        "        $stream  = $client.GetStream()",
        "        $reader  = New-Object System.IO.StreamReader($stream)",
        "        $reqLine = $reader.ReadLine()",
        "",
        "        if ($reqLine -match '^(GET|POST)\\s+(/[^\\s\\?]*)\\??[^\\s]*\\s+HTTP') {",
        "            $method  = $Matches[1]",
        "            $urlPath = [System.Uri]::UnescapeDataString($Matches[2])",
        "            if ($urlPath -eq '/') { $urlPath = '/index.html' }",
        "",
        "            if ($method -eq 'POST' -and $urlPath -eq '/api/submit') {",
        "                $hdrs = @{}",
        "                while ($line = $reader.ReadLine()) {",
        "                    if ($line -eq '') { break }",
        "                    if ($line -match '^([^:]+):\\s*(.*)$') { $hdrs[$Matches[1].ToLower()] = $Matches[2].Trim() }",
        "                }",
        "                $cLen = 0",
        "                if ($hdrs.ContainsKey('content-length')) { [int]::TryParse($hdrs['content-length'], [ref]$cLen) | Out-Null }",
        "                $body = ''",
        "                if ($cLen -gt 0) {",
        "                    $buf  = New-Object System.Char[] $cLen",
        "                    $read = $reader.Read($buf, 0, $cLen)",
        "                    $body = New-Object System.String($buf, 0, $read)",
        "                }",
        "                try {",
        "                    $rec     = $body | ConvertFrom-Json",
        "                    $csvPath = Join-Path $currentDir 'results.csv'",
        "                    $enc     = New-Object System.Text.UTF8Encoding($true)",
        "                    # 計算題目數",
        "                    $qCount = ($rec.PSObject.Properties | Where-Object { $_.Name -match '^q\\d+_answer$' }).Count",
        "                    # 若 CSV 不存在，建立標頭列",
        "                    if (-not (Test-Path $csvPath)) {",
        "                        $hdr = '時間戳記,姓名,對題數,得分'",
        "                        for ($i = 1; $i -le $qCount; $i++) {",
        "                            $qt   = $rec.\"q${i}_question\"",
        "                            $hdr += ',' + (EscCsv \"第${i}題: $qt\")",
        "                        }",
        "                        [System.IO.File]::WriteAllText($csvPath, \"$hdr`r`n\", $enc)",
        "                    }",
        "                    # 組資料列",
        "                    $row  = (EscCsv $rec.timestamp) + ',' + (EscCsv $rec.name)",
        "                    $row += ',' + (EscCsv \"$($rec.correctCount) / $($rec.total)\")",
        "                    $row += ',' + (EscCsv \"$($rec.score) 分\")",
        "                    for ($i = 1; $i -le $qCount; $i++) {",
        "                        $row += ',' + (EscCsv $rec.\"q${i}_answer\")",
        "                    }",
        "                    [System.IO.File]::AppendAllText($csvPath, \"$row`r`n\", $enc)",
        "                    $rb  = [System.Text.Encoding]::UTF8.GetBytes('{\"status\":\"ok\"}')",
        "                    $rh  = \"HTTP/1.1 200 OK`r`nContent-Type: application/json`r`nContent-Length: $($rb.Length)`r`nAccess-Control-Allow-Origin: *`r`nConnection: close`r`n`r`n\"",
        "                    $stream.Write([System.Text.Encoding]::UTF8.GetBytes($rh), 0, $rh.Length)",
        "                    $stream.Write($rb, 0, $rb.Length)",
        "                    Write-Host \"[收到] $($rec.name) -- 得分: $($rec.score)分  已寫入 results.csv\" -ForegroundColor Yellow",
        "                } catch {",
        "                    $eb = [System.Text.Encoding]::UTF8.GetBytes('{\"status\":\"error\"}')",
        "                    $eh = \"HTTP/1.1 500 Error`r`nContent-Length: $($eb.Length)`r`nConnection: close`r`n`r`n\"",
        "                    $stream.Write([System.Text.Encoding]::UTF8.GetBytes($eh), 0, $eh.Length)",
        "                    $stream.Write($eb, 0, $eb.Length)",
        "                    Write-Host \"[錯誤] $($_.Exception.Message)\" -ForegroundColor Red",
        "                }",
        "            } else {",
        "                $fp = Join-Path $currentDir $urlPath",
        "                if (Test-Path $fp -PathType Leaf) {",
        "                    $ext = [System.IO.Path]::GetExtension($fp).ToLower()",
        "                    $ct  = switch ($ext) {",
        "                        '.html' { 'text/html; charset=utf-8' }",
        "                        '.mp3'  { 'audio/mpeg' }",
        "                        '.wav'  { 'audio/wav' }",
        "                        '.css'  { 'text/css' }",
        "                        '.js'   { 'application/javascript' }",
        "                        '.json' { 'application/json' }",
        "                        default { 'application/octet-stream' }",
        "                    }",
        "                    $bytes = [System.IO.File]::ReadAllBytes($fp)",
        "                    $rh    = \"HTTP/1.1 200 OK`r`nContent-Type: $ct`r`nContent-Length: $($bytes.Length)`r`nAccess-Control-Allow-Origin: *`r`nConnection: close`r`n`r`n\"",
        "                    $stream.Write([System.Text.Encoding]::UTF8.GetBytes($rh), 0, $rh.Length)",
        "                    $stream.Write($bytes, 0, $bytes.Length)",
        "                } else {",
        "                    $e4 = [System.Text.Encoding]::UTF8.GetBytes('404 Not Found')",
        "                    $h4 = \"HTTP/1.1 404 Not Found`r`nContent-Length: $($e4.Length)`r`nConnection: close`r`n`r`n\"",
        "                    $stream.Write([System.Text.Encoding]::UTF8.GetBytes($h4), 0, $h4.Length)",
        "                    $stream.Write($e4, 0, $e4.Length)",
        "                }",
        "            }",
        "        }",
        "        $stream.Close(); $client.Close()",
        "    } catch { Write-Host \"[例外] $_\" -ForegroundColor DarkGray }",
        "}",
    ]
    ps1_content = "\n".join(ps1_lines)
    ps1_path = job_dir / "serve_intranet.ps1"
    with open(ps1_path, "w", encoding="utf-8-sig") as f:  # BOM 讓 PowerShell 正確讀 UTF-8
        f.write(ps1_content)

    # ── 3. 寫入 bat

    bat_content = """@echo off
chcp 65001 >nul 2>&1
echo.
echo  ============================================
echo   員工教育訓練測驗系統 — 啟動中...
echo  ============================================
echo.
echo  請稍候，伺服器啟動後，瀏覽器將自動開啟。
echo  若未自動開啟，請手動在瀏覽器輸入伺服器顯示的網址。
echo  *** 請勿關閉此視窗，關閉即停止服務 ***
echo.
powershell -ExecutionPolicy Bypass -File "%~dp0serve_intranet.ps1"
pause
"""
    bat_path = job_dir / "點我啟動內網伺服器(Windows免安裝).bat"
    with open(bat_path, "w", encoding="utf-8-sig") as f:
        f.write(bat_content)

    # ── 4. 寫入 apps_script_code.gs（Google Sheets 雲端同步腳本）
    gs_content = f"""// ============================================================
// 📋 Google Apps Script — 員工教育訓練測驗成績收集系統
// SOP：{title}
// ============================================================

const SHEET_NAME = '作答紀錄';

function doPost(e) {{
  try {{
    const data = JSON.parse(e.postData.contents);
    const ss = SpreadsheetApp.getActiveSpreadsheet();
    let sheet = ss.getSheetByName(SHEET_NAME);
    if (!sheet) {{ sheet = ss.insertSheet(SHEET_NAME); }}
    const name         = data.name || '未知';
    const score        = data.score !== undefined ? data.score : 0;
    const correctCount = data.correctCount !== undefined ? data.correctCount : 0;
    const total        = data.total !== undefined ? data.total : 0;
    const timestamp    = data.timestamp || new Date().toLocaleString('zh-TW', {{ timeZone: 'Asia/Taipei' }});
    const qAnswers = [];
    let qIndex = 1;
    while (data['q' + qIndex + '_answer'] !== undefined) {{
      qAnswers.push(data['q' + qIndex + '_answer']);
      qIndex++;
    }}
    if (sheet.getLastRow() === 0) {{
      const headers = ['時間戳記', '姓名', '對題數', '得分'];
      for (let i = 1; i < qIndex; i++) {{
        const q = data['q' + i + '_question'] || ('第 ' + i + ' 題');
        headers.push(q);
      }}
      sheet.appendRow(headers);
      const range = sheet.getRange(1, 1, 1, headers.length);
      range.setBackground('#4F46E5');
      range.setFontColor('#FFFFFF');
      range.setFontWeight('bold');
    }}
    const rowData = [timestamp, name, correctCount + ' / ' + total, score + ' 分'];
    qAnswers.forEach(ans => rowData.push(ans));
    sheet.appendRow(rowData);
    return ContentService
      .createTextOutput(JSON.stringify({{ status: 'ok', message: '已成功存入雲端試算表！' }}))
      .setMimeType(ContentService.MimeType.JSON);
  }} catch (err) {{
    return ContentService
      .createTextOutput(JSON.stringify({{ status: 'error', message: err.toString() }}))
      .setMimeType(ContentService.MimeType.JSON);
  }}
}}
"""
    gs_path = job_dir / "apps_script_code.gs"
    with open(gs_path, "w", encoding="utf-8") as f:
        f.write(gs_content)

    # ── 5. 寫入 README.md
    readme_content = f"""# 員工教育訓練系統 — 說明書
## SOP：{title}

本資料夾由「SOP 互動簡報產生系統」自動產生，包含簡報閱讀、語音朗讀、測驗及本機/雲端成績收集。

---

## 📁 檔案說明
1. **`index.html`**：主測驗網頁（直接用瀏覽器開啟即可播放簡報並作答）。
2. **`serve_intranet.ps1`**：本機伺服器（接收成績並存入 results.csv）。
3. **`點我啟動內網伺服器(Windows免安裝).bat`**：一鍵啟動。
4. **`apps_script_code.gs`**：Google Sheets 雲端同步腳本（選用）。

---

## 🚀 使用方式（本機內網）
1. 解壓縮此資料夾。
2. 雙擊執行 **`點我啟動內網伺服器(Windows免安裝).bat`**。
3. 視窗會顯示內網網址（如：`http://192.168.1.100:18080/index.html`）。
4. 將網址傳給同仁（需連同一 Wi-Fi）。
5. 同仁完成作答後，成績自動存入 **`results.csv`**（可直接用 Excel 開啟）。

## ☁️ 可選：Google Sheets 雲端同步
1. 建立新的 Google 試算表。
2. 在「擴充功能 > Apps Script」中貼上 `apps_script_code.gs` 的程式碼並部署為網頁應用程式。
3. 在測驗頁面右上角「系統設定」中貼入部署網址。
"""
    readme_path = job_dir / "README.md"
    with open(readme_path, "w", encoding="utf-8") as f:
        f.write(readme_content)

    # ── 6. 預建 results.csv（含標頭，供 Excel 直接開啟）
    quiz_list = job_data.get("quiz", [])
    csv_headers = ["時間戳記", "姓名", "對題數", "得分"]
    for qi, q in enumerate(quiz_list):
        q_text = q.get("question", f"第{qi+1}題").replace('"', '""')
        csv_headers.append(f'第{qi+1}題: {q_text}')
    # 用 BOM UTF-8 確保 Excel 開啟不亂碼
    csv_header_line = ",".join(f'"{h}"' for h in csv_headers) + "\r\n"
    csv_path = job_dir / "results.csv"
    with open(csv_path, "w", encoding="utf-8-sig") as f:
        f.write(csv_header_line)

    # ── 6.5 預建 quiz_paper.csv（紙本考卷參考）
    quiz_paper_path = job_dir / "quiz_paper.csv"
    with open(quiz_paper_path, "w", encoding="utf-8-sig") as f:
        f.write("題號,題目,選項A,選項B,選項C,選項D,解答\n")
        letters = ["A", "B", "C", "D", "E"]
        for qi, q in enumerate(quiz_list):
            q_text = q.get("question", "").replace('"', '""')
            opts = q.get("options", [])
            opt_texts = [(opts[i].replace('"', '""') if i < len(opts) else "") for i in range(4)]
            ans_idx = q.get("answer", 0)
            if isinstance(ans_idx, str) and ans_idx.isdigit():
                ans_idx = int(ans_idx)
            ans_letter = letters[ans_idx] if isinstance(ans_idx, int) and 0 <= ans_idx < len(letters) else ""
            ans_text = opts[ans_idx].replace('"', '""') if isinstance(ans_idx, int) and 0 <= ans_idx < len(opts) else ""
            f.write(f'"{qi+1}","{q_text}","{opt_texts[0]}","{opt_texts[1]}","{opt_texts[2]}","{opt_texts[3]}","{ans_letter}. {ans_text}"\n')

    # ── 7. 打包 ZIP
    zip_filename = f"{job_id}_教育訓練測驗套件.zip"
    zip_path = job_dir / zip_filename
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.write(index_path,  arcname="index.html")
        zf.write(ps1_path,    arcname="serve_intranet.ps1")
        zf.write(bat_path,    arcname="點我啟動內網伺服器(Windows免安裝).bat")
        zf.write(gs_path,     arcname="apps_script_code.gs")
        zf.write(readme_path, arcname="README.md")
        zf.write(csv_path,    arcname="results.csv")
        zf.write(quiz_paper_path, arcname="紙本考卷與解答.csv")
        for af in audio_files_to_pack:
            zf.write(job_dir / af, arcname=af)

    import time
    return {"url": f"/jobs/{job_id}/{zip_filename}?t={int(time.time())}"}
