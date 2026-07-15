import os
import re
import sys
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Ensure UTF-8 output
try:
    sys.stdout.reconfigure(encoding='utf-8')
except AttributeError:
    pass

# Dynamically locate root directory and dependent folders (parent of portal_tools)
tools_dir = os.path.dirname(os.path.abspath(__file__))
root_dir = os.path.dirname(tools_dir)

# ----------------- Configuration & Themes -----------------
BG_COLOR = RGBColor(11, 14, 20)        # #0B0E14 (Premium Dark)
TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF
TEXT_LIGHT_GRAY = RGBColor(208, 211, 212) # #D0D3D4
TEXT_MUTED = RGBColor(160, 170, 180)   # #A0AAB4
CARD_BG = RGBColor(22, 27, 34)         # #161B22 (Card Fill)
CARD_BORDER = RGBColor(48, 54, 61)     # #30363D (Card Border)
CODE_BG = RGBColor(1, 4, 9)            # #010409 (Terminal Black)
ACCENT_GOLD = RGBColor(255, 214, 0)    # #FFD600 (Gold accent)

# Accent colors for categories
THEMES = {
    "第一類_核心網頁與互動系統": {
        "accent": RGBColor(0, 229, 255),      # Neon Cyan
        "label": "第一類：核心網頁與互動系統",
        "bg_accent": RGBColor(0, 40, 50)
    },
    "第二類_生產管理與API串接": {
        "accent": RGBColor(0, 230, 118),      # Neon Green
        "label": "第二類：生產管理與API串接",
        "bg_accent": RGBColor(0, 40, 20)
    },
    "第三類_AI代理與指南企劃": {
        "accent": RGBColor(255, 145, 0),      # Neon Orange
        "label": "第三類：AI代理與指南企劃",
        "bg_accent": RGBColor(50, 25, 0)
    },
    "default": {
        "accent": RGBColor(212, 175, 55),     # Gold
        "label": "專案使用說明書",
        "bg_accent": RGBColor(40, 35, 10)
    }
}

# Image paths for specific projects
IMAGE_DIR = os.path.join(root_dir, "images")
SPECIAL_PROJECT_IMAGES = {
    "hongsheng-web": {
        "cover": os.path.join(IMAGE_DIR, "hs_cover.png"),
        "roles": os.path.join(IMAGE_DIR, "hs_roles.png"),
        "green": os.path.join(IMAGE_DIR, "hs_scan_green.png"),
        "red": os.path.join(IMAGE_DIR, "hs_scan_red.png"),
        "qc": os.path.join(IMAGE_DIR, "hs_qc_dashboard.png")
    },
    "test": {
        "cover": os.path.join(IMAGE_DIR, "cover.png"),
        "day1": os.path.join(IMAGE_DIR, "day1.png"),
        "day2": os.path.join(IMAGE_DIR, "day2.png"),
        "day3": os.path.join(IMAGE_DIR, "day3.png"),
        "day4": os.path.join(IMAGE_DIR, "day4.png"),
        "day5": os.path.join(IMAGE_DIR, "day5.png"),
        "map": os.path.join(IMAGE_DIR, "map.png")
    }
}

# ----------------- Markdown Parser -----------------
def parse_markdown(filepath):
    """
    Parses our standard project manual markdown files into structured data.
    """
    with open(filepath, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    title = "未命名專案"
    github_link = ""
    sections = []
    
    current_section = None
    in_code_block = False
    code_content = []
    code_lang = ""
    
    for line in lines:
        line_str = line.strip()
        
        # 1. Parse Title
        if line.startswith('# '):
            title = line[2:].strip()
            # Clean up title if it has trailing parts
            title = re.sub(r'\s*[-—]\s*操作說明書', '', title)
            title = re.sub(r'\s*[-—]\s*說明書', '', title)
            continue
            
        # 2. Parse GitHub Link
        if line.startswith('>') and 'GitHub' in line:
            # Extract URL from markdown link [text](url)
            match = re.search(r'\[.*?\]\((.*?)\)', line)
            if match:
                github_link = match.group(1)
            continue
            
        # 3. Parse Code Blocks
        if line_str.startswith('```'):
            if in_code_block:
                in_code_block = False
                if current_section:
                    current_section['items'].append({
                        'type': 'code',
                        'lang': code_lang,
                        'content': '\n'.join(code_content)
                    })
                code_content = []
                code_lang = ""
            else:
                in_code_block = True
                code_lang = line_str[3:].strip()
            continue
            
        if in_code_block:
            code_content.append(line.rstrip('\n'))
            continue
            
        # 4. Parse Sections
        if line.startswith('## '):
            if current_section:
                sections.append(current_section)
            sec_title = line[3:].strip()
            current_section = {
                'title': sec_title,
                'items': []
            }
            continue
            
        # 5. Parse content within section
        if current_section is not None:
            if not line_str:
                continue
            # Check list items
            if line_str.startswith('- ') or line_str.startswith('* '):
                # Bullet list item
                text = line_str[2:].strip()
                # Remove bold markers like **text** -> text
                text = re.sub(r'\*\*(.*?)\*\*概念?', r'\1', text)
                text = text.replace('**', '')
                current_section['items'].append({
                    'type': 'bullet',
                    'content': text
                })
            elif re.match(r'^\d+\.\s+', line_str):
                # Numbered list item
                match = re.match(r'^(\d+)\.\s+(.*)', line_str)
                num = match.group(1)
                text = match.group(2)
                text = text.replace('**', '')
                current_section['items'].append({
                    'type': 'numbered',
                    'num': num,
                    'content': text
                })
            else:
                # Plain paragraph
                text = line_str.replace('**', '')
                current_section['items'].append({
                    'type': 'paragraph',
                    'content': text
                })
                
    if current_section:
        sections.append(current_section)
        
    return title, github_link, sections

# ----------------- PPTX Drawer Utilities -----------------
def set_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, accent_color, subtitle_text=None):
    # Main Header Container
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Microsoft JhengHei'
        p2.font.size = Pt(12)
        p2.font.color.rgb = accent_color
        p2.space_before = Pt(5)

def draw_card(slide, left, top, width, height, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# ----------------- Individual Slide Builders -----------------
def build_cover_slide(prs, title, github_link, theme, category_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    set_dark_background(slide)
    accent = theme["accent"]
    
    # 1. Left Color Accent Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    
    # 2. Main Title Text Box
    tf = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5)).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    # Category tag
    p_tag = tf.paragraphs[0]
    p_tag.text = theme["label"].upper()
    p_tag.font.name = 'Microsoft JhengHei'
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = accent
    p_tag.space_after = Pt(20)
    
    # Large Project Title
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.name = 'Microsoft JhengHei'
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(15)
    
    # Subtitle / Manual Label
    p_sub = tf.add_paragraph()
    p_sub.text = "系統操作與維護說明書  |  System Operations Manual"
    p_sub.font.name = 'Microsoft JhengHei'
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_after = Pt(40)
    
    # GitHub URL
    if github_link:
        p_git = tf.add_paragraph()
        p_git.text = f"🌐 GitHub: {github_link}"
        p_git.font.name = 'Consolas'
        p_git.font.size = Pt(11)
        p_git.font.color.rgb = accent

def build_intro_slide(prs, section, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    accent = theme["accent"]
    build_title = "01 " + section['title']
    add_title(slide, build_title, accent, "專案核心定位與簡介概述")
    
    # Gather paragraphs
    paragraphs = [item['content'] for item in section['items'] if item['type'] == 'paragraph']
    bullets = [item['content'] for item in section['items'] if item['type'] in ('bullet', 'numbered')]
    
    # Draw Background Card for text
    draw_card(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    
    # Add Text Frame inside card area
    tf = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(4.0)).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    # Title badge in card
    p_badge = tf.paragraphs[0]
    p_badge.text = "🎯 系統定位 & 簡介"
    p_badge.font.name = 'Microsoft JhengHei'
    p_badge.font.size = Pt(16)
    p_badge.font.bold = True
    p_badge.font.color.rgb = accent
    p_badge.space_after = Pt(15)
    
    for para in paragraphs:
        p = tf.add_paragraph()
        p.text = para
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT_GRAY
        p.space_after = Pt(15)
        p.line_spacing = 1.2
        
    for bull in bullets[:4]: # Limit to prevent overflow
        p = tf.add_paragraph()
        p.text = f"•  {bull}"
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT_GRAY
        p.space_after = Pt(8)

def build_features_slide(prs, section, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    accent = theme["accent"]
    add_title(slide, "02 " + section['title'], accent, "系統具備之核心功能與運作特色")
    
    bullets = [item['content'] for item in section['items'] if item['type'] in ('bullet', 'numbered', 'paragraph')]
    
    if not bullets:
        return
        
    # Grid Layout: 3 Columns
    col_width = Inches(3.6)
    col_height = Inches(4.5)
    top_pos = Inches(2.0)
    left_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]
    
    # Process features in groups of 3
    num_cols = min(len(bullets), 3)
    
    for i in range(num_cols):
        left_pos = left_positions[i]
        bullet_text = bullets[i]
        
        # Split title and description if contains ':' or '：'
        split_chars = [':', '：']
        feature_title = f"功能特色 {i+1}"
        feature_desc = bullet_text
        
        for char in split_chars:
            if char in bullet_text:
                parts = bullet_text.split(char, 1)
                feature_title = parts[0].strip()
                feature_desc = parts[1].strip()
                break
                
        # Draw Card
        draw_card(slide, left_pos, top_pos, col_width, col_height)
        
        # Add colored top border stripe to card
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, col_width, Inches(0.12))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        
        # Add Text Frame inside card
        tf = slide.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.4), col_width - Inches(0.5), col_height - Inches(0.6)).text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        # Number badge
        p_num = tf.paragraphs[0]
        p_num.text = f"0{i+1}."
        p_num.font.name = 'Consolas'
        p_num.font.size = Pt(24)
        p_num.font.bold = True
        p_num.font.color.rgb = accent
        p_num.space_after = Pt(10)
        
        # Feature Title
        p_title = tf.add_paragraph()
        p_title.text = feature_title
        p_title.font.name = 'Microsoft JhengHei'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.space_after = Pt(15)
        
        # Feature Desc
        p_desc = tf.add_paragraph()
        p_desc.text = feature_desc
        p_desc.font.name = 'Microsoft JhengHei'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.line_spacing = 1.3

def build_tech_stack_slide(prs, section, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    accent = theme["accent"]
    add_title(slide, "03 " + section['title'], accent, "系統所採用的核心開發語言與架構組件")
    
    # Collect items
    raw_items = []
    for item in section['items']:
        if item['type'] in ('bullet', 'numbered', 'paragraph'):
            # Split by commas or bullets
            content = item['content']
            parts = re.split(r'[,，、]', content)
            raw_items.extend([p.strip() for p in parts if p.strip()])
            
    # Remove duplicates and clean
    techs = []
    for t in raw_items:
        t_clean = re.sub(r'^[•\-\d\.\s]+', '', t)
        # Remove bold formatting if any
        t_clean = t_clean.replace('*', '')
        if t_clean and t_clean not in techs and len(t_clean) < 30:
            techs.append(t_clean)
            
    if not techs:
        # Fallback to general techs if not parsed
        techs = ["HTML5", "CSS3", "JavaScript (ES6)", "Python", "Git"]
        
    # Draw Tech Badges (Grid system)
    card_width = Inches(2.2)
    card_height = Inches(1.1)
    
    x_start = Inches(0.8)
    y_start = Inches(2.2)
    x_gap = Inches(2.4)
    y_gap = Inches(1.3)
    
    cols = 5
    for idx, tech in enumerate(techs[:15]): # Limit to 15 badges
        row = idx // cols
        col = idx % cols
        
        bx = x_start + (col * x_gap)
        by = y_start + (row * y_gap)
        
        # Draw badge shape
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, card_width, card_height)
        badge.fill.solid()
        badge.fill.fore_color.rgb = CARD_BG
        badge.line.color.rgb = accent
        badge.line.width = Pt(1.5)
        
        tf = badge.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = tech
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(2)
        
        # Subtext badge details
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = "CORE COMPONENT"
        p2.font.name = 'Consolas'
        p2.font.size = Pt(8)
        p2.font.color.rgb = TEXT_MUTED

    # Draw Architecture flow diagram at bottom
    if len(techs) > 0:
        draw_card(slide, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.5))
        tf_flow = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(11.0), Inches(1.1)).text_frame
        p_flow = tf_flow.paragraphs[0]
        p_flow.text = "🔄 系統架構流向："
        p_flow.font.name = 'Microsoft JhengHei'
        p_flow.font.size = Pt(13)
        p_flow.font.bold = True
        p_flow.font.color.rgb = accent
        p_flow.space_after = Pt(8)
        
        p_flow_desc = tf_flow.add_paragraph()
        p_flow_desc.text = " [前端使用者 UI 介面]  ───── ( API / 業務邏輯控制 ) ─────► [後端資料庫 / 雲端儲存與備份]"
        p_flow_desc.font.name = 'Microsoft JhengHei'
        p_flow_desc.font.size = Pt(14)
        p_flow_desc.font.bold = True
        p_flow_desc.font.color.rgb = TEXT_WHITE

def build_steps_slide(prs, section, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    accent = theme["accent"]
    add_title(slide, "04 " + section['title'], accent, "本機安裝部署與日常操作指引步驟")
    
    steps = [item['content'] for item in section['items'] if item['type'] in ('bullet', 'numbered', 'paragraph')]
    if not steps:
        return
        
    # Draw vertical timeline style or horizontal cards
    if len(steps) <= 4:
        # Horizontal flow cards
        card_width = Inches(2.7)
        card_height = Inches(4.5)
        top_pos = Inches(2.0)
        left_positions = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
        
        for i, step in enumerate(steps[:4]):
            left_pos = left_positions[i]
            
            # Card
            draw_card(slide, left_pos, top_pos, card_width, card_height)
            
            # Step circle badge
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_pos + Inches(0.2), top_pos - Inches(0.3), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = accent
            circle.line.fill.background()
            
            ctf = circle.text_frame
            cp = ctf.paragraphs[0]
            cp.text = str(i+1)
            cp.font.name = 'Consolas'
            cp.font.size = Pt(16)
            cp.font.bold = True
            cp.font.color.rgb = BG_COLOR
            cp.alignment = PP_ALIGN.CENTER
            
            # Text Inside
            tf = slide.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.6), card_width - Inches(0.4), card_height - Inches(0.8)).text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            p = tf.paragraphs[0]
            p.text = f"步驟 {i+1}"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(15)
            
            p_desc = tf.add_paragraph()
            p_desc.text = step
            p_desc.font.name = 'Microsoft JhengHei'
            p_desc.font.size = Pt(11)
            p_desc.font.color.rgb = TEXT_LIGHT_GRAY
            p_desc.line_spacing = 1.3
            
            # Connector arrow (except last)
            if i < len(steps[:4]) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_pos + card_width + Inches(0.05), top_pos + Inches(2.0), Inches(0.2), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = CARD_BORDER
                arrow.line.fill.background()
    else:
        # Vertical list inside a nice card
        draw_card(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        tf = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(4.2)).text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        for i, step in enumerate(steps[:6]): # Show up to 6 steps
            p = tf.add_paragraph()
            p.text = f" {i+1} ️⃣   {step}"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(12)
            p.line_spacing = 1.2

def build_extra_slide(prs, section, theme):
    """
    Builds a slide for extra headings found in the markdown files
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    accent = theme["accent"]
    add_title(slide, section['title'], accent, "系統其他重要規格與說明")
    
    # Separate content types
    paragraphs = []
    bullets = []
    code_blocks = []
    
    for item in section['items']:
        if item['type'] == 'paragraph':
            paragraphs.append(item['content'])
        elif item['type'] in ('bullet', 'numbered'):
            bullets.append(item['content'])
        elif item['type'] == 'code':
            code_blocks.append(item)
            
    # Layout logic: if there is a code block, show split layout
    if code_blocks:
        code_item = code_blocks[0]
        # Left Panel (Text)
        draw_card(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
        tf_left = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2)).text_frame
        tf_left.word_wrap = True
        tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = Inches(0)
        
        p_badge = tf_left.paragraphs[0]
        p_badge.text = "📝 說明描述"
        p_badge.font.name = 'Microsoft JhengHei'
        p_badge.font.size = Pt(14)
        p_badge.font.bold = True
        p_badge.font.color.rgb = accent
        p_badge.space_after = Pt(15)
        
        for para in paragraphs[:3]:
            p = tf_left.add_paragraph()
            p.text = para
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_LIGHT_GRAY
            p.space_after = Pt(10)
            
        for bull in bullets[:4]:
            p = tf_left.add_paragraph()
            p.text = f"• {bull}"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(5)
            
        # Right Panel (Code Terminal)
        code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = CODE_BG
        code_box.line.color.rgb = CARD_BORDER
        code_box.line.width = Pt(1.5)
        
        # Add a green status dot to terminal top
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(1.95), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(39, 174, 96) # Green
        dot.line.fill.background()
        
        tf_code = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.2)).text_frame
        tf_code.word_wrap = True
        tf_code.margin_left = tf_code.margin_right = tf_code.margin_top = tf_code.margin_bottom = Inches(0)
        
        p_code = tf_code.paragraphs[0]
        # Truncate code to fit
        code_lines = code_item['content'].split('\n')
        truncated_code = '\n'.join(code_lines[:15]) # Limit to 15 lines
        if len(code_lines) > 15:
            truncated_code += "\n... (以下省略) ..."
            
        p_code.text = truncated_code
        p_code.font.name = 'Consolas'
        p_code.font.size = Pt(9)
        p_code.font.color.rgb = RGBColor(57, 255, 20) # Matrix green
    else:
        # Full screen layout (Text card)
        draw_card(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        tf = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(4.2)).text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        first = True
        for para in paragraphs[:4]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = para
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_LIGHT_GRAY
            p.space_after = Pt(12)
            p.line_spacing = 1.2
            
        for bull in bullets[:6]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"•  {bull}"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(8)

# ----------------- Special Project Specific Adjusters -----------------
def add_images_to_presentation(prs, project_id):
    """
    Injects high-quality image slides for specific projects that have design screenshots.
    """
    if project_id not in SPECIAL_PROJECT_IMAGES:
        return
        
    imgs = SPECIAL_PROJECT_IMAGES[project_id]
    
    if project_id == "hongsheng-web":
        # 1. Inject cover slide graphic
        cover_slide = prs.slides[0]
        if os.path.exists(imgs["cover"]):
            cover_slide.shapes.add_picture(imgs["cover"], Inches(5.8), Inches(0), Inches(7.533), Inches(7.5))
            
        # 2. Add a Roles and Status panel slide
        if os.path.exists(imgs["roles"]):
            slide_roles = prs.slides.add_slide(prs.slide_layouts[6])
            set_dark_background(slide_roles)
            add_title(slide_roles, "角色權限與系統登入", RGBColor(0, 229, 255), "操作人員身份對應與權限級別")
            slide_roles.shapes.add_picture(imgs["roles"], Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
            
            # Description text on the right
            draw_card(slide_roles, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
            tf = slide_roles.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.0)).text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = "👥 操作角色權限配置"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 229, 255)
            p.space_after = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = "• 現場操作員 (Operator)：\n  僅能執行條碼掃描與對刷比對，無放行權限。\n\n• 品質管制員 (QC Inspector)：\n  在對刷成功後，需使用 QC 識別證刷卡授權，方能解鎖系統進行放行。\n\n• 系統管理員 (Admin)：\n  可進行配對碼維護、Firebase 連線設定及重置警報狀態。"
            p2.font.name = 'Microsoft JhengHei'
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_LIGHT_GRAY
            p2.line_spacing = 1.3

        # 3. Add Scan Result & QC Dashboard visual slides
        if os.path.exists(imgs["green"]) and os.path.exists(imgs["red"]):
            slide_scan = prs.slides.add_slide(prs.slide_layouts[6])
            set_dark_background(slide_scan)
            add_title(slide_scan, "掃描防錯與即時比對反饋", RGBColor(0, 229, 255), "系統條碼比對成功與警告界面")
            
            slide_scan.shapes.add_picture(imgs["green"], Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.8))
            slide_scan.shapes.add_picture(imgs["red"], Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.8))
            
            # Summary text at the bottom
            draw_card(slide_scan, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1))
            tf = slide_scan.shapes.add_textbox(Inches(1.2), Inches(5.9), Inches(11.0), Inches(0.9)).text_frame
            p = tf.paragraphs[0]
            p.text = "🟢 成功配對（左）：兩端條碼一致，指示燈變綠，傳送數據。\n🔴 異常警報（右）：條碼不匹配，觸發紅色閃爍、蜂鳴警告，拒絕放行。"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            
        if os.path.exists(imgs["qc"]):
            slide_qc = prs.slides.add_slide(prs.slide_layouts[6])
            set_dark_background(slide_qc)
            add_title(slide_qc, "QC 檢驗放行儀表板", RGBColor(0, 229, 255), "品質檢驗與系統最終解鎖放行面板")
            slide_qc.shapes.add_picture(imgs["qc"], Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
            
            draw_card(slide_qc, Inches(7.8), Inches(1.8), Inches(4.733), Inches(4.8))
            tf = slide_qc.shapes.add_textbox(Inches(8.1), Inches(2.2), Inches(4.2), Inches(4.0)).text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = "📊 QC 授權中心"
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 229, 255)
            p.space_after = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = "• 待審查任務：呈現目前已對刷成功、等待品質核准的槽車櫃任務。\n\n• 安全放行鍵：QC 審核數據無誤後，點選放行，系統會透過 API 將狀態寫回 Firebase 雲端資料庫，現場出貨流程即正式核准。"
            p2.font.name = 'Microsoft JhengHei'
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_LIGHT_GRAY
            p2.line_spacing = 1.3

# ----------------- Master Generator -----------------
def generate_pptx_for_all():
    base_dir = os.path.join(root_dir, "說明書")
    categories = [
        "第一類_核心網頁與互動系統",
        "第二類_生產管理與API串接",
        "第三類_AI代理與指南企劃"
    ]
    
    generated_count = 0
    
    # 1. Gather all md files recursively in the categories
    md_files = []
    for cat in categories:
        cat_path = os.path.join(base_dir, cat)
        if not os.path.exists(cat_path):
            print(f"Warning: Category path {cat_path} does not exist. Skipping.")
            continue
            
        for root, dirs, files in os.walk(cat_path):
            for file in files:
                if file.endswith('.md'):
                    filepath = os.path.join(root, file)
                    md_files.append((filepath, cat))
                    
    print(f"Found {len(md_files)} markdown manuals to convert.")
    
    for filepath, cat in md_files:
        try:
            print(f"Processing: {filepath}")
            # Parse MD
            title, github_link, sections = parse_markdown(filepath)
            
            # Setup PPTX
            prs = Presentation()
            prs.slide_width = Inches(13.333) # 16:9
            prs.slide_height = Inches(7.5)
            
            # Choose Theme
            theme = THEMES.get(cat, THEMES["default"])
            
            # Extract project ID
            filename = os.path.basename(filepath)
            project_id = filename.split('_')[0]
            
            # Slide 1: Cover
            build_cover_slide(prs, title, github_link, theme, cat)
            
            # Map sections to slides
            intro_sec = None
            features_sec = None
            tech_sec = None
            steps_sec = None
            extra_secs = []
            
            for sec in sections:
                sec_title = sec['title'].lower()
                if any(x in sec_title for x in ['簡介', '定位', '概述', '說明']):
                    intro_sec = sec
                elif any(x in sec_title for x in ['功能', '特色', '功能特色']):
                    features_sec = sec
                elif any(x in sec_title for x in ['技術', '開發技術', '技術棧']):
                    tech_sec = sec
                elif any(x in sec_title for x in ['步驟', '操作', '執行', '安裝']):
                    steps_sec = sec
                else:
                    extra_secs.append(sec)
            
            # Add Slides in order
            if intro_sec:
                build_intro_slide(prs, intro_sec, theme)
            if features_sec:
                build_features_slide(prs, features_sec, theme)
            if tech_sec:
                build_tech_stack_slide(prs, tech_sec, theme)
            if steps_sec:
                build_steps_slide(prs, steps_sec, theme)
                
            # Add Extra slides
            for sec in extra_secs:
                if len(sec['items']) > 0:
                    build_extra_slide(prs, sec, theme)
                    
            # Inject special high-quality images
            add_images_to_presentation(prs, project_id)
            
            # Generate Output Path
            clean_name = filename.replace('.md', '.pptx')
            local_out = os.path.join(os.path.dirname(filepath), clean_name)
            prs.save(local_out)
            
            print(f"Success! Saved to {local_out}")
            generated_count += 1
            
        except Exception as e:
            print(f"Error processing {filepath}: {str(e)}")
            import traceback
            traceback.print_exc()
            
    print(f"\nCompleted! Successfully generated {generated_count} PPTX manuals.")

if __name__ == "__main__":
    generate_pptx_for_all()
