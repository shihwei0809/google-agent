import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_edu_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    BG_COLOR = RGBColor(9, 13, 22)
    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    COLOR_PRIMARY = RGBColor(56, 189, 248) 
    
    BG_LIGHT = RGBColor(248, 250, 252)
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_DARK_MUTED = RGBColor(71, 85, 105)
    CARD_LIGHT = RGBColor(241, 245, 249)
    CARD_DARK = RGBColor(15, 23, 42)
    CARD_BLUE = RGBColor(22, 28, 45)
    
    FONT_FAMILY = "Microsoft JhengHei"
    
    def set_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, x, y, w, h, text, font_size, color, bold=False):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        return txBox

    # 1. Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide1, BG_COLOR)
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY; bar.line.fill.background()
    add_textbox(slide1, Inches(2), Inches(2.2), Inches(8), Inches(1), "C10200-EDU-01 教育訓練作業管理辦法", 40, TEXT_WHITE, True)
    add_textbox(slide1, Inches(2), Inches(3.2), Inches(8), Inches(0.8), "4.0版 重點宣導", 32, COLOR_PRIMARY, True)
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.2), Inches(6), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = TEXT_MUTED; line.line.fill.background()
    add_textbox(slide1, Inches(2), Inches(4.5), Inches(8), Inches(0.5), "把訓練做成可追蹤、可驗證、可持續改善的管理流程", 20, TEXT_WHITE, False)
    add_textbox(slide1, Inches(2), Inches(6.5), Inches(8), Inches(0.5), "適用：全公司從業人員之職前、在職教育訓練及相關作業", 14, TEXT_MUTED, False)

    # 2. Purpose
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide2, BG_LIGHT)
    add_textbox(slide2, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4), "01 | 先理解目的", 14, COLOR_PRIMARY, True)
    add_textbox(slide2, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8), "這份辦法要解決的，是「能力有沒有真的到位」", 36, TEXT_DARK, True)
    add_textbox(slide2, Inches(0.6), Inches(1.6), Inches(10), Inches(0.4), "訓練不只是在上課，而是從需求、執行到資格與成效的完整閉環。", 16, TEXT_DARK_MUTED, False)
    left_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.5), Inches(5.5), Inches(4))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = CARD_LIGHT; left_box.line.fill.background()
    tf = left_box.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.4)
    p = tf.paragraphs[0]; p.text = "管理目的"; p.font.name = FONT_FAMILY; p.font.size = Pt(20); p.font.color.rgb = COLOR_PRIMARY; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = "讓員工了解公司理念、目標與作業\n流程，具備本職所需的專業知識與\n技能。"; p2.font.name = FONT_FAMILY; p2.font.size = Pt(24); p2.font.color.rgb = TEXT_DARK; p2.font.bold = True; p2.space_before = Pt(20)
    p3 = tf.add_paragraph(); p3.text = "最終目標"; p3.font.name = FONT_FAMILY; p3.font.size = Pt(14); p3.font.color.rgb = TEXT_DARK_MUTED; p3.space_before = Pt(40)
    p4 = tf.add_paragraph(); p4.text = "品質 | 效率 | 環境 | 安全衛生"; p4.font.name = FONT_FAMILY; p4.font.size = Pt(20); p4.font.color.rgb = COLOR_PRIMARY
    rights = [("誰適用？", "全公司所屬從業人員"),("涵蓋什麼？", "職前、在職訓練與相關作業"),("如何落地？", "內訓、外訓、證照、資格驗證與系統留存")]
    for i, (rt, rb) in enumerate(rights):
        y = 2.8 + (i * 1.2)
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(y), Inches(0.1), Inches(0.6))
        bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY; bar.line.fill.background()
        add_textbox(slide2, Inches(7.1), Inches(y-0.1), Inches(5), Inches(0.4), rt, 18, TEXT_DARK, True)
        add_textbox(slide2, Inches(7.1), Inches(y+0.25), Inches(5), Inches(0.4), rb, 14, TEXT_DARK_MUTED, False)

    # 3. Flowchart
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide3, BG_LIGHT)
    add_textbox(slide3, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4), "02 | 管理主線", 14, COLOR_PRIMARY, True)
    add_textbox(slide3, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8), "每一堂訓練，都要走完這條管理鏈", 36, TEXT_DARK, True)
    add_textbox(slide3, Inches(0.6), Inches(1.6), Inches(10), Inches(0.4), "從職能差距出發，最後回到成效與稽核。", 16, TEXT_DARK_MUTED, False)
    steps = [("01", "找需求", "策略 / KPI / 職能差距", CARD_LIGHT, TEXT_DARK),
             ("02", "排計畫", "需求調查→年度計畫", CARD_LIGHT, TEXT_DARK),
             ("03", "做訓練", "內訓、外訓、職前/變更工作前", CARD_LIGHT, TEXT_DARK),
             ("04", "驗資格", "測驗、實作、證照與名冊", CARD_LIGHT, TEXT_DARK),
             ("05", "留證據", "ERP履歷、簽到、心得、成效", CARD_DARK, TEXT_WHITE)]
    start_x = 0.6; w = 2.2
    for i, (num, title, desc, bg_c, text_c) in enumerate(steps):
        x = start_x + (i * 2.5)
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(w), Inches(2.5))
        box.fill.solid(); box.fill.fore_color.rgb = bg_c; box.line.fill.background()
        num_c = COLOR_PRIMARY if bg_c == CARD_LIGHT else COLOR_PRIMARY
        add_textbox(slide3, Inches(x+0.2), Inches(2.8), Inches(1.8), Inches(0.4), num, 24, num_c, True)
        add_textbox(slide3, Inches(x+0.2), Inches(3.4), Inches(1.8), Inches(0.4), title, 20, text_c, True)
        desc_box = add_textbox(slide3, Inches(x+0.2), Inches(4.0), Inches(1.8), Inches(0.8), desc, 12, TEXT_DARK_MUTED if bg_c == CARD_LIGHT else TEXT_WHITE, False)
        desc_box.text_frame.word_wrap = True
        if i < 4:
            arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+w+0.05), Inches(3.6), Inches(0.2), Inches(0.2))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_PRIMARY; arrow.line.fill.background()
    add_textbox(slide3, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5), "判斷標準：訓練完成 ≠ 上完課；必須能證明「人員受訓、能力驗證、紀錄可查」。", 18, COLOR_PRIMARY, True)

    def create_dark_content_slide(section_num, section_title, main_title, content_list):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(s, BG_COLOR)
        add_textbox(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4), f"{section_num} | {section_title}", 14, COLOR_PRIMARY, True)
        add_textbox(s, Inches(0.6), Inches(0.8), Inches(10), Inches(0.8), main_title, 32, TEXT_WHITE, True)
        
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(12.1), Inches(5))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BLUE; box.line.fill.background()
        box.line.color.rgb = COLOR_PRIMARY; box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.4)
        for pt in content_list:
            p = tf.add_paragraph()
            if pt.startswith("  -"):
                p.text = pt.replace("  -", "•").strip()
                p.level = 1
                p.font.size = Pt(14)
                p.font.color.rgb = TEXT_MUTED
            elif pt.startswith("-"):
                p.text = pt.replace("-", "■").strip()
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_WHITE
            else:
                p.text = pt
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_WHITE
            p.font.name = FONT_FAMILY
            p.space_after = Pt(10)

    create_dark_content_slide("03", "權責區分", "總務課與各單位的分工", [
        "- 總務課：",
        "  - 教育訓練制度訂修、全公司年度計畫彙總與成效評估。",
        "  - 督導各單位實施情形、辦理專業管理類與外派研習。",
        "  - 新進人員職前訓練安排與場所維護。",
        "- 各單位：",
        "  - 單位內計畫擬定、內部(管理/技術)訓練辦理。",
        "  - 單位內新進人員訓練及實施情形彙總。"
    ])

    create_dark_content_slide("04", "證照與計畫", "法規證照管理與需求調查", [
        "- 訓練需求評估：",
        "  - 各單位依「工作敘述表」、年度目標與 KPI 評估職能差距提報。",
        "  - 總務課於每年年底調查需求，制定《年度教育訓練計畫表》。",
        "- 法規證照管理分類：",
        "  - 【A類】需向主管機關報備，及內部診斷人員。",
        "  - 【B類】不須向主管機關報備，僅存公司備查。",
        "  - 【C類】內部程序書或指導書規定者。",
        "- 安全衛生回訓規定：",
        "  - 一般人員 3 小時 / 特殊作業(如電焊、缺氧)增列 3 小時。"
    ])

    create_dark_content_slide("05", "教育訓練設計", "講師資格與教材編選", [
        "- 內部講師資格：",
        "  - 需由組長級以上，或擔任相關職務滿 3 年者擔任。",
        "  - 必須受過相關訓練或具備合格證書(經、副理為當然講師)。",
        "- 外部講師資格：",
        "  - 邀請機關團體專業人士或企管顧問。",
        "- 教材編選與存查：",
        "  - 內訓：講師編撰或沿用公司存檔教材。",
        "  - 外訓：參加者提供教材供傳閱或列入圖書備查，外聘講師教材需存參。"
    ])

    create_dark_content_slide("06", "教育訓練執行", "內部與外部訓練執行要點", [
        "- 內訓執行程序：",
        "  - 課前：變更工作前需受安衛訓練並填寫《教育訓練紀錄表》，並填《上課人員報名表》。",
        "  - 課中：務必簽到《課程簽到表》。",
        "  - 課後：填寫《課後意見調查表》送總務課，受訓資料鍵入 ERP 系統存檔。",
        "- 外訓執行程序：",
        "  - 申請：於 QRisk 查詢證照需求，ERP 填寫《外派訓練申請單》。",
        "  - 結訓 10 日內：必須上系統填寫《心得報告》。",
        "  - 證照處理：法規特定證照正本留存管理部並登錄 QRisk，其餘繳回管理部。",
        "  - 費用與加班：非上班日可報加班，免付費課程亦須填寫申請表(費用為0)。"
    ])

    create_dark_content_slide("07", "考核與稽核", "資格驗證、費用與紀錄保存", [
        "- 訓練考核方式：",
        "  - 口試、筆試、技能實作測驗、學前/學後測驗比較。",
        "  - 教育訓練考核成績將列入「個人年度績效考核」。",
        "- 特定資格驗證：",
        "  - 校正人員：受過儀器校驗訓練並考試合格。",
        "  - 檢驗人員：品管主管實施程序書訓練且實作合格。",
        "- 費用與稽核：",
        "  - 內部講師鐘點費：每小時 200 元 (由總務課填寫申請單)。",
        "  - 稽核：每年至少執行一次內部診斷(ISO標準)，並於管理審查中報告。",
        "  - 保存：免付費研討會採無紙化存檔，自辦教材自行保存。"
    ])

    output_path = r"D:\GOOGLE ANGET\教育訓練-C0588-教材\C10200-EDU-01 教育訓練作業管理辦法(Codex樣式完整版).pptx"
    prs.save(output_path)

create_edu_ppt()

